import json
from pathlib import Path
from typing import Dict, List
import os
import re
import math
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
      
from pprint import pprint

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE   
from pptx.util import Inches

import json
import difflib
from pptx import Presentation
from SlidesAgent.apply_color import *



COLOR_WHITE = RGBColor(0, 0, 0) 
THEME_COLOR = RGBColor(185, 210, 153) 
 

TEMPLATE_ID_ALIASES = {
    "T1_LeftImage": "T3_ImageLeft",
    "T1_RightImage": "T2_ImageRight",
    "T4_ImageLeft": "T3_ImageLeft",
    "T4_ImageRight": "T2_ImageRight",
    "T3_ImageTop": "T4_ImageTop",
}

EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def normalize_template_id(template_id: str, layout_names: List[str]) -> str:
    normalized = TEMPLATE_ID_ALIASES.get(template_id, template_id)
    if normalized in layout_names:
        return normalized

    close_match = difflib.get_close_matches(normalized, layout_names, n=1, cutoff=0.88)
    if close_match:
        print(f"[layout_filler] normalizing template_id '{template_id}' -> '{close_match[0]}'")
        return close_match[0]

    return normalized


def _insert_picture_keep_ratio(ph, img_path: Path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image

    slide = ph.part.slide   
    ph_left, ph_top = ph.left, ph.top
    ph_width, ph_height = ph.width, ph.height
 
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
 
    if ph_width / ph_height > aspect:
        new_h = ph_height
        new_w = int(new_h * aspect)
    else:
        new_w = ph_width
        new_h = int(new_w / aspect)
 
    new_left = ph_left + int((ph_width - new_w) / 2)
    new_top = ph_top + int((ph_height - new_h) / 2)
 
    pic = slide.shapes.add_picture(str(img_path), new_left, new_top, width=new_w, height=new_h)
 
    # --- move picture *behind* all text placeholders ---
    spTree = slide.shapes._spTree
    spTree.remove(pic.element)         # temporarily take it out
    spTree.insert(2, pic.element)      # index 0=background,1=layout 

    # --- finally, remove the now‑unused picture placeholder itself ---
    ph.element.getparent().remove(ph.element)
  
from pptx.util import Pt
from PIL import Image

def insert_image_below_content(slide, img_path: Path):
    """
    Insert an image below the lowest existing shape (text or image) on the slide.
    Centered horizontally. Resizes if not enough space.
    """
     
    with Image.open(img_path) as img:
        width_px, height_px = img.size
    aspect_ratio = width_px / height_px

    # slide size（EMUs）
    slide_width = slide.part.slide_layout.part.package.presentation_part.slide_width
    slide_height = slide.part.slide_layout.part.package.presentation_part.slide_height

    
    target_width = slide_width * 0.6
    target_height = target_width / aspect_ratio

     
    lowest_bottom = 0
    for shape in slide.shapes:
        bottom = shape.top + shape.height
        if bottom > lowest_bottom:
            lowest_bottom = bottom

    margin = Pt(20)
    available_space = slide_height - lowest_bottom - margin

    if available_space < target_height:
        target_height = available_space
        target_width = target_height * aspect_ratio
        if target_height <= 0:
            print("Not enough space to insert image:", img_path)
            return

    
    left = (slide_width - target_width) // 2
    top = lowest_bottom + margin

    slide.shapes.add_picture(str(img_path), left, top, width=int(target_width), height=int(target_height))

 

TEXT_TYPES = {
    PH_TYPE.TITLE,
    PH_TYPE.CENTER_TITLE,
    PH_TYPE.SUBTITLE,
    PH_TYPE.BODY,
}


def _set_paragraph_font_size(paragraph, size_pt: float) -> None:
    paragraph.font.size = Pt(size_pt)
    for run in paragraph.runs:
        run.font.size = Pt(size_pt)


def _enable_shrink_to_fit(text_frame, max_size_pt: float | None = None) -> None:
    if text_frame is None:
        return
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    if max_size_pt is not None:
        for paragraph in text_frame.paragraphs:
            _set_paragraph_font_size(paragraph, max_size_pt)


def _compact_text_frame(text_frame, margin_pt: float = 3.0) -> None:
    if text_frame is None:
        return
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Pt(margin_pt)
    text_frame.margin_right = Pt(margin_pt)
    text_frame.margin_top = Pt(margin_pt)
    text_frame.margin_bottom = Pt(margin_pt)


def _paragraph_text(paragraph) -> str:
    if paragraph is None:
        return ""
    text = paragraph.text or ""
    if text:
        return text
    return "".join(run.text for run in paragraph.runs if getattr(run, "text", ""))


def _paragraph_font_size_pt(paragraph, fallback: float) -> float:
    if paragraph is None:
        return fallback
    if paragraph.font is not None and paragraph.font.size is not None:
        return max(float(paragraph.font.size.pt), 1.0)
    for run in paragraph.runs:
        if run.font is not None and run.font.size is not None:
            return max(float(run.font.size.pt), 1.0)
    return fallback


def _estimate_wrapped_line_count(text: str, usable_width_pt: float, font_size_pt: float) -> int:
    text = (text or "").strip()
    if not text:
        return 1
    # Conservative estimate for proportional fonts in PowerPoint.
    chars_per_line = max(4, int(usable_width_pt / max(font_size_pt * 0.5, 1.0)))
    line_count = 0
    for logical_line in text.splitlines() or [""]:
        line_count += max(1, math.ceil(max(len(logical_line), 1) / chars_per_line))
    return line_count


def _text_frame_fits_estimate(shape, text_frame, min_size_pt: float) -> float:
    width_pt = max(shape.width / EMU_PER_PT - 2 * 3.0, 24.0)
    height_pt = max(shape.height / EMU_PER_PT - 2 * 3.0, 24.0)
    total_height_pt = 0.0

    paragraphs = list(text_frame.paragraphs or [])
    for idx, paragraph in enumerate(paragraphs):
        para_text = _paragraph_text(paragraph)
        font_size_pt = max(_paragraph_font_size_pt(paragraph, min_size_pt), min_size_pt)
        level = max(int(getattr(paragraph, "level", 0) or 0), 0)
        indent_penalty_pt = min(level * 18.0, width_pt * 0.45)
        usable_width_pt = max(width_pt - indent_penalty_pt, 24.0)
        line_count = _estimate_wrapped_line_count(para_text, usable_width_pt, font_size_pt)
        line_height_pt = font_size_pt * (1.18 if level == 0 else 1.12)
        total_height_pt += line_count * line_height_pt
        if idx < len(paragraphs) - 1:
            total_height_pt += max(font_size_pt * 0.18, 1.5)

    return total_height_pt - height_pt


def _text_frame_fill_ratio(shape, text_frame, min_size_pt: float) -> float:
    height_pt = max(shape.height / EMU_PER_PT - 2 * 3.0, 24.0)
    overflow_pt = _text_frame_fits_estimate(shape, text_frame, min_size_pt)
    content_height_pt = max(0.0, height_pt + overflow_pt)
    return content_height_pt / height_pt if height_pt > 0 else 1.0


def _shrink_text_frame_to_fit(shape, text_frame, min_size_pt: float) -> None:
    if shape is None or text_frame is None:
        return
    original_vertical_anchor = getattr(text_frame, "vertical_anchor", None)
    original_alignments = [getattr(paragraph, "alignment", None) for paragraph in text_frame.paragraphs]
    paragraphs = [p for p in text_frame.paragraphs if _paragraph_text(p).strip()]
    if not paragraphs:
        _compact_text_frame(text_frame)
        _enable_shrink_to_fit(text_frame)
        if original_vertical_anchor is not None:
            text_frame.vertical_anchor = original_vertical_anchor
        return

    _compact_text_frame(text_frame)
    overflow_pt = _text_frame_fits_estimate(shape, text_frame, min_size_pt)
    if overflow_pt <= 0:
        _enable_shrink_to_fit(text_frame)
        if original_vertical_anchor is not None:
            text_frame.vertical_anchor = original_vertical_anchor
        for paragraph, alignment in zip(text_frame.paragraphs, original_alignments):
            paragraph.alignment = alignment
        return

    current_sizes = [_paragraph_font_size_pt(p, min_size_pt) for p in paragraphs]
    max_current = max(current_sizes) if current_sizes else min_size_pt
    trial_size = max_current

    while trial_size > min_size_pt:
        scale = trial_size / max_current if max_current > 0 else 1.0
        for paragraph, original_size in zip(paragraphs, current_sizes):
            _set_paragraph_font_size(paragraph, max(min_size_pt, original_size * scale))
        if _text_frame_fits_estimate(shape, text_frame, min_size_pt) <= 0:
            break
        trial_size -= 1.0

    _enable_shrink_to_fit(text_frame)
    if original_vertical_anchor is not None:
        text_frame.vertical_anchor = original_vertical_anchor
    for paragraph, alignment in zip(text_frame.paragraphs, original_alignments):
        paragraph.alignment = alignment


def _grow_text_frame_to_fill(
    shape,
    text_frame,
    *,
    min_size_pt: float,
    max_size_pt: float,
    target_min_fill: float,
    target_max_fill: float,
) -> None:
    if shape is None or text_frame is None:
        return
    original_vertical_anchor = getattr(text_frame, "vertical_anchor", None)
    original_alignments = [getattr(paragraph, "alignment", None) for paragraph in text_frame.paragraphs]
    paragraphs = [p for p in text_frame.paragraphs if _paragraph_text(p).strip()]
    if not paragraphs:
        return

    current_fill = _text_frame_fill_ratio(shape, text_frame, min_size_pt)
    if current_fill >= target_min_fill:
        return

    current_sizes = [_paragraph_font_size_pt(p, min_size_pt) for p in paragraphs]
    max_current = max(current_sizes) if current_sizes else min_size_pt
    trial_size = min(max_current + 1.0, max_size_pt)
    best_sizes = list(current_sizes)
    best_fill = current_fill

    while trial_size <= max_size_pt:
        scale = trial_size / max_current if max_current > 0 else 1.0
        for paragraph, original_size in zip(paragraphs, current_sizes):
            _set_paragraph_font_size(paragraph, min(max_size_pt, original_size * scale))

        overflow_pt = _text_frame_fits_estimate(shape, text_frame, min_size_pt)
        fill_ratio = _text_frame_fill_ratio(shape, text_frame, min_size_pt)

        if overflow_pt > 0 or fill_ratio > target_max_fill:
            break

        best_sizes = [_paragraph_font_size_pt(p, min_size_pt) for p in paragraphs]
        best_fill = fill_ratio
        if fill_ratio >= target_min_fill:
            break
        trial_size += 1.0

    if best_fill > current_fill:
        for paragraph, best_size in zip(paragraphs, best_sizes):
            _set_paragraph_font_size(paragraph, best_size)
    else:
        for paragraph, original_size in zip(paragraphs, current_sizes):
            _set_paragraph_font_size(paragraph, original_size)

    _enable_shrink_to_fit(text_frame)
    if original_vertical_anchor is not None:
        text_frame.vertical_anchor = original_vertical_anchor
    for paragraph, alignment in zip(text_frame.paragraphs, original_alignments):
        paragraph.alignment = alignment


def _shape_rect(shape) -> tuple[int, int, int, int]:
    return (int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height))


def _rects_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int], padding: int = 0) -> bool:
    return not (
        a[2] <= b[0] + padding
        or a[0] >= b[2] - padding
        or a[3] <= b[1] + padding
        or a[1] >= b[3] - padding
    )


def _slide_dimensions(slide) -> tuple[int, int]:
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def _text_shapes_with_content(slide):
    shapes = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if not tf or not any(_paragraph_text(p).strip() for p in tf.paragraphs):
            continue
        shapes.append(shape)
    return shapes


def _expand_sparse_slide_pictures(slide) -> None:
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        return

    text_shapes = _text_shapes_with_content(slide)
    text_fill_ratios = []
    for shape in text_shapes:
        tf = shape.text_frame
        min_size_pt = 12.0 if "title" in (getattr(shape, "name", "") or "").lower() else 10.0
        text_fill_ratios.append(min(_text_frame_fill_ratio(shape, tf, min_size_pt), 1.0))

    avg_text_fill = sum(text_fill_ratios) / len(text_fill_ratios) if text_fill_ratios else 0.0
    slide_width, slide_height = _slide_dimensions(slide)
    slide_area = max(slide_width * slide_height, 1)
    picture_area_ratio = sum(int(p.width) * int(p.height) for p in pictures) / slide_area

    is_sparse = avg_text_fill < 0.42 or picture_area_ratio < 0.18
    if not is_sparse:
        return

    slide_margin = int(0.12 * EMU_PER_INCH)
    blocker_padding = int(0.08 * EMU_PER_INCH)

    for pic in sorted(pictures, key=lambda shp: shp.width * shp.height, reverse=True):
        center_x = pic.left + pic.width / 2
        center_y = pic.top + pic.height / 2
        start_w = int(pic.width)
        start_h = int(pic.height)
        best = (int(pic.left), int(pic.top), start_w, start_h)
        max_scale = 1.75 if len(pictures) == 1 else 1.4
        scale = 1.05

        while scale <= max_scale:
            new_w = int(start_w * scale)
            new_h = int(start_h * scale)
            new_left = int(center_x - new_w / 2)
            new_top = int(center_y - new_h / 2)
            rect = (new_left, new_top, new_left + new_w, new_top + new_h)

            if (
                rect[0] < slide_margin
                or rect[1] < slide_margin
                or rect[2] > slide_width - slide_margin
                or rect[3] > slide_height - slide_margin
            ):
                break

            blocked = False
            for other in slide.shapes:
                if other == pic:
                    continue
                if other.shape_type == MSO_SHAPE_TYPE.PICTURE or other in text_shapes:
                    if _rects_overlap(rect, _shape_rect(other), padding=blocker_padding):
                        blocked = True
                        break
            if blocked:
                break

            best = (new_left, new_top, new_w, new_h)
            scale += 0.05

        pic.left, pic.top, pic.width, pic.height = best


def _set_shape_text_with_fit(
    shape,
    text: str,
    max_size_pt: float,
    min_size_pt: float = 12.0,
    alignment=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = text or ""
    _compact_text_frame(tf)
    _enable_shrink_to_fit(tf, max_size_pt=max_size_pt)
    _shrink_text_frame_to_fit(shape, tf, min_size_pt=min_size_pt)
    tf.vertical_anchor = vertical_anchor
    for paragraph in tf.paragraphs:
        paragraph.alignment = alignment


def _contents_font_size(section_count: int) -> int:
    if section_count <= 5:
        return 36
    if section_count <= 7:
        return 30
    if section_count <= 9:
        return 26
    if section_count <= 11:
        return 20
    if section_count <= 13:
        return 18
    return 16


def _populate_contents_frame(tf, sections, font_size_pt: int) -> None:
    tf.clear()
    _enable_shrink_to_fit(tf)
    for i, sec in enumerate(sections):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = sec
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        _set_paragraph_font_size(p, font_size_pt)


def _header_font_size(text: str) -> int:
    n = len((text or "").strip())
    if n <= 24:
        return 24
    if n <= 36:
        return 20
    if n <= 52:
        return 18
    return 16


def _normalize_bullets(raw_bullets) -> list[dict]:
    normalized: list[dict] = []
    for bullet in (raw_bullets or []):
        if isinstance(bullet, dict):
            raw_text = str(bullet.get("text") or "").strip()
            raw_sub = bullet.get("sub") or []
            if not isinstance(raw_sub, list):
                raw_sub = [raw_sub]
            sub_items = [str(item).strip() for item in raw_sub if str(item).strip()]

            # Some plans emit {"sub": [...]} without a top-level text field.
            # Promote the first sub-bullet to the main bullet so rendering can continue.
            if not raw_text and sub_items:
                raw_text = sub_items[0]
                sub_items = sub_items[1:]

            if raw_text or sub_items:
                normalized.append({"text": raw_text, "sub": sub_items})
            continue

        text = str(bullet).strip()
        if text:
            normalized.append({"text": text, "sub": []})
    return normalized


def _body_font_sizes(slide_info: Dict) -> tuple[int, int]:
    bullets = _normalize_bullets(slide_info.get("bullets") or [])
    bullet_count = len(bullets)
    sub_count = sum(len(b.get("sub") or []) for b in bullets)
    total_chars = sum(len((b.get("text") or "").strip()) for b in bullets)
    total_chars += sum(len(str(s).strip()) for b in bullets for s in (b.get("sub") or []))
    visual_count = len(slide_info.get("images") or []) + len(slide_info.get("tables") or []) + len(slide_info.get("formulas") or [])

    lvl0, lvl1 = 24, 22
    if visual_count > 0:
        lvl0, lvl1 = 20, 18
    if bullet_count >= 3 or sub_count >= 4 or total_chars >= 180:
        lvl0 -= 2
        lvl1 -= 2
    if bullet_count >= 4 or sub_count >= 7 or total_chars >= 260:
        lvl0 -= 2
        lvl1 -= 2
    return max(lvl0, 16), max(lvl1, 14)


def _should_use_top_visual_layout(slide_info: Dict, visual_paths: List[Path]) -> bool:
    template_id = str(slide_info.get("template_id") or "")
    if template_id not in {"T2_ImageRight", "T3_ImageLeft"}:
        return False
    if len(visual_paths) != 1:
        return False
    if slide_info.get("formulas"):
        return False

    visual_path = visual_paths[0]
    try:
        with Image.open(visual_path) as img:
            width_px, height_px = img.size
    except Exception:
        return False

    aspect_ratio = width_px / max(height_px, 1)
    bullet_count = len(slide_info.get("bullets") or [])
    sub_count = sum(len(b.get("sub") or []) for b in (slide_info.get("bullets") or []))
    is_wide_visual = aspect_ratio >= 1.25
    has_meaningful_text = bullet_count >= 2 or sub_count >= 2
    return is_wide_visual and has_meaningful_text


VISUAL_TEMPLATE_IDS = {
    "T2_ImageRight",
    "T3_ImageLeft",
    "T4_ImageTop",
    "T5_TwoImages",
    "T5_TwoImages2",
    "T7_2x2_TopImage",
    "T8_2x2_BottomImage",
    "T9_2x2_AltTextImg",
    "T10_4Img_2x2Grid",
    "T11_3Img_TopTextBottom",
    "T12_3Img_BottomTextTop",
    "T13_3Img",
    "T14_ImageRight_1Formula",
    "T15_ImageLeft_1Formula",
    "T16_1Img_2formula_TopTextBottom",
    "T17_2Img_1formula_TopTextBottom",
    "T18_2formula_TopTextBottom",
}


def _fallback_text_only_template(template_id: str, slide_info: Dict) -> str:
    if template_id not in VISUAL_TEMPLATE_IDS:
        return template_id
    has_requested_visuals = bool(slide_info.get("images") or slide_info.get("tables") or slide_info.get("formulas"))
    if has_requested_visuals:
        return template_id
    return "T1_TextOnly"


def _fallback_missing_template_id(
    template_id: str,
    slide_info: Dict,
    visual_paths: List[Path],
    layout_names: List[str],
) -> str:
    if template_id in layout_names:
        return template_id

    normalized = str(template_id or "")
    visual_count = len(visual_paths)
    has_text = bool(slide_info.get("bullets"))

    candidates: list[str] = []

    if "MultiVisual" in normalized or "multi" in normalized.lower():
        if visual_count >= 2:
            candidates.extend(["T5_TwoImages2", "T5_TwoImages", "T10_4Img_2x2Grid"])
        elif visual_count == 1:
            candidates.extend(["T4_ImageTop", "T2_ImageRight", "T3_ImageLeft"])
        else:
            candidates.append("T1_TextOnly")

    if "ImageTop" in normalized:
        candidates.extend(["T4_ImageTop", "T2_ImageRight", "T3_ImageLeft"])
    if "ImageRight" in normalized:
        candidates.extend(["T2_ImageRight", "T4_ImageTop", "T3_ImageLeft"])
    if "ImageLeft" in normalized:
        candidates.extend(["T3_ImageLeft", "T4_ImageTop", "T2_ImageRight"])

    if visual_count >= 2:
        candidates.extend(["T5_TwoImages2", "T5_TwoImages"])
    elif visual_count == 1:
        candidates.extend(["T4_ImageTop", "T2_ImageRight", "T3_ImageLeft"])
    elif not has_text:
        candidates.append("T1_TextOnly")

    seen = set()
    for candidate in candidates:
        if candidate in seen:
            continue
        seen.add(candidate)
        if candidate in layout_names:
            print(f"[layout_filler] fallback template_id '{template_id}' -> '{candidate}'")
            return candidate

    return template_id

def find_text_placeholders(slide): 
    """Return (part_num_ph, subsection_ph, body_ph) by position."""
    txt_ph = [
        s for s in slide.shapes
        if (
            s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            s.has_text_frame and
            s.placeholder_format.type in TEXT_TYPES   
        )
    ]

    if len(txt_ph) < 2:
        raise ValueError("Not enough text placeholders on this slide")

     
    txt_ph.sort(key=lambda s: s.top)
 
    first_row = sorted(txt_ph[:2], key=lambda s: s.left)
    part_ph, title_ph = first_row
    body_ph = txt_ph[2] if len(txt_ph) >= 3 else None

    return part_ph, title_ph, body_ph

def get_content(sec_title, sub_title,outline):
    _, sub = _best_match(outline, sec_title, sub_title)
    return sub.get("content", "") if sub else ""

def set_font_color(paragraph, theme_color):
    paragraph.font.fill.solid()
    if theme_color is not None:
        paragraph.font.fill.fore_color.theme_color = theme_color
    else:
        paragraph.font.fill.fore_color.rgb = RGBColor(255, 105, 180)  # pink fallback

def insert_visuals_auto(slide, visuals: list[Path]):
    """
    Automatically insert visuals (images, tables, formulas) into all available
    picture placeholders on the given slide.
    """
    # Find all picture placeholder shapes
    picture_placeholders = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
           "Picture" in shape.name
    ]

    if len(visuals) > len(picture_placeholders):
        print(f"Warning: not enough picture placeholders on slide (needed {len(visuals)}, found {len(picture_placeholders)})")
        # remaining = visuals[len(picture_placeholders):]
        # for img_path in remaining:
        #     insert_image_below_content(slide, Path(img_path))

    used_placeholders = []

    #  Insert images one by one
    for img_path, ph in zip(visuals, picture_placeholders):
        _insert_picture_keep_ratio(ph, Path(img_path))
        used_placeholders.append(ph)

    # Remove any unused picture placeholders so PowerPoint does not display
    # "Insert Picture" boxes when the planner chose a visual template without
    # enough visuals to fill it.
    for ph in picture_placeholders[len(used_placeholders):]:
        parent = ph.element.getparent()
        if parent is not None:
            parent.remove(ph.element)


def _placeholder_by_name(slide, name: str):
    """Return placeholder shape whose .name == name."""

    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f'Placeholder "{name}" not found on slide master.')
 
 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def scan_layout_placeholders(template_path: str):
    prs = Presentation(template_path)
    layout_map = {}

    for layout in prs.slide_layouts:
        layout_name = layout.name
        placeholder_names = []

        for shape in layout.shapes:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type)
            try:
                shape_type_name = MSO_SHAPE_TYPE(shape_type).name
            except ValueError:
                pass

            placeholder_names.append({
                "name": shape.name,
                "type": shape_type_name
            })

        layout_map[layout_name] = placeholder_names

    return layout_map


from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_theme_color_from_title(prs, layout_index=2):
    slide_layout = prs.slide_layouts[layout_index]
    temp_slide = prs.slides.add_slide(slide_layout)
 
    for shape in temp_slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and 
            shape.placeholder_format.type == 1  # 1 == TITLE
        ):
            para = shape.text_frame.paragraphs[0]
            font_color = para.font.color
            if font_color.type == 2:  # 2 == THEME
                theme_color = font_color.theme_color 
                xml_slides = prs.slides._sldIdLst  
                xml_slides.remove(xml_slides[-1])  # remove last slide
                return theme_color
    return None
 
def resolve_visual_paths(slide_info, args):
    import os
    import re
    import json
    from pathlib import Path

    paper = args.paper_name
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = Path(prefix) / paper
 
    images_json_path = Path(prefix) / f"{paper}_images.json"
    tables_json_path = Path(prefix) / f"{paper}_tables.json"
    images_json = json.load(open(images_json_path, "r", encoding="utf-8")) if images_json_path.exists() else {}
    tables_json = json.load(open(tables_json_path, "r", encoding="utf-8")) if tables_json_path.exists() else {}
 
    if args.formula_mode == 3:
        formulas_dir = Path("contents") / paper / "formula_images"
    else:
        formulas_dir = base_dir
 
    def _norm_digits(s: str) -> str:
 
        m = re.search(r'(\d+)(?!.*\d)', str(s))
        if not m:
            return str(s)
        v = m.group(1).lstrip("0")
        return v if v != "" else "0"

    def _stem_name(p: str) -> str:
        return Path(str(p)).name

    def _extract_figure_numbers(caption: str):
         
        out = set()
        if not caption:
            return out
        caps = str(caption)

        #  Figure 6–7 / Fig. 10-12 / 图 3-5
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        #  Figure 6 / Fig. 6 / 图6
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _extract_table_numbers(caption: str):
         
        out = set()
        if not caption:
            return out
        caps = str(caption)

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _build_caption_index_images(images_dict: dict):
         
        idx = {}
        for k, meta in (images_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for f in _extract_figure_numbers(caps):
                idx.setdefault(f, set()).add(str(k))
        return idx

    def _build_caption_index_tables(tables_dict: dict):
         
        idx = {}
        for k, meta in (tables_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for t in _extract_table_numbers(caps):
                idx.setdefault(t, set()).add(str(k))
        return idx

    fignum_to_imgkeys = _build_caption_index_images(images_json)
    tbnum_to_tbkeys  = _build_caption_index_tables(tables_json)

     
    images_json = {str(k): v for k, v in images_json.items()}
    tables_json = {str(k): v for k, v in tables_json.items()}

    def _resolve_and_check(path1: Path) -> Path:
        
        if path1.exists():
            return path1
 
        p2 = base_dir / path1.name
        if p2.exists():
            return p2
 
        candidates = [p for p in base_dir.glob("*") if p.name.lower() == path1.name.lower()]
        if candidates:
            return candidates[0]
 
        existing = "\n".join(sorted(p.name for p in base_dir.glob("*.*"))[:80])
        raise FileNotFoundError(
            f"Missing visual file: {path1}\n"
            f"Tried: {path1}, {p2}. Base dir: {base_dir}\n"
            f"Existing (first 80):\n{existing}"
        )

    def _fallback_visual_file_by_number(kind: str, raw_id: str, name: str) -> Path | None:
        norm = _norm_digits(raw_id) or _norm_digits(name)
        if not norm or not norm.isdigit():
            return None

        direct_patterns = []
        if kind == "image":
            direct_patterns = [
                f"*-picture-{norm}.png",
                f"*-{norm}.png",
                f"image_{int(norm):06d}_*.png",
                f"image_{norm}.png",
            ]
        elif kind == "table":
            direct_patterns = [
                f"*-table-{norm}.png",
                f"table_{norm}.png",
            ]

        candidates: list[Path] = []
        artifacts_dir = base_dir / f"{args.paper_name}-with-image-refs_artifacts"
        for pattern in direct_patterns:
            candidates.extend(sorted(base_dir.glob(pattern)))
            if artifacts_dir.exists():
                candidates.extend(sorted(artifacts_dir.glob(pattern)))
            candidates.extend(sorted(base_dir.glob(f"**/{pattern}")))

        deduped: list[Path] = []
        seen = set()
        for path in candidates:
            key = str(path.resolve())
            if key in seen or not path.exists():
                continue
            seen.add(key)
            deduped.append(path)

        return deduped[0] if deduped else None

    def _dedupe_paths(paths):
        seen = set()
        deduped = []
        for path in paths:
            key = str(path.resolve()) if path.exists() else str(path)
            if key in seen:
                continue
            seen.add(key)
            deduped.append(path)
        return deduped

    def _match_by_id_or_caption(img_id: str, name: str, mapping: dict, cap_index: dict, kind: str):
          
        img_id_norm = _norm_digits(img_id)
        name_norm   = _norm_digits(_stem_name(name))
 
        if img_id in mapping:
            return mapping[img_id], img_id
        if img_id_norm in mapping:
            return mapping[img_id_norm], img_id_norm
        if name_norm in mapping:
            return mapping[name_norm], name_norm
 
        if img_id_norm.isdigit():
            num = int(img_id_norm)
            cand = sorted(cap_index.get(num, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand:
                k = cand[0]
                return mapping[k], k
 
        if name_norm.isdigit():
            num2 = int(name_norm)
            cand2 = sorted(cap_index.get(num2, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand2:
                k2 = cand2[0]
                return mapping[k2], k2
 
        avail = sorted(mapping.keys(), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
        raise KeyError(
            f"{kind.capitalize()} id not found. "
            f"given_id='{img_id}', file='{name}', "
            f"norm(id)={img_id_norm}, norm(file)={name_norm}. "
            f"Available keys (first 40): {avail[:40]}"
        )

    # ----------------- images -----------------
    image_paths = []
    for img_str in slide_info.get("images", []):
        name = _stem_name(img_str)
 
        #   'picture|image|fig|table|formula-<num>.ext'
        m = re.search(r'(?:picture|image|fig|table|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        img_id = m.group(1) if m else _norm_digits(name)

       
        try:
            rec, used_key = _match_by_id_or_caption(img_id, name, images_json, fignum_to_imgkeys, kind="image")
            target = rec.get("image_path") or rec.get("path") or rec.get("file")
            if not target:
                raise KeyError(f"Image record has no path field. key={used_key}, record={rec}")
            img_path = _resolve_and_check(Path(target))
        except KeyError:
            fallback = _fallback_visual_file_by_number("image", img_id, name)
            if fallback is None:
                raise
            print(f"[layout_filler] fallback image match for '{name}' -> '{fallback.name}'")
            img_path = fallback
        image_paths.append(img_path)

    # ----------------- tables -----------------
    table_paths = []
    for tb_str in slide_info.get("tables", []):
        name = _stem_name(tb_str)
        m = re.search(r'(?:table|tab|picture|image|fig|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        tb_id = m.group(1) if m else _norm_digits(name)

        try:
            rec, used_key = _match_by_id_or_caption(tb_id, name, tables_json, tbnum_to_tbkeys, kind="table")
            target = rec.get("table_path") or rec.get("image_path") or rec.get("path") or rec.get("file")
            if not target:
                raise KeyError(f"Table record has no path field. key={used_key}, record={rec}")
            tb_path = _resolve_and_check(Path(target))
        except KeyError:
            fallback = _fallback_visual_file_by_number("table", tb_id, name)
            if fallback is None:
                print(f"[layout_filler] skipping missing table asset: {name}")
                continue
            print(f"[layout_filler] fallback table match for '{name}' -> '{fallback.name}'")
            tb_path = fallback
        table_paths.append(tb_path)

    # ----------------- formulas -----------------
    formula_paths = []
    for fname in slide_info.get("formulas", []):
        if args.formula_mode == 3:
            final_path = formulas_dir / fname
        else: 
            resolved_formula_path = resolve_formula_mode1_path(fname, args)
            if resolved_formula_path is None:
                print(f"[layout_filler] skipping non-asset formula entry: {fname}")
                continue
            final_path = Path(resolved_formula_path)
        try:
            formula_paths.append(_resolve_and_check(Path(final_path)))
        except FileNotFoundError:
            fallback = _fallback_visual_file_by_number("image", _norm_digits(fname), fname)
            if fallback is not None:
                print(f"[layout_filler] fallback formula match for '{fname}' -> '{fallback.name}'")
                formula_paths.append(fallback)
            else:
                print(f"[layout_filler] skipping missing formula asset: {fname}")
    print("formula_paths",formula_paths)
    return _dedupe_paths(image_paths + table_paths + formula_paths)

 
 


  
import re
from pathlib import Path
import difflib

 
def _extract_idx(val):
   
    if isinstance(val, int):
        return val
    m = re.findall(r"\d+", str(val))
    return int(m[-1]) if m else None

def _nums_from_files(files):
    """['image_7.png', 'table_1.png'] -> {7, 1}"""
    out = set()
    for f in files:
        m = re.findall(r"\d+", str(f))
        if m:
            out.add(int(m[-1]))
    return out

def _best_match(data, sec_title: str, sub_title: str, min_ratio: float = 0.55):
    best_sec, best_sec_score = None, 0.0
    for sec in data.get("sections", []):
        s = difflib.SequenceMatcher(None, sec.get("title","").lower(), (sec_title or "").lower()).ratio()
        if s > best_sec_score:
            best_sec, best_sec_score = sec, s
    if not best_sec or best_sec_score < min_ratio:
        return None, None

    best_sub, best_sub_score = None, 0.0
    for sub in best_sec.get("subsections", []):
        s = difflib.SequenceMatcher(None, sub.get("title","").lower(), (sub_title or "").lower()).ratio()
        if s > best_sub_score:
            best_sub, best_sub_score = sub, s
    if not best_sub or best_sub_score < min_ratio:
        return None, None
    return best_sec, best_sub
 
def _collect_reasons_for_kind(sec, sub, files, figs_data, kind: str):
    
    _, sub_d = _best_match(figs_data, sec, sub)
    if not sub_d:
        return ""

    want = _nums_from_files(files)
    if not want:
        return ""

    #   N -> value_idx
    pairs = []  # [(N, value_idx)]
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith(kind):  # 'image' or 'table'
            continue
        N = _extract_idx(k) 
        v_idx = _extract_idx(v)   
        if N is None or v_idx is None:
            continue
        if v_idx in want:
            pairs.append((N, v_idx))

    reasons = []
     
    only_one_asset = (len(pairs) == 1)

    for N, _v_idx in pairs:
         
        candidate_keys = [
            f"reason{N}",
        ]
        if kind == "table":
            candidate_keys.append(f"reasonT{N}")
        else:
            candidate_keys.extend([f"reasonI{N}", f"reasonImg{N}"])
         
        if only_one_asset:
            candidate_keys.append("reason")

         
        found = None
        for rk in candidate_keys:
            if rk in sub_d and isinstance(sub_d[rk], str) and sub_d[rk].strip():
                found = sub_d[rk].strip()
                break

         
        if not found and "reason" in sub_d and isinstance(sub_d["reason"], str) and sub_d["reason"].strip():
            found = sub_d["reason"].strip()

        if found:
            reasons.append(found)

    return "\n".join(reasons)

def get_image_reasons(sec, sub, image_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, image_files, figs_data, kind="image")

def get_table_reasons(sec, sub, table_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, table_files, figs_data, kind="table")

def get_formula_reasons(sec, sub, formula_files, formula_data) -> str:
     
    _, sub_d = _best_match(formula_data, sec, sub)
    if not sub_d:
        return ""
 
    pairs = []
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith("formula"):
            continue
        N = _extract_idx(k)   
        if N is None:
            continue
        rkey = f"reason{N}"
        rtxt = sub_d.get(rkey)
        if isinstance(rtxt, str) and rtxt.strip():
            pairs.append((N, rtxt.strip()))

    
    try:
        pairs.sort(key=lambda x: int(x[0]))
    except Exception: 
        pass

   
    if not pairs and isinstance(sub_d.get("reason"), str) and sub_d["reason"].strip():
        if formula_files and len(formula_files) == 1:
            return sub_d["reason"].strip()
 
    if formula_files and len(pairs) > len(formula_files):
        pairs = pairs[:len(formula_files)]
 
    out, seen = [], set()
    for _, r in pairs:
        if r not in seen:
            out.append(r)
            seen.add(r)

    return "\n".join(out)

def resolve_formula_mode1_path(fname: str, args) -> Path:
    """
    Extract formula index i from fname (like "formula_4.png"),
    and generate the path:
    <{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-i.png
    """ 

    stem = Path(fname).stem   
    match = re.search(r'(\d+)(?!.*\d)', stem) 
    # match = re.search(r'(?i)\bformula(?:[_\-\s]*)?(\d+)\b', stem)
    if not match:
        return None
    
    i = match.group(1)
    path_str = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-{i}.png"
    return Path(path_str)

import json
from pathlib import Path
from typing import Dict, Any, List

def _is_T1_textonly(s: Dict[str, Any]) -> bool:
     
    return (
        s.get("template_id") == "T1_TextOnly"
        and not s.get("images")
        and not s.get("tables")
        and not s.get("formulas")
    )


def _bullet_text_length(slide: Dict[str, Any]) -> int:
    total = 0
    for bullet in slide.get("bullets") or []:
        total += len((bullet.get("text") or "").strip())
        total += sum(len(str(sub).strip()) for sub in (bullet.get("sub") or []))
    return total


def _should_merge_T1_pair(left: Dict[str, Any], right: Dict[str, Any]) -> bool:
    left_bullets = left.get("bullets") or []
    right_bullets = right.get("bullets") or []
    left_subsection = (left.get("subsection") or "").strip()
    right_subsection = (right.get("subsection") or "").strip()
    left_len = _bullet_text_length(left)
    right_len = _bullet_text_length(right)

    if not left_subsection or not right_subsection:
        return False
    if len(left_bullets) < 2 or len(right_bullets) < 2:
        return False
    if left_len < 120 or right_len < 120:
        return False
    return True

def pair_T1_to_T19(plan_path: str, write_back: bool = True) -> int:
 
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    out: List[Dict[str, Any]] = []

    i, n, made = 0, len(slides), 0
    while i < n:
        cur = slides[i]
        if (
            i + 1 < n
            and _is_T1_textonly(cur)
            and _is_T1_textonly(slides[i + 1])
            and cur.get("section") == slides[i + 1].get("section")
            and _should_merge_T1_pair(cur, slides[i + 1])
        ):
            left, right = cur, slides[i + 1]
            out.append({
                "section": cur.get("section"),
                "template_id": "T19_2Text",
                "columns": [
                    {
                        "subsection": left.get("subsection", "") or "",
                        "bullets": left.get("bullets", []) or []
                    },
                    {
                        "subsection": right.get("subsection", "") or "",
                        "bullets": right.get("bullets", []) or []
                    }
                ]
            })
            made += 1
            i += 2
            continue
 
        out.append(cur)
        i += 1

    plan["slides"] = out
    if write_back:
        p.write_text(json.dumps(plan, ensure_ascii=False, indent=4), encoding="utf-8")
    return made

def validate_no_consecutive_T1(plan_path: str) -> List[int]:
     
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    bad_idxs: List[int] = []
    for i in range(len(slides) - 1):
        a, b = slides[i], slides[i + 1]
        if _is_T1_textonly(a) and _is_T1_textonly(b) and a.get("section") == b.get("section"):
            bad_idxs.append(i)
    return bad_idxs
from pptx.util import Pt

def _clear_text_frame(tf):
    if tf.paragraphs:
        tf.paragraphs[0].text = ""
         
        while len(tf.paragraphs) > 1:
            p = tf._element.p_lst[-1]
            p.getparent().remove(p)
    else:
        tf.clear()

def _fill_bullets(
    shape_or_tf,
    bullets,
    lvl0_size=24,
    lvl1_size=24,
    min_size_pt: float = 10.0,
    max_size_pt: float | None = None,
    target_min_fill: float = 0.46,
    target_max_fill: float = 0.82,
):
    if getattr(shape_or_tf, "has_text_frame", False):
        shape = shape_or_tf
        tf = shape.text_frame
    else:
        shape = None
        tf = shape_or_tf
    if max_size_pt is None:
        max_size_pt = min(max(lvl0_size, lvl1_size) + 4.0, 28.0)
    _clear_text_frame(tf)
    _compact_text_frame(tf)
    for b in _normalize_bullets(bullets):
        p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = (b.get("text") or "").strip()
        p.level = 0
        _set_paragraph_font_size(p, lvl0_size)
        
        for s in (b.get("sub") or []):
            sp = tf.add_paragraph()
            sp.text = str(s).strip()
            sp.level = 1
            _set_paragraph_font_size(sp, lvl1_size)
    _enable_shrink_to_fit(tf)
    if shape is not None:
        _shrink_text_frame_to_fit(shape, tf, min_size_pt=min_size_pt)
        _grow_text_frame_to_fill(
            shape,
            tf,
            min_size_pt=min_size_pt,
            max_size_pt=max_size_pt,
            target_min_fill=target_min_fill,
            target_max_fill=target_max_fill,
        )


def _stabilize_slide_text_shapes(slide) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if not tf or not any(_paragraph_text(p).strip() for p in tf.paragraphs):
            continue
        name = (getattr(shape, "name", "") or "").lower()
        min_size_pt = 12.0 if "title" in name else 10.0
        _shrink_text_frame_to_fit(shape, tf, min_size_pt=min_size_pt)
    _expand_sparse_slide_pictures(slide)
from pptx.enum.shapes import PP_PLACEHOLDER

def _get_placeholder(slide, name): 
    for shp in slide.shapes:
        if getattr(shp, "name", "").strip().lower() == name.strip().lower():
            return shp
    return None

def _ph_text_n(slide, n:int): 
    targets = {f"text placeholder {n}", f"文本占位符 {n}"}
    for shp in slide.shapes:
        nm = getattr(shp, "name", "").strip().lower()
        if nm in targets:
            return shp 
        if ("placeholder" in nm or "占位符" in nm) and nm.endswith(f" {n}"):
            return shp
    return None
def _ph_by_idx(slide, idx:int):
    for ph in getattr(slide, "placeholders", []):
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def fill_T19_2Text(slide, slide_info, section_no_text):
     
    print("slide_info:")
    print(slide_info)
    part_ph = (
        _get_placeholder(slide, "Part")
        or _get_placeholder(slide, "Text Placeholder 2")
        or _ph_text_n(slide, 2)
    )
   
    title_bar = _get_placeholder(slide, "Text Placeholder 1") or _ph_text_n(slide, 1)
    title_bar = _ph_by_idx(slide, 1)  or _ph_text_n(slide, 1)

     
    section_title = slide_info.get("section") 
    
    title_bar.text = section_title
    tf = title_bar.text_frame
    tf.clear()                        
    tf.paragraphs[0].text = section_title
    _enable_shrink_to_fit(tf, max_size_pt=24)
    print("[AFTER] title_bar text =", repr(tf.text))

    lt = (
        _get_placeholder(slide, "Left Title")
        or _get_placeholder(slide, "Text Placeholder 3")
        or _ph_text_n(slide, 3)
    )
    lb = (
        _get_placeholder(slide, "Left Body")
        or _get_placeholder(slide, "Text Placeholder 4")
        or _ph_text_n(slide, 4)
    )
    rt = (
        _get_placeholder(slide, "Right Title")
        or _get_placeholder(slide, "Text Placeholder 5")
        or _ph_text_n(slide, 5)
    )
    rb = (
        _get_placeholder(slide, "Right Body")
        or _get_placeholder(slide, "Text Placeholder 6")
        or _ph_text_n(slide, 6)
    )

    
    cols = slide_info.get("columns") or []
    left  = cols[0] if len(cols) > 0 else {}
    right = cols[1] if len(cols) > 1 else {}
 
    if part_ph is not None and getattr(part_ph, "has_text_frame", False):
        _set_shape_text_with_fit(part_ph, f"{section_no_text}", max_size_pt=28)

    if title_bar is not None and getattr(title_bar, "has_text_frame", False):
        title_txt = slide_info.get("section", "") or slide_info.get("title", "")
        _set_shape_text_with_fit(title_bar, title_txt, max_size_pt=_header_font_size(title_txt))

    if lt is not None and getattr(lt, "has_text_frame", False):
        _set_shape_text_with_fit(lt, left.get("subsection", "") or left.get("title", "") or "", max_size_pt=22)

    if rt is not None and getattr(rt, "has_text_frame", False):
        _set_shape_text_with_fit(rt, right.get("subsection", "") or right.get("title", "") or "", max_size_pt=22)

    if lb is not None and getattr(lb, "has_text_frame", False):
        _fill_bullets(lb, left.get("bullets"))

    if rb is not None and getattr(rb, "has_text_frame", False):
        _fill_bullets(rb, right.get("bullets"))

     
    missing = []
    for name, ph in [
        ("Part(2)", part_ph), ("Left Title(3)", lt), ("Left Body(4)", lb),
        ("Right Title(5)", rt), ("Right Body(6)", rb),
    ]:
        if ph is None:
            missing.append(name)
    if missing:
        print(f"[WARN] T19_2Text 模板缺少占位：{', '.join(missing)}")



def delete_slide(prs: Presentation, slide_index: int) -> None:
    """
    Delete slide by index in python-pptx (works by removing sldId and dropping rel).
    """
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    sldId_elems = list(sldIdLst)

    if slide_index < 0 or slide_index >= len(sldId_elems):
        raise IndexError(f"slide_index out of range: {slide_index}, total={len(sldId_elems)}")

    sldId = sldId_elems[slide_index]
    rId = sldId.rId
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)

  
# debug_list_placeholders(slide)

def generate_pptx_from_plan( 
    args,
    template: Path | int 
):
 
     
    figs_json_path  =  f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_figures.json"
    formula_json_path = f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_formula_match.json"
    paper_outline_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_raw_content.json' 
    with open(paper_outline_json, "r", encoding="utf-8") as f: outline_json  = json.load(f)
    with open(figs_json_path, encoding="utf-8") as f: figs_data   = json.load(f)
    with open(formula_json_path, encoding="utf-8") as f: formula_data   = json.load(f)
    
    variant_suffix = getattr(
        args,
        "output_variant_suffix",
        "_personalized" if getattr(args, "use_author_preferences", False) else "_baseline",
    )
    plan_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_slide_plan{variant_suffix}.json'
     
    made = pair_T1_to_T19(plan_json)   
    print(f"[plan] T1->T19 pairs made: {made}")
      
    plan: Dict = json.loads(Path(plan_json).read_text(encoding="utf-8"))


    title    = outline_json["metadata"]["title"]
    subtitle = outline_json["metadata"]["author"]
  
    template_path = f"utils/slides_template/slides{template}_template.pptx"
    prs = Presentation(template_path)
      
    theme_color = extract_theme_color_from_title(prs)
    print("Theme color:", theme_color)
     
    # test 
    print("Available slide layouts in template:")
    for layout in prs.slide_layouts:
        print("-", layout.name)
 
    # ---------- cover ----------
    cover_layout = prs.slide_layouts.get_by_name("Title Slide")
    cover = prs.slides.add_slide(cover_layout)

    _set_shape_text_with_fit(
        _placeholder_by_name(cover, "Title 1"),
        title,
        max_size_pt=34,
        min_size_pt=18,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _set_shape_text_with_fit(
        _placeholder_by_name(cover, "Subtitle 2"),
        subtitle,
        max_size_pt=22,
        min_size_pt=14,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    # ---------- Contents ----------
    outline_layout = prs.slide_layouts.get_by_name("Mulu")
    outline = prs.slides.add_slide(outline_layout)

    contents_ph = _placeholder_by_name(outline, "Text Placeholder 1")
    tf = contents_ph.text_frame
    tf.clear()
    _enable_shrink_to_fit(tf)
    seen = set() 
    unique_sections = []
 
    for slide in plan["slides"]:
        sec = slide["section"]
        if sec not in seen:
            seen.add(sec)
            unique_sections.append(sec)
 
    contents_font_size = _contents_font_size(len(unique_sections))
    _populate_contents_frame(tf, unique_sections, contents_font_size)
            
    # ---------- Body ----------
    current_section  = None
    section_counter  = 0
    for slide_info in plan["slides"]: 
        # ---------- Content ----------
        if slide_info["section"] != current_section:
            current_section = slide_info["section"]
            section_counter += 1 
            section_layout = prs.slide_layouts.get_by_name("dan_mulu")
            sec_slide = prs.slides.add_slide(section_layout)
            for shape in sec_slide.shapes:
                print(f"Shape: {shape.name}")
            _set_shape_text_with_fit(
                _placeholder_by_name(sec_slide, "Text Placeholder 2"),
                f"PART {section_counter:02d}",
                max_size_pt=38,
                min_size_pt=18,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )
            _set_shape_text_with_fit(
                _placeholder_by_name(sec_slide, "Title 1"),
                current_section,
                max_size_pt=36,
                min_size_pt=18,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )
  
        visuals = resolve_visual_paths(slide_info, args)
        layout_names = [layout.name for layout in prs.slide_layouts]
        template_id = str(slide_info["template_id"])
        template_id = _fallback_text_only_template(template_id, slide_info)
        if _should_use_top_visual_layout(slide_info, visuals):
            template_id = "T4_ImageTop"
        template_id = normalize_template_id(template_id, layout_names)
        template_id = _fallback_missing_template_id(template_id, slide_info, visuals, layout_names)
        slide_info["template_id"] = template_id
        layout = prs.slide_layouts.get_by_name(template_id)
          
        if layout is None:
            raise ValueError(f" Template layout '{template_id}' not found in template.")
        slide = prs.slides.add_slide(layout)
    
        if template_id == "T19_2Text":
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.is_placeholder:
                    if shape.has_text_frame:
                        print(f" Name: {shape.name}")
                        print(f"  Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}")
                        print(f"  Text: '{shape.text_frame.text.strip()}'")
    
            fill_T19_2Text(slide, slide_info, section_no_text=f"{section_counter:02d}")
            continue

        part_ph, title_ph, body_ph = find_text_placeholders(slide)
  
        _set_shape_text_with_fit(part_ph, f"{section_counter:02d}", max_size_pt=28)
        _set_shape_text_with_fit(title_ph, slide_info["subsection"], max_size_pt=_header_font_size(slide_info["subsection"]))
  
   
        # bullets + sub-bullets
        if body_ph:
            tf = body_ph.text_frame
            _clear_text_frame(tf)
            tf.word_wrap = True
            lvl0_size, lvl1_size = _body_font_sizes(slide_info)
            bullets = _normalize_bullets(slide_info.get("bullets") or [])
            # if tf.paragraphs:
            #     tf.paragraphs[0].text = ""   
            # else:
            #     tf.clear() 
            for bullet in bullets:
                p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
                p.text, p.level = bullet.get("text", ""), 0
                _set_paragraph_font_size(p, lvl0_size)
                for sub in bullet.get("sub", []):
                    sp = tf.add_paragraph()
                    sp.text, sp.level = sub, 1
                    _set_paragraph_font_size(sp, lvl1_size)
            _enable_shrink_to_fit(tf)
            _shrink_text_frame_to_fit(body_ph, tf, min_size_pt=10.0)
            has_visuals = bool(slide_info.get("images") or slide_info.get("tables") or slide_info.get("formulas"))
            _grow_text_frame_to_fill(
                body_ph,
                tf,
                min_size_pt=10.0,
                max_size_pt=30.0 if not has_visuals else 26.0,
                target_min_fill=0.58 if not has_visuals else 0.46,
                target_max_fill=0.86 if not has_visuals else 0.82,
            )
         
         
        insert_visuals_auto(slide, visuals)
        _stabilize_slide_text_shapes(slide)
 

        # ---------- note ----------
        notes_chunks = [] 
        txt = get_content(slide_info["section"], slide_info["subsection"], outline_json)
        if txt: notes_chunks.append(txt)
        if slide_info.get("images"):
            img_r = get_image_reasons(slide_info["section"], slide_info["subsection"],
                                      slide_info["images"], figs_data)
            notes_chunks.append(img_r)
        if slide_info.get("tables"):
            tb_r = get_table_reasons(slide_info["section"], slide_info["subsection"],
                                     slide_info["tables"], figs_data)
            notes_chunks.append(tb_r)
        if slide_info.get("formulas"):
            fm_r = get_formula_reasons(
                slide_info["section"],
                slide_info["subsection"],
                slide_info["formulas"],
                formula_data  ,
                )
                
            notes_chunks.append(fm_r)
                    
        if notes_chunks:
            nframe = slide.notes_slide.notes_text_frame
            if nframe.text and not nframe.text.endswith("\n"):
                nframe.text += "\n"
            nframe.text += "\n\n".join(notes_chunks)

 
    thanks_layout = prs.slide_layouts.get_by_name("Last_page")
    thanks = prs.slides.add_slide(thanks_layout)
    title_ph = _placeholder_by_name(thanks, "Title 1")
    _set_shape_text_with_fit(
        title_ph,
        "THANKS!",
        max_size_pt=42,
        min_size_pt=20,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
 
    run = title_ph.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    for slide in prs.slides:
        _stabilize_slide_text_shapes(slide)
    output_pptx = (
        f'contents/{args.paper_name}/'
        f'{args.model_name_t}_{args.model_name_v}_output_slides{variant_suffix}.pptx'
    ) 
    prs.save(str(output_pptx))
    delete_slide(prs, 0)    
    prs.save(str(output_pptx))
    
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = Path(prefix) / args.paper_name
    target_name = f"{args.paper_name}-with-image-refs_artifacts"
    
    artifacts_dir = base_dir / target_name
    
    if not artifacts_dir.exists():
        hits = list(base_dir.glob(f"**/{target_name}"))
        if hits: 
            artifacts_dir = max(hits, key=lambda p: len(p.parts))

    img_candidates = sorted(artifacts_dir.glob("image_*.png"))
    theme_imgs = [str(p) for p in img_candidates[:2]]


    if len(theme_imgs) < 1:
        print(f"[warn] No image_*.png found under: {artifacts_dir}. Skip theming.")
    else:
 
        theme_hex, base_hex = pick_theme_color(
            images=theme_imgs,
            prefer_dark=True,
            min_v=0.10,
            max_v=0.99,
            return_base_hex=True,
        )
        print("[theme] base_hex :", base_hex)
        print("[theme] theme_hex:", theme_hex)
        print("[theme] imgs:", theme_imgs)

        
        output_pptx_path = Path(output_pptx)
        themed_pptx_path = output_pptx_path.with_name(output_pptx_path.stem + "_themed" + output_pptx_path.suffix)
        print("colored path")
       
        set_one_theme_color(
            pptx_in=str(output_pptx_path),
            pptx_out=str(themed_pptx_path),
            color_hex=theme_hex,
            target_key="dk2",    
        )

        print(f"[ok] themed pptx saved: {themed_pptx_path}")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Fill PPTX from slides plan.")
    parser.add_argument(
        "--plan",  
        default="SlideGen/contents/STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems/<4o_4o>_slide_plan.json",  
        help="slides_plan.json"
        )
    parser.add_argument(
        "--paper_name",  
        default="STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems"  
        )
    parser.add_argument(
        "--template",  
        type=int,
        default=3,
        help="Template number, e.g. 3 for slides3_template.pptx"
    )
    parser.add_argument(
        "--out", 
        default="output.pptx" 
        )
    parser.add_argument(
        "--model_name_t", 
        default="4o" 
        )
    parser.add_argument(
        "--model_name_v", 
        default="4o" 
        )
        
    args = parser.parse_args()

    layout_info = scan_layout_placeholders("SlideGen/utils/slides_template/slides3_template.pptx")

    for layout_name, shapes in layout_info.items():
        print(f"\n Layout: {layout_name}")
        for s in shapes:
            print(f"  - {s['name']} ({s['type']})")

 

    generate_pptx_from_plan(args.plan, args.template, args.out)
     
    print(f" Saved to {args.out}")
