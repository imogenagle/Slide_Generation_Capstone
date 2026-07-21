"""Layout binder + user-template renderer.

Sits between SlideGen's refiner and rendering. When --template_path is set,
SlideGen runs its full agent pipeline as normal; this module then:
  1. Renders the user template's slides to images.
  2. Has a vision LLM describe each user-template slide.
  3. Calls an LLM that picks, for each SlideGen plan slide, which user-template
     slide to clone (the binder agent).
  4. Renders the output PPTX by duplicating chosen user-template slides and
     filling them with SlideGen's content (titles, bullets, figures).

Content (titles, bullets, figure choices) comes straight from SlideGen's plan;
this module only decides which template slide hosts each piece and renders it.
"""

from __future__ import annotations

import base64
import copy
import hashlib
import json
import math
import os
import re
import subprocess
import sys
import time
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Optional

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from slidegen_openai_utils import build_openai_client, resolve_direct_model_name


SOFFICE_BIN = "/opt/homebrew/bin/soffice"


DESCRIBE_SYSTEM_PROMPT = """You are analyzing a single slide from a PowerPoint template.
Look at the slide image and return a JSON object describing its structure.

Output schema (return ONLY the JSON, no prose):
{
  "slide_type": "title" | "section_divider" | "content_bullets" | "image_focus" | "two_column" | "comparison" | "data_table" | "thanks" | "blank" | "instruction" | "other",
  "description": "1-2 sentences describing what this slide visually shows",
  "is_meta": true | false,
  "text_regions": [
    {"role": "title|subtitle|body|caption|footer|other", "approx_position": "top|center|bottom|left|right"}
  ],
  "image_regions": [
    {"approx_position": "top|center|bottom|left|right", "size": "small|medium|large"}
  ],
  "suitable_for": ["short tag describing what content this slide layout could host"]
}

Set "is_meta": true when the slide is a template-instruction page (e.g. "delete this note",
"how to use this template", "click here in PowerPoint"). Otherwise set "is_meta": false.
"""


BIND_SYSTEM_PROMPT = """You are a layout binder. SlideGen has produced a deck plan for a research paper. The user has uploaded a PowerPoint template. Your job is to pick, for each slide in SlideGen's plan, which slide of the user template to clone as its visual host.

You will be given:
  1. SlideGen's plan slides (each has section, subsection, template_id hint, bullet count, and counts for images / tables / formulas — tables and formulas are rendered as images on the slide, so a slide may carry multiple visuals).
  2. Deck-level metadata: the deck title and the subtitle (author string), used only to fill the cover slide. There is NO per-slide metadata for contents / section dividers / thanks — choose those indices purely from the user-template descriptions in (3).
  3. A JSON list describing each slide of the user template (slide_type, image_regions, is_meta, and `has_picture_placeholder` — a programmatically-verified flag for whether the slide has a real picture placeholder a figure can be dropped into).

SlideGen's `template_id` hints encode the structural need:
  - T1_TextOnly                       -> bullets only, no visuals
  - T2_ImageRight                     -> bullets left, visual right
  - T3_ImageLeft                      -> visual left, bullets right
  - T4_ImageTop                       -> visual top, bullets below
  - T16/T17/T18 (formula variants)    -> include a formula image + optionally an image
  - T19_2Text                         -> two text columns
  - Other T* codes                    -> similar text+visual variants

Rules:
  - NEVER pick a user-template slide where `is_meta` is true.
  - For the cover, prefer slide_type "title".
  - For the contents/agenda, prefer "content_bullets" or "two_column". DO NOT return null unless the template literally has zero non-meta slides with body text — even an imperfect match is better than skipping the agenda entirely.
  - For section dividers, prefer "section_divider"; fall back to "title". Return null ONLY if the template has nothing close.
  - For thanks, prefer "thanks"; fall back to "title". Return null ONLY if the template has nothing close.
  - For body slides:
      - HARD REQUIREMENT: if the SlideGen slide has any visuals (n_images + n_tables + n_formulas > 0), you MUST pick a user-template slide whose `has_picture_placeholder` is true. A figure dropped onto a slide with no picture placeholder lands on top of the bullet text. Only ignore this if NO non-meta template slide has `has_picture_placeholder` true.
      - Among `has_picture_placeholder` slides, prefer slide_type "image_focus", "two_column", or "comparison".
      - If it has multiple visuals (>= 2) -> strongly prefer "image_focus" or "comparison" (more room).
      - If it's text-only (total_visuals == 0) -> prefer "content_bullets" or "two_column"; `has_picture_placeholder` does not matter.
      - For data_table content (n_tables > 0 and n_bullets is low) -> prefer "data_table" if available, else an "image_focus" slide with a picture placeholder.
      - Match figure-position hint (right/left/top) to image_regions when possible.
  - It is OK to reuse the same user-template slide for many body slides (multiple bullet-style slides can all clone the same content_bullets layout).
  - Pick the SAME section_divider layout for every section transition.

Output schema (return ONLY the JSON, no prose):
{
  "cover_slide_index": <int>,
  "contents_slide_index": <int or null>,
  "section_divider_slide_index": <int or null>,
  "thanks_slide_index": <int or null>,
  "body_assignments": [<int>, <int>, ...]
}

`body_assignments` must have length equal to the number of SlideGen plan slides, in the same order. Each entry is the user-template slide index to clone for that plan slide.
"""


# ---------- Logging ---------- #


@contextmanager
def _stage(num: int, label: str):
    print(f"[layout_binder] (stage {num}) {label}...")
    import time
    t0 = time.time()
    try:
        yield
    finally:
        print(f"[layout_binder]   (stage {num} took {time.time() - t0:.1f}s)")


# ---------- Stage 1: render template to images ---------- #


def render_template_to_images(template_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    if not Path(SOFFICE_BIN).exists():
        raise FileNotFoundError(
            f"LibreOffice not found at {SOFFICE_BIN}. "
            "Install with `brew install --cask libreoffice` or update SOFFICE_BIN in layout_binder.py."
        )
    try:
        subprocess.run(
            [SOFFICE_BIN, "--headless", "--convert-to", "pdf",
             "--outdir", str(output_dir), str(template_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
    except subprocess.CalledProcessError as e:
        stderr = e.stderr.decode("utf-8", errors="replace") if e.stderr else "(no stderr)"
        raise RuntimeError(
            f"LibreOffice failed to convert template to PDF (exit {e.returncode}). "
            f"Stderr: {stderr.strip()[:500]}"
        ) from e
    except subprocess.TimeoutExpired as e:
        raise RuntimeError(
            f"LibreOffice timed out after {e.timeout}s converting {template_path.name}."
        ) from e

    pdf_path = output_dir / f"{template_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"LibreOffice exited cleanly but produced no PDF at {pdf_path}."
        )

    images = convert_from_path(pdf_path, dpi=120)
    image_paths: list[Path] = []
    for i, img in enumerate(images):
        path = output_dir / f"slide_{i:03d}.png"
        img.save(path, "PNG")
        image_paths.append(path)
    return image_paths


# ---------- Stage 2: describe template slides ---------- #


def _parse_json_from_response(text: str) -> dict[str, Any]:
    text = text.strip()
    if text.startswith("```"):
        text = text.split("```", 2)[1]
        if text.startswith("json"):
            text = text[4:]
        text = text.rsplit("```", 1)[0].strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        try:
            from json_repair import repair_json
            return json.loads(repair_json(text))
        except Exception:
            return {}


def describe_template_slides(image_paths: list[Path]) -> list[dict[str, Any]]:
    client = build_openai_client()
    deployment = resolve_direct_model_name(os.environ.get("AZURE_DEPLOYMENT_NAME", ""))
    descriptions: list[dict[str, Any]] = []
    for i, img_path in enumerate(image_paths):
        with img_path.open("rb") as f:
            img_b64 = base64.b64encode(f.read()).decode()
        parsed = None
        for attempt in (1, 2, 3):
            try:
                response = client.chat.completions.create(
                    model=deployment,
                    messages=[
                        {"role": "system", "content": DESCRIBE_SYSTEM_PROMPT},
                        {"role": "user", "content": [
                            {"type": "text", "text": f"This is slide {i+1} of the template."},
                            {"type": "image_url",
                             "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                        ]},
                    ],
                    max_completion_tokens=800,
                )
                raw = response.choices[0].message.content or ""
                parsed = _parse_json_from_response(raw)
                break
            except Exception as e:
                # A transient API error on one slide must not discard every
                # already-paid vision call — retry, then degrade to defaults.
                print(f"[layout_binder]   WARN: describe failed for slide {i+1} "
                      f"(attempt {attempt}): {e}")
                if attempt < 3:
                    time.sleep(2 * attempt)
        if parsed is None:
            parsed = {}
        parsed.setdefault("slide_type", "other")
        parsed.setdefault("description", "")
        parsed.setdefault("is_meta", False)
        parsed.setdefault("text_regions", [])
        parsed.setdefault("image_regions", [])
        parsed.setdefault("suitable_for", [])
        parsed["slide_index"] = i
        descriptions.append(parsed)
        meta_tag = " [META]" if parsed.get("is_meta") else ""
        print(f"[layout_binder]   slide {i+1}: {parsed.get('slide_type', '?')}{meta_tag} - {parsed.get('description', '')[:70]}")
    return descriptions


def scan_picture_placeholders(template_path: Path) -> dict[int, int]:
    """Return {slide_index: count of usable picture-anchor spots} for the template.

    Programmatic ground truth — far more reliable than the vision LLM's
    slide_type guess for deciding which slides can host a figure. A slide has
    an anchor spot when it carries a picture placeholder (empty or filled with
    a sample photo) or a substantial sample picture a figure can replace — the
    same predicate the renderer's anchor finder uses, so the binder's
    has_picture_placeholder flag matches what rendering can actually do.
    """
    prs = Presentation(str(template_path))
    return {i: len(_anchor_boxes_for_slide(slide, prs.slide_width, prs.slide_height))
            for i, slide in enumerate(prs.slides)}


# ---------- Stage 3: bind SlideGen plan slides to user-template slides ---------- #


def _summarize_plan_for_binder(slidegen_plan: dict) -> list[dict]:
    """Produce a compact per-slide summary for the binder LLM."""
    summary = []
    for i, s in enumerate(slidegen_plan.get("slides", [])):
        n_images = len(s.get("images") or [])
        n_tables = len(s.get("tables") or [])
        n_formulas = len(s.get("formulas") or [])
        # Two-column (T19) slides store their bullets under columns[].bullets,
        # not top-level bullets — count both so the binder doesn't see 0.
        n_bullets = len(_flatten_slidegen_bullets(s.get("bullets") or []))
        for col in (s.get("columns") or []):
            if isinstance(col, dict):
                n_bullets += len(_flatten_slidegen_bullets(col.get("bullets") or []))
        summary.append({
            "plan_index": i,
            "section": s.get("section", ""),
            "subsection": s.get("subsection", ""),
            "template_id": s.get("template_id", ""),
            "n_bullets": n_bullets,
            "n_images": n_images,
            "n_tables": n_tables,
            "n_formulas": n_formulas,
            "total_visuals": n_images + n_tables + n_formulas,
        })
    return summary


def _coerce_index(v) -> Optional[int]:
    """Best-effort conversion of an LLM-returned slide index to an int."""
    if isinstance(v, bool):
        return None
    if isinstance(v, int):
        return v
    if isinstance(v, float) and v.is_integer():
        return int(v)
    if isinstance(v, str) and v.strip().lstrip("-").isdigit():
        return int(v.strip())
    return None


def _first_non_meta(descriptions: list[dict]) -> int:
    for d in descriptions:
        if not d.get("is_meta"):
            return d.get("slide_index", 0)
    return 0


def _plan_slide_visuals(plan_slide: dict) -> int:
    return (len(plan_slide.get("images") or [])
            + len(plan_slide.get("tables") or [])
            + len(plan_slide.get("formulas") or []))


def _hint_position(template_id: str) -> Optional[str]:
    t = (template_id or "").upper()
    if "RIGHT" in t:
        return "right"
    if "LEFT" in t:
        return "left"
    if "TOP" in t:
        return "top"
    return None


def _choose_pic_slide(plan_slide: dict, pic_slides: list[int],
                      desc_by_idx: dict, used_pic: list[int]) -> int:
    """Pick the picture-anchor template slide that best matches this plan
    slide's needs (figure-position hint, visual count), instead of funneling
    every figure slide onto one identical layout."""
    n_visuals = _plan_slide_visuals(plan_slide)
    hint = _hint_position(plan_slide.get("template_id", ""))

    def score(idx: int) -> int:
        d = desc_by_idx.get(idx) or {}
        s = 0
        if n_visuals >= 2:
            if (d.get("n_picture_placeholders") or 0) >= 2:
                s += 4
            if d.get("slide_type") in ("image_focus", "comparison"):
                s += 2
        if hint:
            for r in (d.get("image_regions") or []):
                if isinstance(r, dict) and r.get("approx_position") == hint:
                    s += 3
                    break
        if d.get("slide_type") in ("image_focus", "two_column", "comparison"):
            s += 1
        if idx in used_pic:
            s += 1  # mild consistency bonus
        return s

    return max(pic_slides, key=score)


def bind_user_template_to_plan(
    slidegen_plan: dict,
    template_descriptions: list[dict],
    deck_meta: dict,
) -> dict:
    """LLM call: pick a user-template slide for each SlideGen plan slide.

    Returns a dict with cover/contents/section_divider/thanks indices plus a
    list of body_assignments parallel to slidegen_plan['slides'].
    """
    client = build_openai_client()
    deployment = resolve_direct_model_name(os.environ.get("AZURE_DEPLOYMENT_NAME", ""))

    plan_summary = _summarize_plan_for_binder(slidegen_plan)

    user_payload = {
        "deck_meta": deck_meta,
        "plan_slides": plan_summary,
        "template_slides": template_descriptions,
    }

    response = client.chat.completions.create(
        model=deployment,
        messages=[
            {"role": "system", "content": BIND_SYSTEM_PROMPT},
            {"role": "user", "content": json.dumps(user_payload)},
        ],
        max_completion_tokens=2000,
    )
    raw = response.choices[0].message.content or ""
    binding = _parse_json_from_response(raw)
    if not binding or "body_assignments" not in binding:
        raise RuntimeError(
            f"Binder returned malformed response. Head: {raw[:300]}"
        )

    n_plan_slides = len(slidegen_plan.get("slides", []))
    assignments = list(binding.get("body_assignments") or [])
    if len(assignments) != n_plan_slides:
        print(
            f"[layout_binder]   WARN: binder returned {len(assignments)} assignments "
            f"for {n_plan_slides} plan slides; padding/truncating."
        )
        if len(assignments) < n_plan_slides:
            # Pad with None — the per-slide repair below picks a sensible host
            # (a title layout's clone usually has no body placeholder).
            assignments.extend([None] * (n_plan_slides - len(assignments)))
        else:
            assignments = assignments[:n_plan_slides]

    # Sanitize everything the LLM returned: coerce to int (it sometimes emits
    # digit strings), require a described slide index, and never allow a meta
    # (template-instruction) slide to be cloned.
    desc_by_idx = {d.get("slide_index"): d for d in template_descriptions}

    def _sanitize(v) -> Optional[int]:
        idx = _coerce_index(v)
        if idx is None or idx not in desc_by_idx:
            return None
        if desc_by_idx[idx].get("is_meta"):
            return None
        return idx

    for key in ("cover_slide_index", "contents_slide_index",
                "section_divider_slide_index", "thanks_slide_index"):
        cleaned = _sanitize(binding.get(key))
        if cleaned != binding.get(key):
            print(f"[layout_binder]   WARN: sanitized {key}: "
                  f"{binding.get(key)!r} -> {cleaned!r}")
        binding[key] = cleaned
    assignments = [_sanitize(a) for a in assignments]

    # Enforce the picture-placeholder rule the prompt asks for: any plan slide
    # with visuals MUST land on a template slide that has a picture anchor,
    # otherwise the figure renders on top of the bullet text. We don't trust the
    # LLM to honor this — verify against the programmatic flag and repair.
    # Invalid/None entries get a per-slide fallback host so no plan slide is
    # ever silently dropped at render time.
    pic_slides = sorted(idx for idx, d in desc_by_idx.items()
                        if idx is not None
                        and d.get("has_picture_placeholder")
                        and not d.get("is_meta"))
    text_fallback = _pick_fallback_slide(
        template_descriptions, ["content_bullets", "two_column", "other"])
    if text_fallback is None:
        text_fallback = _first_non_meta(template_descriptions)

    plan_slides = slidegen_plan.get("slides", [])
    used_pic = [a for a in assignments if a in pic_slides]
    repaired_fig = repaired_none = repaired_text = 0
    for i, plan_slide in enumerate(plan_slides):
        n_visuals = _plan_slide_visuals(plan_slide)
        n_bullets = len(_flatten_slidegen_bullets(plan_slide.get("bullets") or []))
        for col in (plan_slide.get("columns") or []):
            if isinstance(col, dict):
                n_bullets += len(_flatten_slidegen_bullets(col.get("bullets") or []))
        if n_visuals > 0 and pic_slides:
            if assignments[i] not in pic_slides:
                assignments[i] = _choose_pic_slide(
                    plan_slide, pic_slides, desc_by_idx, used_pic)
                used_pic.append(assignments[i])
                repaired_fig += 1
        elif assignments[i] is None:
            assignments[i] = text_fallback
            repaired_none += 1
        elif (n_visuals == 0 and n_bullets >= 3 and text_fallback is not None
              and (desc_by_idx.get(assignments[i]) or {}).get("slide_type") == "image_focus"):
            # Enforce the prompt's text-only rule programmatically: bullet-heavy
            # slides hosted on full-bleed photo/quote layouts render as text
            # overlapping a stock photo.
            assignments[i] = text_fallback
            repaired_text += 1
    if repaired_fig:
        print(f"[layout_binder]   repaired {repaired_fig} figure slide(s) "
              f"onto picture-anchor layouts")
    if repaired_none:
        print(f"[layout_binder]   repaired {repaired_none} invalid assignment(s) "
              f"onto text layout {text_fallback}")
    if repaired_text:
        print(f"[layout_binder]   moved {repaired_text} text-only slide(s) "
              f"off image-focus layouts")
    if not pic_slides:
        print("[layout_binder]   WARN: template has no picture-anchor "
              "slides; figures will use beside-text fallback placement")

    binding["body_assignments"] = assignments
    return binding


# ---------- Stage 4: render output PPTX ---------- #


# Placeholder type IDs (PP_PLACEHOLDER):
#   1 TITLE, 2 BODY, 3 CENTER_TITLE, 4 SUBTITLE, 7 OBJECT
_TITLE_TYPES = {1, 3}
_SUBTITLE_TYPES = {4}
_SKIP_TYPES = {13, 14, 15, 16}  # slide number / header / footer / date


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _copy_referenced_rels(src_part, new_part, element) -> None:
    """Re-create on `new_part` every relationship that XML copied from
    `src_part` references (r:embed / r:link / r:id / ...), rewriting the
    attributes to the new rIds. Without this, deep-copied picture XML keeps
    rIds that only exist on the source slide part and every template image
    renders as a broken/missing picture.
    """
    for el in element.iter():
        for attr, old_rid in list(el.attrib.items()):
            if not attr.startswith(_R_NS) or not old_rid:
                continue
            try:
                rel = src_part.rels[old_rid]
            except KeyError:
                continue  # already dangling in the source; nothing to map
            if rel.is_external:
                new_rid = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rid = new_part.relate_to(rel.target_part, rel.reltype)
            el.set(attr, new_rid)


def _remove_shape_element(slide, element) -> None:
    """Remove a shape's XML from `slide` and drop any relationships that no
    other XML on the slide still references (keeps the package free of orphan
    image/chart parts)."""
    rids = {v for el in element.iter()
            for k, v in el.attrib.items() if k.startswith(_R_NS)}
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)
    if not rids:
        return
    remaining = {v for el in slide.element.iter()
                 for k, v in el.attrib.items() if k.startswith(_R_NS)}
    for rid in rids - remaining:
        try:
            slide.part.drop_rel(rid)
        except (KeyError, ValueError):
            pass


def _duplicate_slide(prs, src_index: int):
    """Append a deep copy of slide `src_index` to the end of `prs`.

    Copies shape XML, the per-slide background / color-map override, and
    re-creates every relationship the copied XML references so no r:embed /
    r:id dangles on the new slide part.
    """
    source = prs.slides[src_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    new_spTree = new_slide.shapes._spTree
    for shape in list(new_slide.shapes):
        new_spTree.remove(shape.element)

    for shape in source.shapes:
        el = copy.deepcopy(shape.element)
        new_spTree.append(el)
        _copy_referenced_rels(source.part, new_slide.part, el)

    # Background and color-map override live OUTSIDE spTree on the source slide.
    src_cSld = source.element.find(f"{_P_NS}cSld")
    new_cSld = new_slide.element.find(f"{_P_NS}cSld")
    if src_cSld is not None and new_cSld is not None:
        src_bg = src_cSld.find(f"{_P_NS}bg")
        if src_bg is not None and new_cSld.find(f"{_P_NS}bg") is None:
            bg = copy.deepcopy(src_bg)
            new_cSld.insert(0, bg)
            _copy_referenced_rels(source.part, new_slide.part, bg)
    src_clrMapOvr = source.element.find(f"{_P_NS}clrMapOvr")
    if src_clrMapOvr is not None:
        old = new_slide.element.find(f"{_P_NS}clrMapOvr")
        if old is not None:
            new_slide.element.remove(old)
        new_slide.element.append(copy.deepcopy(src_clrMapOvr))

    return new_slide


def _delete_slides_at(prs, indices: list[int]) -> None:
    """Remove slides AND their presentation-part relationships so the deleted
    template slides (and any media only they reference) are not serialized
    into the saved package."""
    sldIdLst = prs.slides._sldIdLst
    sldId_elems = list(sldIdLst)
    for i in sorted(indices, reverse=True):
        sldId = sldId_elems[i]
        rId = sldId.get(f"{_R_NS}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except (KeyError, ValueError):
                pass
        sldIdLst.remove(sldId)


# Substrings (lowercased) that mark a shape as leftover template sample content.
_SAMPLE_TEXT_MARKERS = (
    "bullet level", "lorem ipsum", "helvetica", "additional content here",
    "delete this slide", "image caption", "caption text here", "caption or",
    "description text here", "presentation title here", "divider slide text",
    "slide title", "firstname lastname", "cnet@uchicago", "url.uchicago",
    "head level", "text box ma", "text level", "click here", "placeholder text",
    "entity name", "school name", "presentation title", "your logo",
    "add a footer", "sample text", "insert your", "type your",
    "back to agenda", "your subtitle",
)


def _collect_shape_text(shape) -> str:
    """Gather all text from a shape, recursing into groups. Lowercased."""
    parts: list[str] = []
    if shape.has_text_frame and shape.text_frame.text.strip():
        parts.append(shape.text_frame.text)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            parts.append(_collect_shape_text(child))
    return " ".join(parts).lower()


def _strip_sample_shapes(slide) -> int:
    """Remove leftover template sample content from a cloned slide.

    Targets non-placeholder shapes only (placeholders are handled by
    _fill_slide_text). Removes sample charts and tables outright, and any
    shape whose text matches a known sample marker (e.g. callout groups with
    'Bullet level 1' / 'Helvetica Bold 24pt'). Brand decoration with no
    sample text is left intact. Returns the number stripped.
    """
    to_remove = []
    for shape in slide.shapes:
        # Sample charts / tables — template examples, never real content here.
        if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
            to_remove.append(shape)
            continue
        if shape.is_placeholder:
            continue  # _fill_slide_text owns placeholders
        text = _collect_shape_text(shape)
        if text and any(marker in text for marker in _SAMPLE_TEXT_MARKERS):
            to_remove.append(shape)
    for shape in to_remove:
        _remove_shape_element(slide, shape._element)
    return len(to_remove)


def _scrub_sample_text_in_shape(shape) -> int:
    """Clear marker-matched sample text runs inside a non-placeholder shape or
    group WITHOUT removing the shape. Used on slide layouts/masters, where
    brand decoration (logos in the same group) must survive. Returns the
    number of text frames cleared."""
    cleared = 0
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            cleared += _scrub_sample_text_in_shape(child)
        return cleared
    if shape.is_placeholder:
        return 0  # layout placeholder prompt text never renders on slides
    if shape.has_text_frame:
        text = shape.text_frame.text.strip().lower()
        if text and any(m in text for m in _SAMPLE_TEXT_MARKERS):
            _clear_text_frame(shape.text_frame)
            cleared += 1
    return cleared


def _scrub_layouts_and_masters(prs) -> int:
    """Clear leftover sample text ('ENTITY NAME', 'presentation title') from
    the layouts/masters the rendered slides inherit from — per-slide stripping
    can't reach text that lives there."""
    layouts = {}
    for slide in prs.slides:
        layout = slide.slide_layout
        layouts[id(layout.element)] = layout  # SlideLayout is unhashable
    masters = {}
    for layout in layouts.values():
        master = layout.slide_master
        masters[id(master.element)] = master
    scrubbed = 0
    for host in list(layouts.values()) + list(masters.values()):
        for shape in host.shapes:
            scrubbed += _scrub_sample_text_in_shape(shape)
    return scrubbed


def _enable_shrink_to_fit(tf) -> None:
    if tf is None:
        return
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_text_frame(tf, text: str) -> None:
    if not tf.paragraphs:
        tf.text = text
        _enable_shrink_to_fit(tf)
        return
    p0 = tf.paragraphs[0]
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    # Hard line-breaks from multi-line template sample titles must go too, or
    # the filled title ships with a dangling blank line.
    for br in list(p0._p.findall(f"{_A_NS}br")):
        p0._p.remove(br)
    if p0.runs:
        p0.runs[0].text = text
    else:
        p0.text = text
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    _enable_shrink_to_fit(tf)


_LATEX_UNICODE = {
    r"\\leq": "≤", r"\\geq": "≥", r"\\le": "≤", r"\\ge": "≥",
    r"\\neq": "≠", r"\\approx": "≈", r"\\sim": "∼",
    r"\\to": "→", r"\\rightarrow": "→", r"\\Rightarrow": "⇒",
    r"\\leftarrow": "←", r"\\Leftarrow": "⇐", r"\\leftrightarrow": "↔",
    r"\\times": "×", r"\\cdot": "·", r"\\pm": "±",
    r"\\alpha": "α", r"\\beta": "β", r"\\gamma": "γ", r"\\delta": "δ",
    r"\\epsilon": "ε", r"\\varepsilon": "ε", r"\\zeta": "ζ", r"\\eta": "η",
    r"\\theta": "θ", r"\\lambda": "λ", r"\\mu": "μ", r"\\nu": "ν",
    r"\\pi": "π", r"\\rho": "ρ", r"\\sigma": "σ", r"\\tau": "τ",
    r"\\phi": "φ", r"\\chi": "χ", r"\\psi": "ψ", r"\\omega": "ω",
    r"\\Gamma": "Γ", r"\\Delta": "Δ", r"\\Theta": "Θ", r"\\Lambda": "Λ",
    r"\\Pi": "Π", r"\\Sigma": "Σ", r"\\Phi": "Φ", r"\\Omega": "Ω",
    r"\\ell": "ℓ", r"\\infty": "∞", r"\\partial": "∂", r"\\nabla": "∇",
    r"\\in": "∈", r"\\subset": "⊂", r"\\subseteq": "⊆", r"\\cup": "∪", r"\\cap": "∩",
    r"\\forall": "∀", r"\\exists": "∃",
}
_MATH_DELIM_RE = re.compile(r"\$+([^$]*?)\$+")
_MATH_PAREN_RE = re.compile(r"\\\((.*?)\\\)", re.S)
_MATH_BRACK_RE = re.compile(r"\\\[(.*?)\\\]", re.S)
_SUBSCRIPTS = {
    "0": "₀", "1": "₁", "2": "₂", "3": "₃", "4": "₄", "5": "₅",
    "6": "₆", "7": "₇", "8": "₈", "9": "₉", "a": "ₐ", "e": "ₑ",
    "h": "ₕ", "i": "ᵢ", "j": "ⱼ", "k": "ₖ", "l": "ₗ", "m": "ₘ",
    "n": "ₙ", "o": "ₒ", "p": "ₚ", "r": "ᵣ", "s": "ₛ", "t": "ₜ",
    "u": "ᵤ", "v": "ᵥ", "x": "ₓ",
}
_SUPERSCRIPTS = {
    "0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴", "5": "⁵",
    "6": "⁶", "7": "⁷", "8": "⁸", "9": "⁹", "n": "ⁿ",
}


def _clean_math(text: str) -> str:
    if not text or ("$" not in text and "\\" not in text):
        return text
    # Strip \(...\) and \[...\] delimiters first, keeping the inner math.
    text = _MATH_PAREN_RE.sub(lambda m: m.group(1), text)
    text = _MATH_BRACK_RE.sub(lambda m: m.group(1), text)
    for cmd, repl in _LATEX_UNICODE.items():
        # (?![a-zA-Z]) not \b: `\ell_p` has no word boundary before `_`, and
        # \b would leave the command for the catch-all delete below.
        text = re.sub(cmd + r"(?![a-zA-Z])", repl, text)
    text = _MATH_DELIM_RE.sub(lambda m: m.group(1), text)
    text = re.sub(
        r"\\(?:mathbf|mathit|mathrm|mathcal|mathbb|boldsymbol|operatorname|text|emph|bf|it)\{([^}]*)\}",
        r"\1", text)
    text = re.sub(r"\\[a-zA-Z]+", "", text)
    text = text.replace("{", "").replace("}", "")
    # Render simple sub/superscripts with unicode ("ℓ_p" -> "ℓₚ", "x^2" -> "x²").
    # Only after a non-alphanumeric char, so snake_case prose is left alone.
    text = re.sub(r"(?<![A-Za-z0-9])\s*_([0-9a-z])(?![A-Za-z0-9])",
                  lambda m: _SUBSCRIPTS.get(m.group(1), "_" + m.group(1)), text)
    text = re.sub(r"\^([0-9n])(?![A-Za-z0-9])",
                  lambda m: _SUPERSCRIPTS.get(m.group(1), "^" + m.group(1)), text)
    text = re.sub(r"  +", " ", text)
    return text.strip()


def _flatten_slidegen_bullets(raw: list) -> list[tuple[str, int]]:
    """Convert SlideGen's nested {text, sub: [...]} bullets to (text, level).

    SlideGen emits `bullets: [{"text": "...", "sub": ["...", "..."]}, ...]`.
    Levels become real PowerPoint outline levels (not two-space fake indents)
    so nesting survives rendering.
    """
    out: list[tuple[str, int]] = []

    def _walk(items, level: int) -> None:
        for item in items or []:
            if isinstance(item, dict):
                text = item.get("text") or ""
                if text:
                    out.append((text, level))
                _walk(item.get("sub") or [], min(level + 1, 4))
            elif isinstance(item, str) and item:
                out.append((item, level))

    _walk(raw, 0)
    return out


def _as_bullet_pairs(raw) -> list[tuple[str, int]]:
    """Normalize a bullets list (strings or (text, level) pairs) to pairs."""
    out: list[tuple[str, int]] = []
    for b in raw or []:
        if isinstance(b, str):
            out.append((b, 0))
        elif isinstance(b, (tuple, list)) and b:
            lvl = int(b[1]) if len(b) > 1 else 0
            out.append((str(b[0]), max(0, min(lvl, 4))))
    return out


def _set_bullets(tf, bullets) -> None:
    """Fill a text frame with bullets. Accepts strings or (text, level) pairs;
    levels are applied as real outline levels."""
    pairs = _as_bullet_pairs(bullets)
    if not pairs:
        return
    _set_text_frame(tf, pairs[0][0])
    # Snapshot paragraph 0's formatting BEFORE touching levels so every bullet
    # starts from the template's sample-paragraph styling.
    template_para_xml = copy.deepcopy(tf.paragraphs[0]._p)
    txBody = tf.paragraphs[0]._p.getparent()
    if pairs[0][1]:
        tf.paragraphs[0].level = pairs[0][1]

    for text, level in pairs[1:]:
        new_p = copy.deepcopy(template_para_xml)
        for r in list(new_p.findall(f"{_A_NS}r")):
            new_p.remove(r)
        for br in list(new_p.findall(f"{_A_NS}br")):
            new_p.remove(br)
        txBody.append(new_p)
        last_para = tf.paragraphs[-1]
        run = last_para.add_run()
        run.text = text
        if level:
            # Drop explicit level-0 indents copied from the sample paragraph so
            # the layout's own per-level list style takes over.
            pPr = last_para._p.find(f"{_A_NS}pPr")
            if pPr is not None:
                for attr in ("marL", "indent"):
                    if attr in pPr.attrib:
                        del pPr.attrib[attr]
            last_para.level = level


def _clear_text_frame(tf) -> None:
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)


def _estimate_fits(lines: list[tuple[str, int]], box_w_in: float, box_h_in: float,
                   size_pt: float, spacing_pt: float = 0.0) -> bool:
    """Rough check: does this text fit the box at size_pt? Uses average glyph
    width ≈ 0.52*size (+ letter-spacing) and line height ≈ 1.24*size. Also
    requires the longest WORD to fit one line — otherwise renderers break
    words mid-word ('Certifie / d')."""
    usable_w = max(box_w_in - 0.2, 0.5)
    usable_h = max(box_h_in - 0.15, 0.3)
    char_w_in = (size_pt * 0.52 + spacing_pt) / 72.0
    # Wide/bold display faces run past the average estimate — check the longest
    # word against a fatter glyph width and a safety margin.
    word_char_w_in = (size_pt * 0.60 + spacing_pt) / 72.0
    total_lines = 0
    max_word = 0
    for text, level in lines:
        eff_w = max(usable_w - 0.28 * level, 0.5)
        chars_per_line = max(int(eff_w / char_w_in), 8)
        total_lines += max(1, math.ceil(len(text) / chars_per_line))
        for word in text.split():
            max_word = max(max_word, len(word))
    if max_word * word_char_w_in > usable_w * 0.94:
        return False
    return total_lines * size_pt * 1.24 / 72.0 <= usable_h


def _letter_spacing_pt(tf) -> float:
    """Explicit letter-spacing (a:rPr spc, 1/100 pt) of the first styled run —
    heavily tracked template titles need it in the width estimate."""
    for p in tf.paragraphs:
        for r in p.runs:
            rPr = r._r.find(f"{_A_NS}rPr")
            if rPr is not None and rPr.get("spc"):
                try:
                    return int(rPr.get("spc")) / 100.0
                except ValueError:
                    pass
    return 0.0


def _inherited_defrpr(slide, ph) -> tuple[Optional[float], float]:
    """(font size pt, letter-spacing pt) a placeholder inherits from its layout
    placeholder's list style, falling back to the master's txStyles. Slide-level
    runs usually carry NO explicit size — the template's 60pt/15pt-tracked
    divider titles live here, invisible to run-level inspection."""
    ptype = ph.placeholder_format.type
    idx = ph.placeholder_format.idx
    candidates = []
    try:
        layout = slide.slide_layout
        for lph in layout.placeholders:
            if lph.placeholder_format.idx == idx or lph.placeholder_format.type == ptype:
                el = lph._element.find(
                    f".//{_A_NS}lstStyle/{_A_NS}lvl1pPr/{_A_NS}defRPr")
                if el is not None:
                    candidates.append(el)
                break
        master = layout.slide_master
        style_tag = "titleStyle" if ptype in (1, 3) else "bodyStyle"
        el = master.element.find(
            f".//{_P_NS}{style_tag}/{_A_NS}lvl1pPr/{_A_NS}defRPr")
        if el is not None:
            candidates.append(el)
    except Exception:
        pass
    size = spc = None
    for el in candidates:
        if size is None and el.get("sz"):
            size = int(el.get("sz")) / 100.0
        if spc is None and el.get("spc"):
            spc = int(el.get("spc")) / 100.0
    return size, (spc or 0.0)


def _apply_fitted_font(tf, box_w, box_h, *, base_default: float, min_pt: float = 11.0,
                       slide=None, ph=None) -> None:
    """Shrink run font sizes so the text actually fits the box.

    python-pptx's TEXT_TO_FIT_SHAPE writes <a:normAutofit/> with no computed
    fontScale, and PowerPoint only recalculates autofit when a box is edited —
    so without explicit sizes, overfull text displays at full size. This sets
    explicit sizes when (and only when) the estimate says the template size
    would overflow. The base size and letter-spacing are resolved through the
    layout/master inheritance chain (templates rarely set them on the run),
    and tracking is scaled down with the font so shrunk text stays dense.
    """
    if not box_w or not box_h:
        return
    lines: list[tuple[str, int]] = []
    base = None
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs)
        lines.append((text, p.level or 0))
        if base is None:
            for r in p.runs:
                if r.font.size is not None:
                    base = r.font.size.pt
                    break
    if not any(t for t, _ in lines):
        return
    spacing = _letter_spacing_pt(tf)
    if (base is None or not spacing) and slide is not None and ph is not None:
        inh_size, inh_spc = _inherited_defrpr(slide, ph)
        base = base or inh_size
        spacing = spacing or inh_spc
    base = base or base_default
    # Letter-spacing shrinks proportionally with the font during fitting.
    size = base
    while size > min_pt and not _estimate_fits(
            lines, box_w / 914400.0, box_h / 914400.0, size,
            spacing * (size / base)):
        size -= 1.0
    if size >= base:
        return  # template size already fits; keep the template's styling
    new_spc = spacing * (size / base)
    for p in tf.paragraphs:
        target = max(size - (2.0 if (p.level or 0) else 0.0), min_pt)
        for r in p.runs:
            r.font.size = Pt(target)
            if spacing:
                rPr = r._r.get_or_add_rPr()
                rPr.set("spc", str(int(new_spc * 100)))


def _add_fallback_body_textbox(slide, prs, has_visuals: bool, min_top: int = 0):
    """Last-resort host for bullets when a cloned slide has no body/object
    placeholder — content must never be silently dropped."""
    sw, sh = prs.slide_width, prs.slide_height
    title_bottom = min_top
    for ph in slide.placeholders:
        if ph.placeholder_format.type in _TITLE_TYPES:
            title_bottom = max(title_bottom, (ph.top or 0) + (ph.height or 0))
    left = int(sw * 0.06)
    top = max(int(sh * 0.22), title_bottom + int(sh * 0.03))
    width = int(sw * (0.44 if has_visuals else 0.88))
    height = max(int(sh * 0.92) - top, int(sh * 0.2))
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def _max_run_size_pt(shape) -> float:
    best = 0.0
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.find(f"{_A_NS}rPr")
            if rPr is not None and rPr.get("sz"):
                best = max(best, int(rPr.get("sz")) / 100.0)
    return best


def _clear_decor_text_overlapping(slide, box_shape, keep_els: set) -> int:
    """Clear template decor text fragments (e.g. a leftover 'sleep.' word or
    giant sample numerals) that sit on top of a text box we just filled —
    otherwise they overprint the real title. Only non-placeholder text shapes
    substantially overlapping the filled box are touched."""
    bl, bt = box_shape.left or 0, box_shape.top or 0
    bw, bh = box_shape.width or 0, box_shape.height or 0
    if not bw or not bh:
        return 0
    cleared = 0
    for sh in slide.shapes:
        if sh.is_placeholder or not sh.has_text_frame:
            continue
        if sh._element in keep_els:
            continue
        if not sh.text_frame.text.strip():
            continue
        sl, st = sh.left or 0, sh.top or 0
        sw, shh = sh.width or 0, sh.height or 0
        if not sw or not shh:
            continue
        ix = max(0, min(sl + sw, bl + bw) - max(sl, bl))
        iy = max(0, min(st + shh, bt + bh) - max(st, bt))
        if ix * iy >= 0.3 * sw * shh:
            _clear_text_frame(sh.text_frame)
            cleared += 1
    return cleared


def _clear_literal_footer_text(slide) -> int:
    """Clear literal sample strings ('presentation title', '20XX') from
    header/footer/date placeholders while preserving field-based content
    (slide numbers, auto dates). Returns count of placeholders cleared."""
    cleared = 0
    for ph in slide.placeholders:
        if ph.placeholder_format.type not in (14, 15, 16):  # header/footer/date
            continue
        if not ph.has_text_frame:
            continue
        changed = False
        for p in ph.text_frame.paragraphs:
            if p._p.find(f"{_A_NS}fld") is not None:
                continue  # field-driven paragraph: keep
            for r in list(p.runs):
                if r.text.strip():
                    r._r.getparent().remove(r._r)
                    changed = True
        if changed:
            cleared += 1
    return cleared


def _fill_slide_text(slide, content: dict, prs=None) -> None:
    """Fill placeholders BY TYPE (not by name) so it works with any template.

    Bullets go to the LARGEST body/object placeholder — templates often carry
    small caption/contact placeholders that must not swallow the body text.
    Two-column content fills two placeholders when the layout offers them, and
    a fallback text box is added when a slide has no text placeholder at all.
    """
    title_text = _clean_math(content.get("title", "") or "")
    subtitle_text = _clean_math(content.get("subtitle", "") or "")
    bullets = [(_clean_math(t), lvl) for t, lvl in _as_bullet_pairs(content.get("bullets") or [])]
    bullets = [(t, lvl) for t, lvl in bullets if t]
    columns = content.get("columns") or []

    used_phs: set[int] = set()

    # Pseudo-placeholders: designer templates (Canva/SlidesCarnival exports)
    # often build slides from plain text boxes / auto shapes with NO
    # placeholders at all. When a placeholder search fails, we fill the
    # template's own large text shapes in place — keeping the designer's
    # typography — instead of dropping content or overlaying a generic box.
    used_pseudo: set = set()  # lxml ELEMENTS of filled shapes (identity-stable while referenced)
    pseudo_title_bottom = 0

    def _pseudo_candidates():
        if prs is None:
            return []
        slide_area = prs.slide_width * prs.slide_height
        out = []
        for sh in slide.shapes:
            if sh.is_placeholder or not sh.has_text_frame:
                continue
            if not sh.text_frame.text.strip():
                continue
            w, h = sh.width or 0, sh.height or 0
            if not w or not h or w * h < 0.012 * slide_area:
                continue  # logos, page numbers, nav links
            out.append(sh)
        return out

    # 1. Title
    title_filled = False
    if title_text:
        for ph in slide.placeholders:
            if ph.placeholder_format.type in _TITLE_TYPES and ph.has_text_frame:
                _set_text_frame(ph.text_frame, title_text)
                _apply_fitted_font(ph.text_frame, ph.width, ph.height,
                                   base_default=32.0, min_pt=16.0,
                                   slide=slide, ph=ph)
                used_phs.add(ph.placeholder_format.idx)
                _clear_decor_text_overlapping(slide, ph, used_pseudo)
                title_filled = True
                break
        if not title_filled and prs is not None:
            # No title placeholder: replace the most title-like text shape
            # (largest font, upper part of the slide) in place.
            cands = [sh for sh in _pseudo_candidates()
                     if (sh.top or 0) < prs.slide_height * 0.6
                     and (sh.width or 0) >= prs.slide_width * 0.22]
            if cands:
                box = max(cands, key=lambda s: (_max_run_size_pt(s),
                                                (s.width or 0) * (s.height or 0)))
                _set_text_frame(box.text_frame, title_text)
                box.text_frame.word_wrap = True
                _apply_fitted_font(box.text_frame, box.width, box.height,
                                   base_default=32.0, min_pt=16.0)
                used_pseudo.add(box._element)
                _clear_decor_text_overlapping(slide, box, used_pseudo)
                # Sample headings are often SPLIT across sibling shapes
                # ("Cats love" / "to sleep."). Clear any other candidate whose
                # vertical band overlaps the chosen title box — leftover
                # fragments there would stack on the real title.
                band_gap = int(prs.slide_height * 0.05)
                b_top = (box.top or 0) - band_gap
                b_bot = (box.top or 0) + (box.height or 0) + band_gap
                for other in cands:
                    if other._element is box._element:
                        continue
                    o_top = other.top or 0
                    o_bot = o_top + (other.height or 0)
                    if o_bot >= b_top and o_top <= b_bot:
                        _clear_text_frame(other.text_frame)
                pseudo_title_bottom = (box.top or 0) + (box.height or 0)
                title_filled = True
                print("[layout_binder]   no title placeholder; replaced template title text shape")

    # 2. Bullets (before subtitle, so they win the single-content slot on
    # Section Header layouts). Candidates sorted largest-first.
    body_phs = [ph for ph in slide.placeholders
                if ph.placeholder_format.type in (7, 2)
                and ph.has_text_frame
                and ph.placeholder_format.idx not in used_phs]
    body_phs.sort(key=lambda p: (p.width or 0) * (p.height or 0), reverse=True)

    if content.get("distribute") and len(bullets) > 1 and len(body_phs) >= 2:
        # Agenda-style layouts carry one small placeholder per item (e.g. seven
        # arrow boxes). Distribute one bullet per placeholder in column-major
        # reading order instead of cramming everything into the first box.
        targets = sorted(body_phs, key=lambda p: (round((p.left or 0) / 914400), p.top or 0))
        k = len(targets)
        sizes = [len(bullets) // k + (1 if i < len(bullets) % k else 0) for i in range(k)]
        start = 0
        for ph, sz in zip(targets, sizes):
            chunk = bullets[start:start + sz]
            start += sz
            if not chunk:
                continue
            _set_bullets(ph.text_frame, chunk)
            _apply_fitted_font(ph.text_frame, ph.width, ph.height,
                               base_default=18.0, slide=slide, ph=ph)
            used_phs.add(ph.placeholder_format.idx)
        bullets = []
    elif columns and len(columns) >= 2 and len(body_phs) >= 2:
        # Two-column plan content onto a multi-placeholder layout: one column
        # per placeholder, left-to-right.
        targets = sorted(body_phs[:len(columns)], key=lambda p: (p.left or 0))
        for col, ph in zip(columns, targets):
            col_bullets: list[tuple[str, int]] = []
            header = _clean_math(col.get("header") or "")
            if header:
                col_bullets.append((header, 0))
            for t, lvl in _as_bullet_pairs(col.get("bullets") or []):
                t = _clean_math(t)
                if t:
                    col_bullets.append((t, min(lvl + 1, 4) if header else lvl))
            if not col_bullets:
                continue
            _set_bullets(ph.text_frame, col_bullets)
            _apply_fitted_font(ph.text_frame, ph.width, ph.height,
                               base_default=18.0, slide=slide, ph=ph)
            used_phs.add(ph.placeholder_format.idx)
        bullets = []  # consumed by the per-column fill
    if bullets:
        def _too_small(shape) -> bool:
            # A tiny caption/stat box can't host real content at readable size
            # — cramming yields microscopic text. Bail to the fallback box.
            if prs is None:
                return False
            w, h = shape.width or 0, shape.height or 0
            if (w * h) / (prs.slide_width * prs.slide_height) >= 0.15:
                return False
            return not _estimate_fits(list(bullets), w / 914400.0, h / 914400.0, 11.0)

        target_ph = body_phs[0] if body_phs else None
        if target_ph is not None and _too_small(target_ph):
            print("[layout_binder]   body placeholder too small for content; using fallback")
            target_ph = None
        if target_ph is not None:
            _set_bullets(target_ph.text_frame, bullets)
            _apply_fitted_font(target_ph.text_frame, target_ph.width, target_ph.height,
                               base_default=18.0, slide=slide, ph=target_ph)
            used_phs.add(target_ph.placeholder_format.idx)
        elif prs is not None:
            # No body placeholder. Prefer the template's own largest text
            # shape (its sample-copy box) so bullets inherit the design; fall
            # back to a generic text box only when the slide has none.
            cands = [sh for sh in _pseudo_candidates()
                     if sh._element not in used_pseudo]
            cands = [sh for sh in cands if not _too_small(sh)]
            if cands:
                box = max(cands, key=lambda s: (s.width or 0) * (s.height or 0))
                _set_bullets(box.text_frame, bullets)
                box.text_frame.word_wrap = True
                _apply_fitted_font(box.text_frame, box.width, box.height,
                                   base_default=18.0)
                used_pseudo.add(box._element)
                print("[layout_binder]   no body placeholder; filled template text shape with bullets")
            else:
                box = _add_fallback_body_textbox(slide, prs, bool(content.get("n_visuals")),
                                                 min_top=pseudo_title_bottom)
                _set_bullets(box.text_frame, bullets)
                _apply_fitted_font(box.text_frame, box.width, box.height,
                                   base_default=16.0)
                print("[layout_binder]   no body placeholder; added fallback text box for bullets")

    # 3. Subtitle
    if subtitle_text:
        chosen = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type in _SUBTITLE_TYPES and ph.has_text_frame:
                chosen = ph
                break
        if chosen is None:
            for ph in slide.placeholders:
                if (ph.placeholder_format.type == 2
                    and ph.has_text_frame
                    and ph.placeholder_format.idx not in used_phs):
                    chosen = ph
                    break
        if chosen is not None:
            _set_text_frame(chosen.text_frame, subtitle_text)
            _apply_fitted_font(chosen.text_frame, chosen.width, chosen.height,
                               base_default=20.0, min_pt=12.0,
                               slide=slide, ph=chosen)
            used_phs.add(chosen.placeholder_format.idx)
            _clear_decor_text_overlapping(slide, chosen, used_pseudo)
        elif prs is not None:
            # No subtitle/body placeholder: reuse the template's next text
            # shape (e.g. the cover tagline box) for the author line.
            cands = [sh for sh in _pseudo_candidates()
                     if sh._element not in used_pseudo]
            if cands:
                box = max(cands, key=lambda s: (s.width or 0) * (s.height or 0))
                _set_text_frame(box.text_frame, subtitle_text)
                box.text_frame.word_wrap = True
                _apply_fitted_font(box.text_frame, box.width, box.height,
                                   base_default=20.0, min_pt=12.0)
                used_pseudo.add(box._element)
                print("[layout_binder]   no subtitle placeholder; replaced template text shape")

    # 4. Clear remaining unused content placeholders so sample text doesn't ship.
    for ph in slide.placeholders:
        ptype = ph.placeholder_format.type
        if ptype in _SKIP_TYPES:
            continue
        if not ph.has_text_frame:
            continue
        if ph.placeholder_format.idx in used_phs:
            continue
        if ptype in (2, 4, 7):
            _clear_text_frame(ph.text_frame)

    # 5. Clear literal sample strings in footer/date placeholders (kept out of
    # _SKIP_TYPES clearing above so real slide-number fields survive).
    _clear_literal_footer_text(slide)

    # 6. On placeholder-poor slides where we filled the template's own text
    # shapes, clear the REMAINING substantial sample-copy shapes ("We believe
    # software should support people...") so template prose doesn't ship next
    # to paper content. Short labels/logos are untouched.
    if used_pseudo:
        cleared = 0
        for sh in _pseudo_candidates():
            if sh._element in used_pseudo:
                continue
            if len(sh.text_frame.text.split()) >= 3:
                _clear_text_frame(sh.text_frame)
                cleared += 1
        if cleared:
            print(f"[layout_binder]   cleared {cleared} leftover sample text shape(s)")


def _fit_image_in_box(slide, image_path: Path, box_left: int, box_top: int,
                      box_w: int, box_h: int) -> None:
    """Place an image into a box, aspect-preserved, centered."""
    from PIL import Image
    with Image.open(str(image_path)) as img:
        fw, fh = img.size
    figure_aspect = fw / fh if fh else 1.0
    box_aspect = box_w / box_h if box_h else 1.0

    if figure_aspect > box_aspect:
        pic = slide.shapes.add_picture(str(image_path), box_left, box_top, width=box_w)
        pic.top = box_top + (box_h - pic.height) // 2
    else:
        pic = slide.shapes.add_picture(str(image_path), box_left, box_top, height=box_h)
        pic.left = box_left + (box_w - pic.width) // 2


def _is_picture_like(shape) -> bool:
    """True for real pictures AND picture placeholders already filled with an
    image (those report shape_type PLACEHOLDER but their element is <p:pic>)."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    return shape._element.tag == f"{_P_NS}pic"


def _anchor_boxes_for_slide(slide, slide_w: int, slide_h: int) -> list:
    """Designed picture spots on a slide, in reading order.

    1. Picture placeholders (empty, or filled with a sample photo).
    2. Otherwise, substantial sample pictures (3%-75% of slide area — smaller
       is a logo, larger is full-bleed decoration; neither is a content spot).
    Anchors with unresolvable/degenerate geometry are skipped rather than
    coalesced to a (0,0,0,0) box. Returns [(box, shape)], box=(l, t, w, h).
    """
    slide_area = slide_w * slide_h
    anchors = []
    for ph in slide.placeholders:
        if ph.placeholder_format.type != 18:
            continue
        w, h = ph.width, ph.height
        if not w or not h:
            continue
        anchors.append(((ph.left or 0, ph.top or 0, w, h), ph))
    if not anchors:
        pics = []
        for s in slide.shapes:
            if s.is_placeholder or not _is_picture_like(s):
                continue
            w, h = s.width or 0, s.height or 0
            area = w * h
            if area < 0.03 * slide_area or area > 0.75 * slide_area:
                continue
            pics.append(((s.left or 0, s.top or 0, w, h), s))
        pics.sort(key=lambda bs: bs[0][2] * bs[0][3], reverse=True)
        anchors = pics[:3]
    anchors.sort(key=lambda bs: (bs[0][1], bs[0][0]))
    return anchors


def _strip_pictures_in_region(slide, box_left: int, box_top: int,
                              box_w: int, box_h: int, overlap_threshold: float = 0.3) -> int:
    """Remove picture-like shapes (including filled picture placeholders)
    whose bounding box significantly overlaps the target visual region.
    Decorative pictures elsewhere (e.g. corner logos) survive. Returns the
    number stripped.

    We strip BEFORE inserting our new visual so two copies of the same area
    don't end up stacked on top of each other.
    """
    to_remove = []
    for shape in slide.shapes:
        if not _is_picture_like(shape):
            continue
        sl, st = shape.left or 0, shape.top or 0
        sw, sh = shape.width or 0, shape.height or 0
        ix1, iy1 = max(sl, box_left), max(st, box_top)
        ix2, iy2 = min(sl + sw, box_left + box_w), min(st + sh, box_top + box_h)
        if ix2 > ix1 and iy2 > iy1:
            inter_area = (ix2 - ix1) * (iy2 - iy1)
            shape_area = sw * sh
            if shape_area > 0 and inter_area / shape_area >= overlap_threshold:
                to_remove.append(shape)
    for shape in to_remove:
        _remove_shape_element(slide, shape._element)
    return len(to_remove)


def _best_grid(n: int, box_w: int, box_h: int, aspects: list[float]) -> tuple[int, float]:
    """Choose a column count (1..n) that maximizes the smallest fitted visual
    area for n visuals packed into the box. Returns (cols, min_fitted_area)."""
    best_cols, best_min_area = 1, -1.0
    for cols in range(1, max(n, 1) + 1):
        rows = math.ceil(n / cols)
        cw, ch = box_w / cols, box_h / rows
        min_area = None
        for a in aspects or [1.33]:
            w = min(cw, ch * a) if a else cw
            h = w / a if a else ch
            area = w * h
            min_area = area if min_area is None else min(min_area, area)
        if min_area is not None and min_area > best_min_area:
            best_cols, best_min_area = cols, min_area
    return best_cols, best_min_area


def _place_visuals_grid(slide, paths: list[Path], box, aspects: list[float]) -> int:
    """Pack visuals into `box` on the best rows×cols grid (aspect-aware)."""
    box_left, box_top, box_w, box_h = box
    n = len(paths)
    if n == 0 or box_w <= 0 or box_h <= 0:
        return 0
    cols, _ = _best_grid(n, box_w, box_h, aspects)
    rows = math.ceil(n / cols)
    gap_x = int(box_w * 0.02) if cols > 1 else 0
    gap_y = int(box_h * 0.02) if rows > 1 else 0
    cell_w = (box_w - gap_x * (cols - 1)) // cols
    cell_h = (box_h - gap_y * (rows - 1)) // rows
    placed = 0
    for i, path in enumerate(paths):
        r, c = divmod(i, cols)
        left = box_left + c * (cell_w + gap_x)
        top = box_top + r * (cell_h + gap_y)
        try:
            _fit_image_in_box(slide, path, left, top, cell_w, cell_h)
            placed += 1
        except Exception as e:
            print(f"[layout_binder]   WARN: insert failed for {path.name}: {e}")
    return placed


def _carve_text_for_figure(slide, prs, fig_left: int) -> None:
    """Shrink or relocate text bodies that would collide with the carved
    right-hand figure column.

    Reads the resolved (inherited) geometry and writes ALL FOUR dims back:
    writing width alone on a placeholder with inherited geometry materializes
    an xfrm with cy=0, collapsing the text box to zero height.
    """
    slide_w = prs.slide_width
    gap = int(slide_w * 0.02)
    min_w = int(slide_w * 0.22)
    for shape in slide.shapes:
        is_body_ph = (shape.is_placeholder
                      and shape.placeholder_format.type in (2, 7)
                      and shape.has_text_frame)
        is_textbox = (not shape.is_placeholder
                      and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                      and shape.has_text_frame
                      and shape.text_frame.text.strip())
        if not (is_body_ph or is_textbox):
            continue
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if l is None or not w or not h:
            continue
        if l + w <= fig_left - gap:
            continue  # already clear of the figure column
        new_w = fig_left - gap - l
        if new_w < min_w:
            l = int(slide_w * 0.05)
            new_w = max(fig_left - gap - l, min_w)
        shape.left, shape.top, shape.width, shape.height = l, t or 0, new_w, h


def _insert_visuals(slide, prs, visual_paths: list[Path]) -> int:
    """Place one or more visuals (image / table / formula PNGs) on the slide.

    Distributes visuals across the template's designed picture spots when they
    exist and are usably sized; otherwise carves a right-hand column beside the
    text (relocating any text that would collide). Returns the count placed.
    """
    if not visual_paths:
        return 0

    from PIL import Image
    aspects: list[float] = []
    for p in visual_paths:
        try:
            with Image.open(str(p)) as img:
                fw, fh = img.size
            aspects.append(fw / fh if fh else 1.33)
        except Exception:
            aspects.append(1.33)

    slide_w, slide_h = prs.slide_width, prs.slide_height
    slide_area = slide_w * slide_h
    n = len(visual_paths)

    anchors = _anchor_boxes_for_slide(slide, slide_w, slide_h)
    carve_box = (int(slide_w * 0.54), int(slide_h * 0.24),
                 int(slide_w * 0.44), int(slide_h * 0.68))

    use_anchors = False
    chunks = []
    k = 0
    if anchors:
        k = min(len(anchors), n)
        sizes = [n // k + (1 if i < n % k else 0) for i in range(k)]
        start = 0
        for sz in sizes:
            chunks.append((visual_paths[start:start + sz], aspects[start:start + sz]))
            start += sz
        anchor_score = None
        for (box, _shape), (paths, asps) in zip(anchors[:k], chunks):
            _, min_area = _best_grid(len(paths), box[2], box[3], asps)
            anchor_score = min_area if anchor_score is None else min(anchor_score, min_area)
        _, carve_score = _best_grid(n, carve_box[2], carve_box[3], aspects)
        # Use the designed spots unless they are pathologically small compared
        # to what the beside-text fallback could offer.
        use_anchors = (anchor_score is not None
                       and (anchor_score >= 0.5 * carve_score
                            or anchor_score >= 0.05 * slide_area))

    placed = 0
    if use_anchors:
        for (box, shape), (paths, asps) in zip(anchors[:k], chunks):
            _remove_shape_element(slide, shape._element)
            stripped = _strip_pictures_in_region(slide, *box, overlap_threshold=0.1)
            if stripped:
                print(f"[layout_binder]   stripped {stripped} extra picture(s) near anchor")
            placed += _place_visuals_grid(slide, paths, box, asps)
        # Unused FILLED picture anchors would keep their stock photos next to
        # the real figures — remove those outright.
        for box, shape in anchors[k:]:
            if shape._element.tag == f"{_P_NS}pic":
                _remove_shape_element(slide, shape._element)
        print(f"[layout_binder]   anchored to {k} template picture spot(s)")
    else:
        box_left, box_top, box_w, box_h = carve_box
        _carve_text_for_figure(slide, prs, box_left)
        stripped = _strip_pictures_in_region(slide, box_left, box_top, box_w, box_h)
        if stripped:
            print(f"[layout_binder]   stripped {stripped} sample picture(s) from fallback region")
        placed += _place_visuals_grid(slide, visual_paths, carve_box, aspects)
        print("[layout_binder]   no usable picture anchor; placed figure(s) beside text")
    return placed


# ---------- Build the render sequence from SlideGen plan + binding ---------- #


def _pick_fallback_slide(
    descriptions: list[dict],
    preferred_types: list[str],
    avoid_idx: Optional[int] = None,
) -> Optional[int]:
    """Find the first non-meta template slide whose slide_type matches one of
    `preferred_types`, optionally avoiding a specific index (e.g. don't reuse
    the cover for everything). Returns the slide_index, or None if nothing
    matches.
    """
    for ptype in preferred_types:
        for d in descriptions:
            if d.get("is_meta"):
                continue
            idx = d.get("slide_index")
            if avoid_idx is not None and idx == avoid_idx:
                continue
            if d.get("slide_type") == ptype:
                return idx
    return None


def _build_render_sequence(
    slidegen_plan: dict,
    binding: dict,
    deck_meta: dict,
    template_descriptions: Optional[list[dict]] = None,
) -> list[dict]:
    """Build the full list of slides to render (cover, contents, dividers, body, thanks).

    Each entry: {kind, template_slide_index, content: {title, subtitle, bullets}, slidegen_slide?}

    If the binder LLM returned null for contents/section_divider/thanks, fall
    back to scanning `template_descriptions` for a usable slide_type.
    """
    seq: list[dict] = []
    descs = template_descriptions or []

    # Cover
    cover_idx = binding.get("cover_slide_index")
    if cover_idx is None:
        cover_idx = _pick_fallback_slide(descs, ["title"])
        if cover_idx is None:
            cover_idx = _first_non_meta(descs)  # never default onto a meta slide
        print(f"[layout_binder]   cover fallback -> slide {cover_idx}")
    seq.append({
        "kind": "cover",
        "template_slide_index": cover_idx,
        "content": {
            "title": deck_meta.get("deck_title", ""),
            "subtitle": deck_meta.get("deck_subtitle", ""),
        },
    })

    # Contents (agenda) — fall back to a content_bullets / two_column layout
    # if the LLM punted to null.
    unique_sections = []
    seen = set()
    for s in slidegen_plan.get("slides", []):
        sec = s.get("section", "")
        if sec and sec not in seen:
            seen.add(sec)
            unique_sections.append(sec)

    contents_idx = binding.get("contents_slide_index")
    if contents_idx is None and unique_sections:
        contents_idx = _pick_fallback_slide(
            descs, ["content_bullets", "two_column", "image_focus"],
            avoid_idx=cover_idx,
        )
        if contents_idx is not None:
            print(f"[layout_binder]   contents fallback -> slide {contents_idx}")
    if contents_idx is not None and unique_sections:
        seq.append({
            "kind": "contents",
            "template_slide_index": contents_idx,
            # distribute: agenda layouts with one placeholder per item get one
            # section per placeholder instead of everything in the first box.
            "content": {"title": "Contents", "bullets": unique_sections,
                        "distribute": True},
        })

    # Body slides, with section-divider injection on section transitions.
    section_div_idx = binding.get("section_divider_slide_index")
    if section_div_idx is None:
        section_div_idx = _pick_fallback_slide(
            descs, ["section_divider", "title"], avoid_idx=cover_idx,
        )
        if section_div_idx is not None:
            print(f"[layout_binder]   section_divider fallback -> slide {section_div_idx}")
    body_assignments = binding.get("body_assignments") or []
    current_section: Optional[str] = None
    section_counter = 0
    for i, plan_slide in enumerate(slidegen_plan.get("slides", [])):
        section = plan_slide.get("section", "")
        if section != current_section:
            current_section = section
            section_counter += 1
            if section_div_idx is not None and section:
                seq.append({
                    "kind": "section_divider",
                    "template_slide_index": section_div_idx,
                    "content": {
                        "title": section,
                        "subtitle": f"PART {section_counter:02d}",
                    },
                })

        tpl_idx = body_assignments[i] if i < len(body_assignments) else (cover_idx or 0)
        bullets = _flatten_slidegen_bullets(plan_slide.get("bullets") or [])

        # Two-column (T19) content lives under columns[].bullets. Carry the
        # columns through for per-placeholder filling, plus a flattened merge
        # as the single-placeholder fallback.
        content_cols = []
        for c in (plan_slide.get("columns") or []):
            if not isinstance(c, dict):
                continue
            header = c.get("subsection") or c.get("header") or c.get("title") or ""
            col_bullets = _flatten_slidegen_bullets(c.get("bullets") or [])
            if header or col_bullets:
                content_cols.append({"header": header, "bullets": col_bullets})
        if bullets:
            content_cols = []  # top-level bullets win; don't double-fill
        elif content_cols:
            for c in content_cols:
                if c["header"]:
                    bullets.append((c["header"], 0))
                    bullets.extend((t, min(lvl + 1, 4)) for t, lvl in c["bullets"])
                else:
                    bullets.extend(c["bullets"])
        seq.append({
            "kind": "body",
            "template_slide_index": tpl_idx,
            "content": {
                "title": plan_slide.get("subsection", "") or plan_slide.get("section", ""),
                "bullets": bullets,
                "columns": content_cols,
                "n_visuals": _plan_slide_visuals(plan_slide),
            },
            # Carry the raw SlideGen slide so the renderer can resolve visuals
            # (images + tables + formulas) via layout_filler.resolve_visual_paths.
            "slidegen_slide": plan_slide,
        })

    # Thanks
    thanks_idx = binding.get("thanks_slide_index")
    if thanks_idx is None:
        thanks_idx = _pick_fallback_slide(
            descs, ["thanks", "title"], avoid_idx=cover_idx,
        )
        if thanks_idx is not None:
            print(f"[layout_binder]   thanks fallback -> slide {thanks_idx}")
    if thanks_idx is not None:
        seq.append({
            "kind": "thanks",
            "template_slide_index": thanks_idx,
            "content": {"title": "Thank You", "subtitle": ""},
        })

    return seq


# ---------- Top-level entry point ---------- #


def _load_slidegen_artifacts(args) -> tuple[dict, dict, dict, dict]:
    """Load SlideGen's plan / raw content / figures / formulas from contents."""
    variant_suffix = "_personalized" if getattr(args, "use_author_preferences", False) else "_baseline"
    contents_dir = REPO_ROOT / "contents" / args.paper_name
    plan_path = contents_dir / f"<{args.model_name_t}_{args.model_name_v}>_slide_plan{variant_suffix}.json"
    raw_path = contents_dir / f"<{args.model_name_t}_{args.model_name_v}>_raw_content.json"
    figs_path = contents_dir / f"<{args.model_name_t}_{args.model_name_v}>_figures.json"
    formulas_path = contents_dir / f"<{args.model_name_t}_{args.model_name_v}>_formula_match.json"

    if not plan_path.exists():
        raise FileNotFoundError(f"SlideGen plan not found: {plan_path}")
    plan = json.loads(plan_path.read_text(encoding="utf-8"))

    raw = {}
    if raw_path.exists():
        raw = json.loads(raw_path.read_text(encoding="utf-8"))

    figures = {}
    if figs_path.exists():
        figures = json.loads(figs_path.read_text(encoding="utf-8"))

    formulas = {}
    if formulas_path.exists():
        formulas = json.loads(formulas_path.read_text(encoding="utf-8"))

    return plan, raw, figures, formulas


def _append_body_notes(
    slide,
    slidegen_slide: dict,
    outline_json: dict,
    figs_data: dict,
    formula_data: dict,
    *,
    get_content,
    get_image_reasons,
    get_table_reasons,
    get_formula_reasons,
) -> bool:
    """Add the same body-slide notes used by layout_filler.generate_pptx_from_plan."""
    notes_chunks: list[str] = []
    section = slidegen_slide.get("section", "")
    subsection = slidegen_slide.get("subsection", "")

    txt = get_content(section, subsection, outline_json)
    if txt:
        notes_chunks.append(txt)

    if slidegen_slide.get("images"):
        img_r = get_image_reasons(section, subsection, slidegen_slide["images"], figs_data)
        if img_r:
            notes_chunks.append(img_r)

    if slidegen_slide.get("tables"):
        tb_r = get_table_reasons(section, subsection, slidegen_slide["tables"], figs_data)
        if tb_r:
            notes_chunks.append(tb_r)

    if slidegen_slide.get("formulas"):
        fm_r = get_formula_reasons(section, subsection, slidegen_slide["formulas"], formula_data)
        if fm_r:
            notes_chunks.append(fm_r)

    if not notes_chunks:
        return False

    nframe = slide.notes_slide.notes_text_frame
    if nframe.text and not nframe.text.endswith("\n"):
        nframe.text += "\n"
    nframe.text += "\n\n".join(notes_chunks)
    return True


def _slugify_template_label(template_path: Path) -> str:
    """Derive a filesystem-safe label from a template path (e.g. 'bee template.pptx' -> 'bee_template')."""
    stem = template_path.stem
    return re.sub(r"[^A-Za-z0-9]+", "_", stem).strip("_").lower()


def _template_fingerprint(path: Path) -> str:
    h = hashlib.md5()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return f"{path.stat().st_size}:{h.hexdigest()}"


def _visible_slide_indices(prs) -> list[int]:
    """python-pptx indices of slides that are NOT hidden. LibreOffice omits
    hidden slides from the PDF export, so PDF-page order only matches pptx
    order after skipping them."""
    return [i for i, slide in enumerate(prs.slides)
            if slide.element.get("show") != "0"]


def bind_and_render(args, user_template_path: str, template_label: Optional[str] = None) -> Optional[Path]:
    """Run binder + render. Assumes SlideGen's full pipeline has already run.

    Reads the cached plan/raw_content/figures from contents/<paper_name>/, calls
    the binder LLM, and writes the output PPTX. Returns the output path on
    success, None on failure.

    `template_label` is a short tag (e.g. "uchicago", "bee") used in artifact
    paths and the output filename so multiple templates can share one
    paper_name without overwriting each other. If None, derived from template
    filename.
    """
    template_path = Path(user_template_path).expanduser().resolve()
    if not template_path.exists():
        print(f"[layout_binder] ERROR: template not found: {template_path}")
        return None
    if template_path.suffix.lower() != ".pptx":
        print(f"[layout_binder] ERROR: template must be .pptx: {template_path}")
        return None

    if template_label is None:
        template_label = _slugify_template_label(template_path)

    try:
        plan, raw_content, figures, formulas = _load_slidegen_artifacts(args)
    except FileNotFoundError as e:
        print(f"[layout_binder] ERROR: {e}")
        return None

    deck_meta = {
        "deck_title": (raw_content.get("metadata") or {}).get("title", args.paper_name),
        "deck_subtitle": (raw_content.get("metadata") or {}).get("author", ""),
    }

    contents_dir = REPO_ROOT / "contents" / args.paper_name
    # Template-specific subdir so different templates against the same
    # paper_name don't collide on _user_template_images / _user_template_binding.
    template_artifacts_dir = contents_dir / f"_user_template__{template_label}"
    template_images_dir = template_artifacts_dir / "images"
    descriptions_cache = template_artifacts_dir / "descriptions.json"
    fingerprint_path = template_artifacts_dir / "template_fingerprint.txt"

    # Invalidate cached images/descriptions when the template FILE changed —
    # the caches are keyed only by label, and stale ones make the binder pick
    # slides from a template that no longer exists. A pre-fingerprint cache is
    # adopted as-is (fingerprint written) rather than thrown away.
    fingerprint = _template_fingerprint(template_path)
    stored_fp = fingerprint_path.read_text(encoding="utf-8").strip() if fingerprint_path.exists() else None
    if stored_fp is not None and stored_fp != fingerprint:
        print(f"[layout_binder]   template file changed since caches were built; regenerating [{template_label}]")
        if template_images_dir.exists():
            for png in template_images_dir.glob("slide_*.png"):
                png.unlink()
        descriptions_cache.unlink(missing_ok=True)
    template_artifacts_dir.mkdir(parents=True, exist_ok=True)
    fingerprint_path.write_text(fingerprint, encoding="utf-8")

    with _stage(1, f"render template to images [{template_label}]"):
        # Cache: skip LibreOffice if the per-slide PNGs already exist.
        existing_pngs = sorted(template_images_dir.glob("slide_*.png")) if template_images_dir.exists() else []
        if existing_pngs:
            image_paths = existing_pngs
            print(f"[layout_binder]   reusing {len(image_paths)} cached template slide images")
        else:
            image_paths = render_template_to_images(template_path, template_images_dir)
            print(f"[layout_binder]   rendered {len(image_paths)} template slide images")

    with _stage(2, f"describe template slides (vision LLM) [{template_label}]"):
        # Cache: skip the vision LLM if descriptions are already on disk.
        if descriptions_cache.exists():
            descriptions = json.loads(descriptions_cache.read_text(encoding="utf-8"))
            print(f"[layout_binder]   reusing cached descriptions for {len(descriptions)} slides")
        else:
            descriptions = describe_template_slides(image_paths)
            descriptions_cache.parent.mkdir(parents=True, exist_ok=True)
            descriptions_cache.write_text(json.dumps(descriptions, indent=2), encoding="utf-8")

        # Descriptions are in PDF-page order; remap to python-pptx slide order
        # when the template contains hidden slides (LibreOffice skips them in
        # the PDF, python-pptx does not — one hidden slide would silently shift
        # every clone after it).
        prs_tpl = Presentation(str(template_path))
        visible = _visible_slide_indices(prs_tpl)
        n_all = len(prs_tpl.slides)
        if len(visible) < n_all:
            if len(descriptions) == len(visible):
                for d, real_idx in zip(descriptions, visible):
                    d["slide_index"] = real_idx
                print(f"[layout_binder]   template has {n_all - len(visible)} hidden slide(s); "
                      f"remapped description indices to pptx order")
            else:
                print(f"[layout_binder]   WARN: {len(descriptions)} descriptions vs "
                      f"{len(visible)} visible slides — indices may be misaligned")

        # Merge in programmatic picture-anchor ground truth (always — the
        # cached descriptions.json may predate this field).
        pic_counts = scan_picture_placeholders(template_path)
        for d in descriptions:
            n = pic_counts.get(d.get("slide_index"), 0)
            d["has_picture_placeholder"] = n > 0
            d["n_picture_placeholders"] = n
        n_pic = sum(1 for d in descriptions if d.get("has_picture_placeholder"))
        print(f"[layout_binder]   {n_pic}/{len(descriptions)} template slides have a picture anchor")

    with _stage(3, f"bind plan slides to template slides [{template_label}]"):
        binding = bind_user_template_to_plan(plan, descriptions, deck_meta)
        print(f"[layout_binder]   binding: cover={binding.get('cover_slide_index')}, "
              f"contents={binding.get('contents_slide_index')}, "
              f"divider={binding.get('section_divider_slide_index')}, "
              f"thanks={binding.get('thanks_slide_index')}, "
              f"body={binding.get('body_assignments')}")
        binding_cache = template_artifacts_dir / "binding.json"
        binding_cache.parent.mkdir(parents=True, exist_ok=True)
        binding_cache.write_text(json.dumps(binding, indent=2), encoding="utf-8")

    with _stage(4, "render output PPTX"):
        # Lazy import to avoid pulling apply_color / heavy filler deps at module load.
        from SlidesAgent.layout_filler import (
            get_content,
            get_formula_reasons,
            get_image_reasons,
            get_table_reasons,
            resolve_visual_paths,
        )

        seq = _build_render_sequence(plan, binding, deck_meta, template_descriptions=descriptions)
        prs = Presentation(str(template_path))
        n_template_slides = len(prs.slides)

        rendered = 0
        visuals_inserted = 0
        notes_inserted = 0
        for entry in seq:
            src = entry["template_slide_index"]
            if src is None or not (0 <= src < n_template_slides):
                print(f"[layout_binder]   WARN: skipping {entry['kind']} (bad template_slide_index={src})")
                continue
            new_slide = _duplicate_slide(prs, src)
            # Strip leftover template sample content (sample charts/tables,
            # callout groups with "Bullet level 1" etc.) before filling.
            stripped = _strip_sample_shapes(new_slide)
            if stripped:
                print(f"[layout_binder]   stripped {stripped} sample shape(s) from {entry['kind']}")
            try:
                _fill_slide_text(new_slide, entry["content"], prs)
            except Exception as e:
                print(f"[layout_binder]   WARN: text fill failed on {entry['kind']}: {e}")

            # Resolve and insert visuals (images + tables + formulas) for body slides.
            slidegen_slide = entry.get("slidegen_slide")
            if slidegen_slide:
                try:
                    if _append_body_notes(
                        new_slide,
                        slidegen_slide,
                        raw_content,
                        figures,
                        formulas,
                        get_content=get_content,
                        get_image_reasons=get_image_reasons,
                        get_table_reasons=get_table_reasons,
                        get_formula_reasons=get_formula_reasons,
                    ):
                        notes_inserted += 1
                except Exception as e:
                    print(f"[layout_binder]   WARN: speaker-note fill failed on {entry['kind']}: {e}")
                try:
                    visual_paths = resolve_visual_paths(slidegen_slide, args)
                except Exception as e:
                    print(f"[layout_binder]   WARN: visual resolution failed on {entry['kind']}: {e}")
                    visual_paths = []
                if visual_paths:
                    placed = _insert_visuals(new_slide, prs, visual_paths)
                    visuals_inserted += placed
            rendered += 1

        _delete_slides_at(prs, list(range(n_template_slides)))

        scrubbed = _scrub_layouts_and_masters(prs)
        if scrubbed:
            print(f"[layout_binder]   scrubbed {scrubbed} sample text shape(s) on layouts/masters")

        variant_suffix = "_personalized" if getattr(args, "use_author_preferences", False) else "_baseline"
        output_pptx = contents_dir / f"{args.model_name_t}_{args.model_name_v}_output_slides{variant_suffix}_user_template__{template_label}.pptx"
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_pptx))
        print(
            f"[layout_binder]   rendered {rendered} slides "
            f"({visuals_inserted} visuals placed, {notes_inserted} notes added) -> {output_pptx}"
        )

    return output_pptx
