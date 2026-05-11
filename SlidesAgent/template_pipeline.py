"""Template-aware slide generation — lean rebuild of Branch B.

Replaces the SlideTailor subprocess pipeline with a single-codebase
implementation that runs entirely in the paper2pptx env using SlideGen's
existing Azure infrastructure.

Five-stage flow:
  1. Parse the user's paper via Docling (markdown export, cached)
  2. Render the user's template to slide images
  3. Vision LLM describes each template slide's structure
  4. LLM plans which template slide each output slide should be based on,
     with title and bullets per slide (cached)
  5. python-pptx writes the final PPTX by duplicating chosen template slides
     and replacing text/image content
"""

from __future__ import annotations

import base64
import copy
import json
import os
import re
import subprocess
import sys
import time
import traceback
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Optional

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from slidegen_openai_utils import build_openai_client, resolve_direct_model_name


SOFFICE_BIN = "/opt/homebrew/bin/soffice"

DESCRIBE_SYSTEM_PROMPT = """You are analyzing a single slide from a PowerPoint template.
Look at the slide image and return a JSON object describing its structure.

Output schema (return ONLY the JSON, no prose):
{
  "slide_type": "title" | "section_divider" | "content_bullets" | "image_focus" | "two_column" | "comparison" | "data_table" | "thanks" | "blank" | "instruction" | "other",
  "description": "1-2 sentences describing what this slide visually shows",
  "is_meta": true | false,
  "text_regions": [
    {"role": "title|subtitle|body|caption|footer|other", "approx_position": "top|center|bottom|left|right"}
  ],
  "image_regions": [
    {"approx_position": "top|center|bottom|left|right", "size": "small|medium|large"}
  ],
  "suitable_for": ["short tag describing what content this slide layout could host"]
}

Set "is_meta": true when the slide is a template-instruction page (e.g. "delete this note",
"how to use this template", "click here in PowerPoint"). Otherwise set "is_meta": false.
"""


PLAN_SYSTEM_PROMPT = """You are planning a presentation deck for a scientific paper, using a user-supplied PowerPoint template.

You will be given:
  1. The user's paper as Markdown.
  2. A JSON list describing each slide of the user's template (its visual structure).
  3. A JSON list of figures available from the paper (each with id, caption, page_no, aspect).
  4. (Optional) An author preference profile JSON describing this author's typical slide style.

Your job is to produce a slide plan: for each output slide, pick which template slide to base it on (by template_slide_index), write the title and bullet content, and OPTIONALLY assign one figure_id from the available figures.

Constraints:
  - Use ONLY template slides where "is_meta" is false. Skip meta/instruction slides entirely.
  - Aim for 8-12 output slides for a typical paper.
  - The first output slide should be a title slide (use a template slide whose slide_type is "title").
  - Include sections in a logical order (introduction/motivation, related work briefly, method, experiments/results, conclusion).
  - For bullet content: short, presentation-friendly phrases. Not full sentences from the paper.
  - `bullets` MUST be a flat list of plain strings. Do NOT use nested objects like {"level1": "...", "items": [...]}. If you need sub-points, prepend a leading "- " or indent with two spaces in the same string.
  - Do NOT include LaTeX math markers ($...$, \\frac, etc.) in bullets. Use plain words or Unicode (≤, ≥, ℓ, ε) — anything in $...$ will render as raw text in PowerPoint.
  - LAYOUT FIT MATTERS. If a slide will have bullets, pick a template slide whose slide_type is "content_bullets", "two_column", "comparison", or "image_focus" — these have body/object regions that fit bulleted content. Do NOT pick "title", "section_divider", or "thanks" layouts for slides with bullets, since those layouts have a single text region or none at all and bullets will be dropped or cramped. Use "section_divider" only for slides whose content is just a one-line section heading with no bullets.

Figure assignment:
  - You MAY assign at most ONE figure_id per slide (use null if no figure fits).
  - Each figure_id may be used on AT MOST ONE slide across the deck — do not repeat figures.
  - Pick figures that are central to the paper (architecture diagrams, results plots, motivating examples). Skip figures that are decorative or only locally relevant.
  - When you assign a figure to a slide, prefer a template_slide_index whose slide_type is "image_focus" or "two_column" so there is room for the figure alongside text.

Author preferences (when provided):
  - The profile contains coarse signals like slide_count_preference, bullet_density_preference, text_density_preference, figure_usage_preference, and prefers_takeaway_slide / prefers_multi_slide_method_section / prefers_multi_slide_results_section.
  - Treat them as soft guidance, not hard rules. Examples:
    - figure_usage_preference="high" → assign figures to more slides (most content slides should have one).
    - bullet_density_preference="low" → fewer bullets per slide (2-3); "high" → more bullets (4-6).
    - prefers_multi_slide_method_section="high" → split method across 2-3 slides instead of compressing into one overview.
    - prefers_takeaway_slide="high" → end with an explicit takeaways/conclusion slide.

Output schema (return ONLY the JSON, no prose):
{
  "deck_title": "concise paper title for the cover slide",
  "deck_subtitle": "authors string for the cover slide",
  "slides": [
    {
      "output_index": 0,
      "template_slide_index": <int>,
      "purpose": "short label, e.g. 'title', 'motivation', 'method overview'",
      "figure_id": <int or null>,
      "content": {
        "title": "title text for this slide",
        "subtitle": "optional subtitle (or empty string)",
        "bullets": ["bullet 1", "bullet 2", "..."]
      },
      "reasoning": "1 sentence on why this template slide fits this content"
    }
  ]
}
"""


# ---------- Logging + pre-flight helpers ---------- #


@contextmanager
def _stage(num, label: str):
    """Context manager that prints a stage header and the elapsed time on exit."""
    print(f"[template-pipeline] (stage {num}) {label}...")
    t0 = time.time()
    try:
        yield
    finally:
        elapsed = time.time() - t0
        print(f"[template-pipeline]   (stage {num} took {elapsed:.1f}s)")


def _check_prerequisites(template_path: Path) -> list[str]:
    """Return a list of human-readable problems; empty list means we're good to go."""
    problems: list[str] = []

    # Stage 2 needs LibreOffice to convert PPTX→PDF.
    if not Path(SOFFICE_BIN).exists():
        problems.append(
            f"LibreOffice binary not found at {SOFFICE_BIN}. "
            "Install with `brew install --cask libreoffice` or update SOFFICE_BIN in template_pipeline.py."
        )

    # Stages 3 and 4 need Azure OpenAI credentials.
    if not os.environ.get("AZURE_OPENAI_API_KEY"):
        problems.append(
            "AZURE_OPENAI_API_KEY is not set. "
            "Add it to .env (see CHANGELOG for the full set of Azure vars)."
        )

    # Template must be a PPTX, not PPT or PDF.
    if template_path.suffix.lower() != ".pptx":
        problems.append(f"Template must be a .pptx file (got {template_path.suffix}).")

    return problems


# ---------- Stage 1: parse paper (markdown + figures, single Docling pass) ---------- #


_IMAGE_SCALE = 5.0  # matches SlidesAgent/parse_raw.py:IMAGE_RESOLUTION_SCALE


def parse_paper(paper_path: Path, output_dir: Path) -> tuple[str, list[dict[str, Any]]]:
    """Parse the user's paper once, exporting both markdown and figure crops.

    Caches:
      - markdown   -> output_dir / "_paper_markdown.md"
      - figures    -> output_dir / "_paper_figures.json" + output_dir / "_paper_figures/*.png"

    Figure JSON entries look like:
      {"id": 1, "caption": "...", "page_no": 2, "image_path": "...", "width": 957, "height": 721, "aspect": 1.33}

    Returns (markdown, figures_list).
    """
    md_cache = output_dir / "_paper_markdown.md"
    figs_cache = output_dir / "_paper_figures.json"
    figs_dir = output_dir / "_paper_figures"

    if md_cache.exists() and figs_cache.exists() and figs_dir.exists():
        print(f"[template-pipeline]   (cached) {md_cache} + {figs_cache}")
        return md_cache.read_text(encoding="utf-8"), json.loads(figs_cache.read_text(encoding="utf-8"))

    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling_core.types.doc.document import BoundingBox

    opts = PdfPipelineOptions()
    opts.images_scale = _IMAGE_SCALE
    opts.generate_page_images = True
    opts.generate_picture_images = True

    converter = DocumentConverter(
        format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
    )
    result = converter.convert(str(paper_path))
    doc = result.document

    markdown = doc.export_to_markdown()
    md_cache.write_text(markdown, encoding="utf-8")

    figs_dir.mkdir(parents=True, exist_ok=True)
    figures: list[dict[str, Any]] = []
    for i, picture in enumerate(doc.pictures, start=1):
        caption = picture.caption_text(doc) or ""
        if not caption.strip():
            continue  # skip uncaptioned pictures (logos, watermarks)
        if not picture.prov:
            continue
        prov = picture.prov[0]
        page = doc.pages.get(prov.page_no)
        if page is None or page.image is None:
            continue
        full_img = page.image.pil_image
        page_size = page.size
        page_w_pt, page_h_pt = page_size.width, page_size.height
        scale = full_img.width / page_w_pt

        bbox = prov.bbox
        pad = 1
        padded = BoundingBox(
            l=bbox.l - pad, r=bbox.r + pad,
            b=bbox.b - pad, t=bbox.t + pad,
            coord_origin=bbox.coord_origin,
        )
        tl_box = padded.to_top_left_origin(page_height=page_h_pt).scaled(scale=scale).as_tuple()
        cropped = full_img.crop(tl_box)
        img_path = figs_dir / f"figure_{i}.png"
        cropped.save(img_path, "PNG")

        figures.append({
            "id": i,
            "caption": caption.strip(),
            "page_no": prov.page_no,
            "image_path": str(img_path.relative_to(REPO_ROOT)) if img_path.is_relative_to(REPO_ROOT) else str(img_path),
            "width": cropped.width,
            "height": cropped.height,
            "aspect": cropped.width / cropped.height if cropped.height else 1.0,
        })

    figs_cache.write_text(json.dumps(figures, indent=2), encoding="utf-8")
    print(f"[template-pipeline]   parsed and cached: {md_cache} + {len(figures)} figures")
    return markdown, figures


# ---------- Stage 2: render template to images ---------- #


def render_template_to_images(template_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [SOFFICE_BIN, "--headless", "--convert-to", "pdf",
             "--outdir", str(output_dir), str(template_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
    except subprocess.CalledProcessError as e:
        # Surface stderr — soffice's actual error message lives there.
        stderr = e.stderr.decode("utf-8", errors="replace") if e.stderr else "(no stderr)"
        raise RuntimeError(
            f"LibreOffice failed to convert template to PDF (exit {e.returncode}). "
            f"Stderr: {stderr.strip()[:500]}"
        ) from e
    except subprocess.TimeoutExpired as e:
        raise RuntimeError(
            f"LibreOffice timed out after {e.timeout}s converting {template_path.name}. "
            "Template may be too large or LibreOffice may be hung."
        ) from e

    pdf_path = output_dir / f"{template_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"LibreOffice exited cleanly but produced no PDF at {pdf_path}. "
            "Check that the template opens correctly in PowerPoint/LibreOffice."
        )

    images = convert_from_path(pdf_path, dpi=120)
    image_paths: list[Path] = []
    for i, img in enumerate(images):
        path = output_dir / f"slide_{i:03d}.png"
        img.save(path, "PNG")
        image_paths.append(path)
    return image_paths


# ---------- Stage 3: vision LLM describes template slides ---------- #


def _parse_json_from_response(text: str) -> dict[str, Any]:
    text = text.strip()
    if text.startswith("```"):
        text = text.split("```", 2)[1]
        if text.startswith("json"):
            text = text[4:]
        text = text.rsplit("```", 1)[0].strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        try:
            from json_repair import repair_json
            return json.loads(repair_json(text))
        except Exception:
            return {}


def describe_template_slides(image_paths: list[Path]) -> list[dict[str, Any]]:
    client = build_openai_client()
    deployment = resolve_direct_model_name(os.environ.get("AZURE_DEPLOYMENT_NAME", ""))
    descriptions: list[dict[str, Any]] = []
    for i, img_path in enumerate(image_paths):
        with img_path.open("rb") as f:
            img_b64 = base64.b64encode(f.read()).decode()
        response = client.chat.completions.create(
            model=deployment,
            messages=[
                {"role": "system", "content": DESCRIBE_SYSTEM_PROMPT},
                {"role": "user", "content": [
                    {"type": "text", "text": f"This is slide {i+1} of the template."},
                    {"type": "image_url",
                     "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                ]},
            ],
            max_completion_tokens=800,
        )
        raw = response.choices[0].message.content or ""
        parsed = _parse_json_from_response(raw)
        parsed.setdefault("slide_type", "other")
        parsed.setdefault("description", "")
        parsed.setdefault("is_meta", False)
        parsed.setdefault("text_regions", [])
        parsed.setdefault("image_regions", [])
        parsed.setdefault("suitable_for", [])
        parsed["slide_index"] = i
        descriptions.append(parsed)
        meta_tag = " [META]" if parsed.get("is_meta") else ""
        print(f"[template-pipeline]   slide {i+1}: {parsed.get('slide_type', '?')}{meta_tag} - {parsed.get('description', '')[:70]}")
    return descriptions


# ---------- Stage 4: plan output slides ---------- #


def plan_output_slides(
    paper_markdown: str,
    template_descriptions: list[dict[str, Any]],
    figures: Optional[list[dict[str, Any]]] = None,
    author_profile: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    client = build_openai_client()
    deployment = resolve_direct_model_name(os.environ.get("AZURE_DEPLOYMENT_NAME", ""))

    # Drop the heavy image_path field from figures shown to the planner — it just
    # needs id+caption+aspect to decide assignment.
    figs_for_prompt = [
        {"id": f["id"], "caption": f["caption"], "page_no": f["page_no"], "aspect": round(f["aspect"], 2)}
        for f in (figures or [])
    ]

    user_payload = {
        "paper_markdown": paper_markdown[:30000],  # rough truncation for context-window safety
        "template_slides": template_descriptions,
        "figures": figs_for_prompt,
        "author_profile": author_profile,  # null if not provided
    }

    response = client.chat.completions.create(
        model=deployment,
        messages=[
            {"role": "system", "content": PLAN_SYSTEM_PROMPT},
            {"role": "user", "content": json.dumps(user_payload)},
        ],
        max_completion_tokens=4000,
    )
    raw = response.choices[0].message.content or ""
    plan = _parse_json_from_response(raw)
    if not plan or "slides" not in plan:
        print(f"[template-pipeline]   WARN: plan missing or malformed; raw response head: {raw[:200]}")
        plan = {"deck_title": "", "deck_subtitle": "", "slides": []}
    return plan


# ---------- Stage 5: render output PPTX ---------- #


# Placeholder type IDs from python-pptx (PP_PLACEHOLDER enum):
#   1 TITLE, 2 BODY, 3 CENTER_TITLE, 4 SUBTITLE, 7 OBJECT, 13 SLIDE_NUMBER, 14 HEADER, 15 FOOTER, 16 DATE
_TITLE_TYPES = {1, 3}        # TITLE, CENTER_TITLE
_SUBTITLE_TYPES = {4}        # SUBTITLE
_BULLET_TYPES = {7, 2}       # OBJECT preferred, BODY as fallback
_SKIP_TYPES = {13, 14, 15, 16}  # slide number / header / footer / date — leave alone


def _duplicate_slide(prs, src_index: int):
    """Append a deep copy of slide `src_index` to the end of `prs`.

    Uses the standard idiom: add a new slide bound to the source's layout, then
    deep-copy every shape from source's spTree into the new slide's spTree.
    The auto-inserted layout placeholders are removed first to avoid duplicates.
    """
    source = prs.slides[src_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Strip the placeholders python-pptx auto-inserted from the layout —
    # we want a faithful copy of the source slide, not a layout instantiation.
    new_spTree = new_slide.shapes._spTree
    for shape in list(new_slide.shapes):
        new_spTree.remove(shape.element)

    for shape in source.shapes:
        new_spTree.append(copy.deepcopy(shape.element))

    return new_slide


def _delete_slides_at(prs, indices: list[int]) -> None:
    """Delete slides at the given indices (in the original ordering)."""
    sldIdLst = prs.slides._sldIdLst
    sldId_elems = list(sldIdLst)
    for i in sorted(indices, reverse=True):
        sldIdLst.remove(sldId_elems[i])


def _enable_shrink_to_fit(tf) -> None:
    """Mirror PowerPoint's 'Shrink text on overflow' so long content stays in-box.

    Without this, deep-copied placeholders keep the source slide's font size
    and overflow wildly when given longer text (e.g. a 50-char paper title in
    a slot designed for a 24-char template heading).
    """
    if tf is None:
        return
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_text_frame(tf, text: str) -> None:
    """Replace text in a text_frame, preserving the first run's formatting if possible."""
    if not tf.paragraphs:
        tf.text = text
        _enable_shrink_to_fit(tf)
        return
    p0 = tf.paragraphs[0]
    # Drop all runs after the first to avoid leftover sample text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    if p0.runs:
        p0.runs[0].text = text
    else:
        p0.text = text
    # Drop any extra paragraphs from the placeholder's sample content
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    _enable_shrink_to_fit(tf)


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ---- Plain-text math cleanup ----
# Map common LaTeX commands → Unicode so bullet/title text doesn't ship as raw $...$.
_LATEX_UNICODE = {
    r"\\leq": "≤", r"\\geq": "≥", r"\\le": "≤", r"\\ge": "≥",
    r"\\neq": "≠", r"\\approx": "≈", r"\\sim": "∼",
    r"\\to": "→", r"\\rightarrow": "→", r"\\Rightarrow": "⇒",
    r"\\leftarrow": "←", r"\\Leftarrow": "⇐", r"\\leftrightarrow": "↔",
    r"\\times": "×", r"\\cdot": "·", r"\\pm": "±",
    r"\\alpha": "α", r"\\beta": "β", r"\\gamma": "γ", r"\\delta": "δ",
    r"\\epsilon": "ε", r"\\varepsilon": "ε", r"\\zeta": "ζ", r"\\eta": "η",
    r"\\theta": "θ", r"\\lambda": "λ", r"\\mu": "μ", r"\\nu": "ν",
    r"\\pi": "π", r"\\rho": "ρ", r"\\sigma": "σ", r"\\tau": "τ",
    r"\\phi": "φ", r"\\chi": "χ", r"\\psi": "ψ", r"\\omega": "ω",
    r"\\Gamma": "Γ", r"\\Delta": "Δ", r"\\Theta": "Θ", r"\\Lambda": "Λ",
    r"\\Pi": "Π", r"\\Sigma": "Σ", r"\\Phi": "Φ", r"\\Omega": "Ω",
    r"\\ell": "ℓ", r"\\infty": "∞", r"\\partial": "∂", r"\\nabla": "∇",
    r"\\in": "∈", r"\\subset": "⊂", r"\\subseteq": "⊆", r"\\cup": "∪", r"\\cap": "∩",
    r"\\forall": "∀", r"\\exists": "∃",
}

_MATH_DELIM_RE = re.compile(r"\$+([^$]*?)\$+")  # $...$ or $$...$$


def _flatten_bullets(raw: list[Any]) -> list[str]:
    """Coerce a possibly-nested bullets value into a flat list of strings.

    The planner is instructed to emit `bullets` as `list[str]`, but occasionally
    returns nested objects like `{"level1": "...", "items": [...]}`. This helper
    flattens those into level-0 + indented sub-bullets so we never crash at fill
    time.
    """
    out: list[str] = []
    for item in raw:
        if isinstance(item, str):
            out.append(item)
        elif isinstance(item, dict):
            # Heuristic: heading-ish key first, then any list-typed value as sub-bullets.
            heading = next((str(v) for k, v in item.items() if isinstance(v, str)), None)
            if heading:
                out.append(heading)
            for v in item.values():
                if isinstance(v, list):
                    for sub in v:
                        if isinstance(sub, str):
                            out.append("  " + sub)
                        elif isinstance(sub, dict):
                            out.extend(_flatten_bullets([sub]))
        elif isinstance(item, list):
            out.extend(_flatten_bullets(item))
        else:
            out.append(str(item))
    return out


def _clean_math(text: str) -> str:
    """Strip basic LaTeX so bullets render as readable plain text in PowerPoint."""
    if not text or "$" not in text and "\\" not in text:
        return text

    # Replace commands first (inside or outside $...$).
    for cmd, repl in _LATEX_UNICODE.items():
        text = re.sub(cmd + r"\b", repl, text)

    # Drop $...$ delimiters but keep their content.
    text = _MATH_DELIM_RE.sub(lambda m: m.group(1), text)

    # Strip a few leftover formatting commands; keep their args.
    text = re.sub(r"\\(?:mathbf|mathit|mathrm|text|emph|bf|it)\{([^}]*)\}", r"\1", text)
    # Drop unhandled bare commands like \frac, \sum (rare in short bullets).
    text = re.sub(r"\\[a-zA-Z]+", "", text)

    return text.strip()


def _set_bullets(tf, bullets: list[str]) -> None:
    """Replace the text frame's contents with one paragraph per bullet,
    preserving the formatting of paragraph 0 from the source placeholder."""
    if not bullets:
        return
    # Seed paragraph 0 from the first bullet, preserving its formatting.
    _set_text_frame(tf, bullets[0])
    template_para_xml = tf.paragraphs[0]._p
    txBody = template_para_xml.getparent()

    for bullet in bullets[1:]:
        new_p = copy.deepcopy(template_para_xml)
        # Strip the copied paragraph's runs so we can write fresh text into it.
        for r in list(new_p.findall(f"{_A_NS}r")):
            new_p.remove(r)
        # Also strip soft line breaks from the template paragraph copy.
        for br in list(new_p.findall(f"{_A_NS}br")):
            new_p.remove(br)
        txBody.append(new_p)
        # The newly appended paragraph is now the last in the text frame.
        last_para = tf.paragraphs[-1]
        run = last_para.add_run()
        run.text = bullet


def _clear_text_frame(tf) -> None:
    """Empty out a text frame so leftover placeholder sample text doesn't show."""
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)


def _fill_slide_text(slide, content: dict[str, Any]) -> None:
    """Fill placeholders on `slide` with title/subtitle/bullets from `content`.

    Matches by placeholder TYPE (choice C). First matching placeholder wins.
    Any text-bearing content placeholder we don't fill is cleared so leftover
    template sample text ("Bullet level 1, Bullet level 2 …") doesn't ship.
    """
    title_text = _clean_math(content.get("title", "") or "")
    subtitle_text = _clean_math(content.get("subtitle", "") or "")
    bullets = [_clean_math(b) for b in _flatten_bullets(content.get("bullets") or [])]
    bullets = [b for b in bullets if b]  # drop ones that became empty after cleanup

    used_phs: set[int] = set()  # placeholder ph.placeholder_format.idx values

    # 1. Title — first TITLE/CENTER_TITLE placeholder.
    if title_text:
        for ph in slide.placeholders:
            if ph.placeholder_format.type in _TITLE_TYPES and ph.has_text_frame:
                _set_text_frame(ph.text_frame, title_text)
                used_phs.add(ph.placeholder_format.idx)
                break

    # 2. Bullets — first unused OBJECT, else first unused BODY. Filled BEFORE
    #    subtitle so that on layouts with only one content slot (e.g. Section
    #    Header: TITLE+BODY), the bullets win the slot rather than getting
    #    dropped by a one-line subtitle.
    if bullets:
        chosen = None
        for ptype in (7, 2):  # OBJECT preferred, then BODY
            for ph in slide.placeholders:
                if (ph.placeholder_format.type == ptype
                    and ph.has_text_frame
                    and ph.placeholder_format.idx not in used_phs):
                    chosen = ph
                    break
            if chosen is not None:
                break
        if chosen is not None:
            _set_bullets(chosen.text_frame, bullets)
            used_phs.add(chosen.placeholder_format.idx)

    # 3. Subtitle — first SUBTITLE; if absent, first remaining unused BODY.
    if subtitle_text:
        chosen = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type in _SUBTITLE_TYPES and ph.has_text_frame:
                chosen = ph
                break
        if chosen is None:
            for ph in slide.placeholders:
                if (ph.placeholder_format.type == 2  # BODY
                    and ph.has_text_frame
                    and ph.placeholder_format.idx not in used_phs):
                    chosen = ph
                    break
        if chosen is not None:
            _set_text_frame(chosen.text_frame, subtitle_text)
            used_phs.add(chosen.placeholder_format.idx)

    # 4. Clear any remaining unused content placeholders so sample text
    #    ("Bullet level 1…", "Helvetica Bold 24pt") doesn't ship in the deck.
    for ph in slide.placeholders:
        ptype = ph.placeholder_format.type
        if ptype in _SKIP_TYPES:
            continue
        if not ph.has_text_frame:
            continue
        if ph.placeholder_format.idx in used_phs:
            continue
        # Subtitle-only slides whose subtitle we already filled are handled.
        # This catches secondary BODY/OBJECT slots on Comparison layouts, etc.
        if ptype in (2, 4, 7):  # BODY, SUBTITLE, OBJECT
            _clear_text_frame(ph.text_frame)


# Box geometry (left, top, width, height) as fractions of slide dimensions, keyed by
# (approx_position, size). Bottom positions clear the top ~22% to leave title room.
_REGION_BOXES: dict[tuple[str, str], tuple[float, float, float, float]] = {
    ("top", "small"):    (0.375, 0.05, 0.25, 0.30),
    ("top", "medium"):   (0.30,  0.05, 0.40, 0.40),
    ("top", "large"):    (0.20,  0.05, 0.60, 0.50),
    ("bottom", "small"): (0.375, 0.55, 0.25, 0.40),
    ("bottom", "medium"):(0.30,  0.42, 0.40, 0.55),
    ("bottom", "large"): (0.10,  0.35, 0.80, 0.60),
    ("center", "small"): (0.375, 0.30, 0.25, 0.40),
    ("center", "medium"):(0.30,  0.22, 0.40, 0.55),
    ("center", "large"): (0.20,  0.18, 0.60, 0.65),
    ("left", "small"):   (0.05,  0.30, 0.30, 0.40),
    ("left", "medium"):  (0.05,  0.22, 0.40, 0.55),
    ("left", "large"):   (0.05,  0.20, 0.45, 0.70),
    ("right", "small"):  (0.65,  0.30, 0.30, 0.40),
    ("right", "medium"): (0.55,  0.22, 0.40, 0.55),
    ("right", "large"):  (0.50,  0.20, 0.45, 0.70),
}

# Used when the template description has no image_regions for the chosen slide.
_FALLBACK_BOX = (0.55, 0.22, 0.40, 0.68)


def _resolve_figure_box(region_hint: Optional[dict]) -> tuple[float, float, float, float]:
    """Pick a normalized (left, top, width, height) box for the figure.

    Prefers the slide's described image_regions[0]; falls back to a right-side box.
    """
    if not region_hint:
        return _FALLBACK_BOX
    pos = str(region_hint.get("approx_position", "")).strip().lower()
    size = str(region_hint.get("size", "medium")).strip().lower()
    if size not in {"small", "medium", "large"}:
        size = "medium"
    return _REGION_BOXES.get((pos, size), _FALLBACK_BOX)


def _insert_figure(
    slide,
    prs,
    figure_path: Path,
    figure_aspect: float,
    region_hint: Optional[dict] = None,
) -> None:
    """Insert a figure on the slide, aspect-preserved.

    v2: uses the template structure JSON's `image_regions[0]` (approx_position +
    size from the vision LLM) when present, falling back to a right-side box.
    """
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    box_l, box_t, box_w_frac, box_h_frac = _resolve_figure_box(region_hint)
    box_left = int(slide_w * box_l)
    box_top = int(slide_h * box_t)
    box_w = int(slide_w * box_w_frac)
    box_h = int(slide_h * box_h_frac)
    box_aspect = box_w / box_h if box_h else 1.0

    if figure_aspect > box_aspect:
        # Figure is wider relative to the box → fit by width, center vertically.
        pic = slide.shapes.add_picture(str(figure_path), box_left, box_top, width=box_w)
        pic.top = box_top + (box_h - pic.height) // 2
    else:
        # Figure is taller relative to the box → fit by height, center horizontally.
        pic = slide.shapes.add_picture(str(figure_path), box_left, box_top, height=box_h)
        pic.left = box_left + (box_w - pic.width) // 2


def render_output_pptx(
    template_path: Path,
    plan: dict[str, Any],
    output_path: Path,
    figures: Optional[list[dict[str, Any]]] = None,
    template_descriptions: Optional[list[dict[str, Any]]] = None,
) -> Path:
    prs = Presentation(str(template_path))
    n_template_slides = len(prs.slides)

    plan_slides = plan.get("slides") or []
    if not plan_slides:
        print("[template-pipeline]   WARN: empty plan; saving template unchanged.")
        prs.save(str(output_path))
        return output_path

    fig_by_id = {f["id"]: f for f in (figures or [])}
    descs_by_idx = {d.get("slide_index", i): d for i, d in enumerate(template_descriptions or [])}

    new_slide_indices = []
    figures_inserted = 0
    for spec in plan_slides:
        src = spec.get("template_slide_index")
        if src is None or not (0 <= src < n_template_slides):
            print(f"[template-pipeline]   WARN: skipping slide with bad template_slide_index={src}")
            continue
        new_slide = _duplicate_slide(prs, src)
        try:
            _fill_slide_text(new_slide, spec.get("content") or {})
        except Exception as e:
            print(f"[template-pipeline]   WARN: text fill failed on output_index={spec.get('output_index')}: {e}")

        figure_id = spec.get("figure_id")
        if figure_id is not None and figure_id in fig_by_id:
            fig = fig_by_id[figure_id]
            try:
                fig_path = Path(fig["image_path"])
                if not fig_path.is_absolute():
                    fig_path = REPO_ROOT / fig_path
                if fig_path.exists():
                    # Look up the template slide's first image_region, if any.
                    desc = descs_by_idx.get(src, {})
                    regions = desc.get("image_regions") or []
                    region_hint = regions[0] if regions else None
                    _insert_figure(new_slide, prs, fig_path, fig.get("aspect", 1.0), region_hint)
                    figures_inserted += 1
                else:
                    print(f"[template-pipeline]   WARN: figure file missing: {fig_path}")
            except Exception as e:
                print(f"[template-pipeline]   WARN: figure insert failed on output_index={spec.get('output_index')}: {e}")

        new_slide_indices.append(len(prs.slides) - 1)

    # Now delete the original template slides (indices 0..n_template_slides-1)
    _delete_slides_at(prs, list(range(n_template_slides)))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[template-pipeline]   rendered {len(new_slide_indices)} slides ({figures_inserted} with figures) -> {output_path}")
    return output_path


# ---------- Top-level entry ---------- #


def generate(args, user_template_path: str) -> Optional[Path]:
    paper_path = Path(args.paper_path).resolve()
    template_path = Path(user_template_path).resolve()

    if not paper_path.exists():
        print(f"[template-pipeline] ERROR: paper not found: {paper_path}")
        return None
    if not template_path.exists():
        print(f"[template-pipeline] ERROR: template not found: {template_path}")
        return None

    problems = _check_prerequisites(template_path)
    if problems:
        print("[template-pipeline] ERROR: pre-flight checks failed:")
        for p in problems:
            print(f"  - {p}")
        return None

    output_dir = REPO_ROOT / "contents" / args.paper_name
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "template_aware_output.pptx"

    print(f"[template-pipeline] paper:    {paper_path}")
    print(f"[template-pipeline] template: {template_path}")
    print(f"[template-pipeline] output:   {output_path}")

    pipeline_start = time.time()

    try:
        with _stage(1, "parsing paper (markdown + figures)"):
            paper_markdown, figures = parse_paper(paper_path, output_dir)
            print(f"[template-pipeline]   {len(figures)} captioned figures available")

        template_images_dir = output_dir / "_template_images"
        structure_path = output_dir / "_template_structure.json"

        if structure_path.exists():
            print(f"[template-pipeline] (stages 2-3 cached) loading {structure_path}")
            descriptions = json.loads(structure_path.read_text(encoding="utf-8"))
        else:
            with _stage(2, "rendering template to images"):
                image_paths = render_template_to_images(template_path, template_images_dir)
                print(f"[template-pipeline]   produced {len(image_paths)} slide images")
            with _stage(3, "describing template slides via vision LLM"):
                descriptions = describe_template_slides(image_paths)
                with structure_path.open("w") as f:
                    json.dump(descriptions, f, indent=2)
                print(f"[template-pipeline]   wrote: {structure_path}")

        author_profile = _load_or_distill_author_profile(args)

        plan_path = output_dir / "_slide_plan.json"
        if plan_path.exists() and not getattr(args, "force_replan", False):
            print(f"[template-pipeline] (stage 4 cached) loading {plan_path}")
            plan = json.loads(plan_path.read_text(encoding="utf-8"))
        else:
            label = "planning output slides" + (" [with author profile]" if author_profile else "")
            with _stage(4, label):
                plan = plan_output_slides(paper_markdown, descriptions, figures=figures, author_profile=author_profile)
                with plan_path.open("w") as f:
                    json.dump(plan, f, indent=2)
                print(f"[template-pipeline]   planned {len(plan.get('slides', []))} output slides; wrote: {plan_path}")

        with _stage(5, "rendering output PPTX"):
            render_output_pptx(
                template_path, plan, output_path,
                figures=figures, template_descriptions=descriptions,
            )
            print(f"[template-pipeline] wrote: {output_path}")

        total = time.time() - pipeline_start
        print(f"[template-pipeline] DONE — total elapsed {total:.1f}s")
        return output_path

    except KeyboardInterrupt:
        print("[template-pipeline] interrupted by user.")
        return None
    except Exception as e:
        print(f"[template-pipeline] ERROR: pipeline failed: {type(e).__name__}: {e}")
        traceback.print_exc()
        return None


def _load_or_distill_author_profile(args) -> Optional[dict[str, Any]]:
    """Load an author profile for planner injection, distilling on demand if needed.

    Mirrors Branch A's flow in new_pipeline_logtime.py: respects --author_profile_path
    if given, else looks under Capstone/profiles/<author_id>.json, else calls
    Capstone.preference_distill.distill_author_profile to build one.
    """
    if not getattr(args, "use_author_preferences", False):
        return None
    author_id = getattr(args, "author_id", None)
    if not author_id:
        print("[template-pipeline]   WARN: --use_author_preferences set without --author_id; skipping profile.")
        return None

    explicit_path = getattr(args, "author_profile_path", None)
    if explicit_path:
        profile_path = Path(explicit_path)
    else:
        profile_path = REPO_ROOT / "Capstone" / "profiles" / f"{author_id}.json"

    if profile_path.exists() and not getattr(args, "force_refresh_preferences", False):
        print(f"[template-pipeline]   loading author profile: {profile_path}")
        return json.loads(profile_path.read_text(encoding="utf-8"))

    try:
        from Capstone.preference_distill import distill_author_profile
        print(f"[template-pipeline]   distilling author profile for {author_id}...")
        profile = distill_author_profile(
            author_id,
            output_dir=profile_path.parent,
            max_papers=getattr(args, "preference_max_papers", 5),
            model=getattr(args, "preference_model", None) or getattr(args, "model_name_t", None),
            force_refresh=getattr(args, "force_refresh_preferences", False),
        )
        profile_path.parent.mkdir(parents=True, exist_ok=True)
        profile_path.write_text(json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")
        return profile
    except Exception as e:
        print(f"[template-pipeline]   WARN: profile distillation failed ({e}); planning without profile.")
        return None
