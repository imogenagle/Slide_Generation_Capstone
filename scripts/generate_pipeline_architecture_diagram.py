from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parents[1] / "docs" / "slidegen_pipeline_architecture.pptx"


BG = RGBColor(248, 250, 252)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(30, 64, 175)
BLUE_LIGHT = RGBColor(219, 234, 254)
GREEN = RGBColor(22, 101, 52)
GREEN_LIGHT = RGBColor(220, 252, 231)
AMBER = RGBColor(146, 64, 14)
AMBER_LIGHT = RGBColor(254, 243, 199)
PURPLE = RGBColor(109, 40, 217)
PURPLE_LIGHT = RGBColor(237, 233, 254)
RED = RGBColor(153, 27, 27)
RED_LIGHT = RGBColor(254, 226, 226)
GRAY = RGBColor(100, 116, 139)
GRAY_LIGHT = RGBColor(241, 245, 249)


def add_box(slide, left, top, width, height, title, body, fill, line, title_size=18, body_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = Pt(title_size)
    try:
        r1.font.color.rgb = TEXT
    except Exception:
        pass

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    for idx, line_text in enumerate(body.split("\n")):
        if idx > 0:
            p = tf.add_paragraph()
        else:
            p = p2
        r = p.add_run()
        r.text = line_text
        r.font.size = Pt(body_size)
        try:
            r.font.color.rgb = MUTED
        except Exception:
            pass
    return shape


def add_label(slide, left, top, width, height, text, size=12, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    try:
        r.font.color.rgb = color
    except Exception:
        pass
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.25):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_elbow(slide, x1, y1, x2, y2, color=GRAY, width=2.25):
    line = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def center_x(shape):
    return shape.left + shape.width / 2


def center_y(shape):
    return shape.top + shape.height / 2


def bottom(shape):
    return shape.top + shape.height


def right(shape):
    return shape.left + shape.width


def build_diagram():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_label(
        slide,
        Inches(0.35),
        Inches(0.12),
        Inches(12.6),
        Inches(0.4),
        "SlideGen Architecture: Baseline vs Retrieval-Personalized Pipeline",
        size=24,
        color=TEXT,
        bold=True,
    )
    add_label(
        slide,
        Inches(0.37),
        Inches(0.5),
        Inches(12.4),
        Inches(0.25),
        "Editable overview of the current pipeline entrypoints, agents, shared extraction stages, personalization branch, and evaluation.",
        size=11,
        color=MUTED,
    )

    # Offline profile branch
    add_label(slide, Inches(0.35), Inches(0.88), Inches(2.5), Inches(0.2), "Offline retrieval profile builder", size=13, color=PURPLE, bold=True)
    hist = add_box(
        slide,
        Inches(0.35),
        Inches(1.1),
        Inches(2.1),
        Inches(0.95),
        "Historical author decks",
        "Retrieved prior slide PNGs / deck folders\nClosest prior talks for the same author",
        PURPLE_LIGHT,
        PURPLE,
    )
    profile = add_box(
        slide,
        Inches(2.75),
        Inches(1.05),
        Inches(2.55),
        Inches(1.05),
        "retrieval_profile_pilot.py",
        "Retrieves related historical deck(s)\nDistills targets: slides, words, bullets, images/tables/formulas, sections, color, font",
        PURPLE_LIGHT,
        PURPLE,
    )
    profile_json = add_box(
        slide,
        Inches(5.62),
        Inches(1.1),
        Inches(2.0),
        Inches(0.95),
        "Retrieval profile JSON",
        "Author-conditioned target profile\nUsed only by personalized branch",
        PURPLE_LIGHT,
        PURPLE,
    )

    add_arrow(slide, right(hist), center_y(hist), profile.left, center_y(profile), color=PURPLE)
    add_arrow(slide, right(profile), center_y(profile), profile_json.left, center_y(profile_json), color=PURPLE)

    # Shared pipeline
    add_label(slide, Inches(0.35), Inches(2.35), Inches(4.0), Inches(0.2), "Shared paper ingestion + extraction", size=13, color=BLUE, bold=True)
    paper = add_box(
        slide,
        Inches(0.35),
        Inches(2.58),
        Inches(1.95),
        Inches(0.95),
        "Target paper PDF",
        "Input research paper\n`--paper_path`",
        BLUE_LIGHT,
        BLUE,
    )
    parse = add_box(
        slide,
        Inches(2.6),
        Inches(2.48),
        Inches(2.6),
        Inches(1.15),
        "parse_raw.py",
        "Docling convert + markdown export\nLLM outliner turns paper into raw slide content / outline JSON",
        BLUE_LIGHT,
        BLUE,
    )
    assets = add_box(
        slide,
        Inches(5.5),
        Inches(2.48),
        Inches(3.1),
        Inches(1.15),
        "Asset extraction agents",
        "gen_image_and_table\nfilter_image_table\ngen_figure_match\ngen_formula_match_v1",
        BLUE_LIGHT,
        BLUE,
    )

    add_arrow(slide, right(paper), center_y(paper), parse.left, center_y(parse), color=BLUE)
    add_arrow(slide, right(parse), center_y(parse), assets.left, center_y(assets), color=BLUE)

    artifact = add_box(
        slide,
        Inches(8.95),
        Inches(2.48),
        Inches(1.9),
        Inches(1.15),
        "Shared artifacts",
        "Raw content JSON\nFigures JSON\nFormula match JSON\nFiltered images/tables",
        BLUE_LIGHT,
        BLUE,
    )
    add_arrow(slide, right(assets), center_y(assets), artifact.left, center_y(artifact), color=BLUE)

    # Split branches
    add_label(slide, Inches(0.35), Inches(3.95), Inches(2.5), Inches(0.2), "Generation branches", size=13, color=GREEN, bold=True)
    baseline = add_box(
        slide,
        Inches(0.45),
        Inches(4.2),
        Inches(3.2),
        Inches(1.2),
        "Baseline branch",
        "new_pipeline.py\n`personalization_mode=standard`\nlayout_agent_xin.py planner",
        GREEN_LIGHT,
        GREEN,
    )
    personalized = add_box(
        slide,
        Inches(4.2),
        Inches(4.2),
        Inches(3.65),
        Inches(1.35),
        "Retrieval-personalized branch",
        "new_pipeline.py\n`personalization_mode=retrieval`\nlayout_agent_xin_retrieval.py planner + repair loop\nConsumes retrieval profile JSON",
        AMBER_LIGHT,
        AMBER,
    )

    add_elbow(slide, center_x(artifact), bottom(artifact), center_x(baseline), baseline.top, color=GREEN)
    add_elbow(slide, center_x(artifact), bottom(artifact), center_x(personalized), personalized.top, color=AMBER)
    add_arrow(slide, center_x(profile_json), bottom(profile_json), center_x(personalized), personalized.top, color=PURPLE)

    planner_note = add_box(
        slide,
        Inches(4.2),
        Inches(5.72),
        Inches(3.65),
        Inches(0.92),
        "Personalization controls",
        "Targets guide slide count, text/bullet density, image/table/formula budgets,\nsection preferences, theme color, and font choices.",
        AMBER_LIGHT,
        AMBER,
        title_size=15,
        body_size=10,
    )
    add_arrow(slide, center_x(personalized), bottom(personalized), center_x(planner_note), planner_note.top, color=AMBER)

    # Rendering
    renderer = add_box(
        slide,
        Inches(8.28),
        Inches(4.28),
        Inches(2.35),
        Inches(1.15),
        "layout_filler.py",
        "Chooses templates\nPlaces text + visual assets\nApplies color/font\nWrites editable PPTX",
        GREEN_LIGHT,
        GREEN,
    )
    add_arrow(slide, right(baseline), center_y(baseline), renderer.left, center_y(renderer) - Pt(8), color=GREEN)
    add_arrow(slide, right(personalized), center_y(personalized), renderer.left, center_y(renderer) + Pt(8), color=AMBER)

    pptx_out = add_box(
        slide,
        Inches(10.95),
        Inches(4.28),
        Inches(2.0),
        Inches(1.15),
        "Outputs",
        "Baseline PPTX\nPersonalized PPTX\nSlide plans + logs",
        GREEN_LIGHT,
        GREEN,
    )
    add_arrow(slide, right(renderer), center_y(renderer), pptx_out.left, center_y(pptx_out), color=GREEN)

    # Evaluation
    add_label(slide, Inches(7.95), Inches(5.92), Inches(2.5), Inches(0.2), "Evaluation", size=13, color=RED, bold=True)
    eval_box = add_box(
        slide,
        Inches(8.05),
        Inches(6.15),
        Inches(4.9),
        Inches(1.0),
        "run_batch_experiment.py -> retrieval evaluators",
        "evaluate_retrieval_alignment_numeric.py\n"
        "evaluate_retrieval_alignment_sections_llm.py\n"
        "evaluate_retrieval_alignment_all.py\n"
        "Compares baseline vs personalized against profile targets",
        RED_LIGHT,
        RED,
    )
    add_arrow(slide, center_x(pptx_out), bottom(pptx_out), center_x(eval_box), eval_box.top, color=RED)

    # Footer note
    add_label(
        slide,
        Inches(0.4),
        Inches(7.08),
        Inches(12.2),
        Inches(0.22),
        "Main entrypoints: `SlidesAgent/new_pipeline.py` for single-paper generation and `Capstone/run_batch_experiment.py` for baseline vs personalized batch experiments.",
        size=10,
        color=GRAY,
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)


if __name__ == "__main__":
    build_diagram()
