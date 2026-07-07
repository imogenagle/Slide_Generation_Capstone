#!/usr/bin/env python3
"""Per-target vote evaluator for retrieval-profile personalization runs.

This script is intentionally retrieval-specific: it compares baseline and
personalized slide plans only on the numeric targets used by the retrieval
personalization path, and determines the overall winner by majority vote
across targets.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import Counter
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE


REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from SlidesAgent.slide_plan_summary import summarize_slide_plan


RETRIEVAL_METRIC_SPECS = {
    "target_slide_count": {
        "label": "slide_count",
        "summary_key": "slide_count",
        "normalizer_floor": 1.0,
    },
    "target_avg_bullets_per_slide": {
        "label": "avg_bullets_per_slide",
        "summary_key": "avg_bullets_per_slide",
        "normalizer_floor": 1.0,
    },
    "target_avg_words_per_slide": {
        "label": "avg_words_per_slide",
        "summary_key": "avg_words_per_slide",
        "normalizer_floor": 8.0,
    },
    "target_image_slide_count": {
        "label": "image_slide_count",
        "summary_key": "image_slide_count",
        "normalizer_floor": 1.0,
    },
    "target_table_slide_count": {
        "label": "table_slide_count",
        "summary_key": "table_slide_count",
        "normalizer_floor": 1.0,
    },
    "target_formula_slide_count": {
        "label": "formula_slide_count",
        "summary_key": "formula_slide_count",
        "normalizer_floor": 1.0,
    },
}

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def build_retrieval_target_summary(profile: dict[str, Any]) -> dict[str, float]:
    numeric = dict(profile.get("numeric_preferences") or {})
    summary: dict[str, float] = {}
    for key in RETRIEVAL_METRIC_SPECS:
        raw_value = numeric.get(key)
        if raw_value is None:
            continue
        try:
            summary[key] = float(raw_value)
        except Exception:
            continue
    return summary


def build_color_target_summary(profile: dict[str, Any]) -> dict[str, Any]:
    color_preferences = dict(profile.get("color_preferences") or {})
    target_theme_hex = str(color_preferences.get("target_theme_hex") or "").strip().upper()
    if not target_theme_hex:
        return {}
    return {
        "target_theme_hex": target_theme_hex,
        "target_base_hex": str(color_preferences.get("target_base_hex") or "").strip().upper() or None,
        "color_source_paper_id": color_preferences.get("color_source_paper_id"),
        "color_source_paper_title": color_preferences.get("color_source_paper_title"),
    }


def build_font_target_summary(profile: dict[str, Any]) -> dict[str, Any]:
    font_preferences = dict(profile.get("font_preferences") or {})
    title_font_name = str(font_preferences.get("title_font_name") or "").strip()
    body_font_name = str(font_preferences.get("body_font_name") or "").strip()
    result: dict[str, Any] = {}
    if title_font_name:
        result["title_font_name"] = title_font_name
    if body_font_name:
        result["body_font_name"] = body_font_name
    return result


def _normalize_hex(value: str | None) -> str | None:
    raw = str(value or "").strip().upper()
    if not raw:
        return None
    if raw.startswith("#"):
        raw = raw[1:]
    if len(raw) != 6:
        return None
    try:
        int(raw, 16)
    except ValueError:
        return None
    return f"#{raw}"


def _hex_to_rgb(value: str | None) -> tuple[int, int, int] | None:
    normalized = _normalize_hex(value)
    if normalized is None:
        return None
    return (
        int(normalized[1:3], 16),
        int(normalized[3:5], 16),
        int(normalized[5:7], 16),
    )


def _color_distance(a: str | None, b: str | None) -> float | None:
    rgb_a = _hex_to_rgb(a)
    rgb_b = _hex_to_rgb(b)
    if rgb_a is None or rgb_b is None:
        return None
    return round(sum((float(x) - float(y)) ** 2 for x, y in zip(rgb_a, rgb_b)) ** 0.5, 4)


def _normalize_font_name(value: str | None) -> str | None:
    text = str(value or "").strip()
    if not text:
        return None
    lowered = text.lower()
    lowered = lowered.replace("-", " ")
    lowered = " ".join(lowered.split())
    return lowered or None


def _font_match_distance(observed: str | None, target: str | None) -> float | None:
    observed_norm = _normalize_font_name(observed)
    target_norm = _normalize_font_name(target)
    if observed_norm is None or target_norm is None:
        return None
    if observed_norm == target_norm:
        return 0.0
    if observed_norm in target_norm or target_norm in observed_norm:
        return 0.25
    observed_tokens = set(observed_norm.split())
    target_tokens = set(target_norm.split())
    if observed_tokens & target_tokens:
        return 0.5
    return 1.0


def infer_themed_pptx_path(plan_path: Path) -> Path | None:
    parent = plan_path.parent
    stem = plan_path.name
    variant = None
    if "_slide_plan_" in stem:
        variant = stem.split("_slide_plan_", 1)[1].rsplit(".json", 1)[0]
    if not variant:
        return None
    matches = sorted(parent.glob(f"*_output_slides_{variant}_themed.pptx"))
    if matches:
        return matches[0]
    fallback = sorted(parent.glob(f"*_output_slides_{variant}.pptx"))
    return fallback[0] if fallback else None


def extract_font_preferences_from_pptx(pptx_path: Path) -> dict[str, Any]:
    if not pptx_path.exists():
        return {}

    prs = Presentation(str(pptx_path))
    title_fonts: Counter[str] = Counter()
    body_fonts: Counter[str] = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER or not getattr(shape, "is_placeholder", False):
                placeholder_type = None
            else:
                placeholder_type = shape.placeholder_format.type

            bucket = body_fonts
            if placeholder_type in {PH_TYPE.TITLE, PH_TYPE.CENTER_TITLE}:
                bucket = title_fonts

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_name = str(getattr(run.font, "name", "") or "").strip()
                    if font_name:
                        bucket[font_name] += len((run.text or "").strip()) or 1
                if not paragraph.runs:
                    para_font_name = str(getattr(paragraph.font, "name", "") or "").strip()
                    if para_font_name:
                        bucket[para_font_name] += len((paragraph.text or "").strip()) or 1

    result: dict[str, Any] = {}
    if title_fonts:
        result["title_font_name"] = title_fonts.most_common(1)[0][0]
    if body_fonts:
        result["body_font_name"] = body_fonts.most_common(1)[0][0]
    return result


def extract_theme_color_from_pptx(pptx_path: Path, *, target_key: str = "dk2") -> str | None:
    if not pptx_path.exists():
        return None
    with ZipFile(pptx_path) as archive:
        theme_paths = [name for name in archive.namelist() if name.startswith("ppt/theme/") and name.endswith(".xml")]
        for theme_path in theme_paths:
            root = etree.fromstring(archive.read(theme_path))
            clr_scheme = root.find(".//a:clrScheme", namespaces=NS)
            if clr_scheme is None:
                continue
            target_elem = clr_scheme.find(f"a:{target_key}", namespaces=NS)
            if target_elem is None:
                continue
            srgb = target_elem.find("a:srgbClr", namespaces=NS)
            if srgb is not None and srgb.get("val"):
                return _normalize_hex(srgb.get("val"))
            sys_clr = target_elem.find("a:sysClr", namespaces=NS)
            if sys_clr is not None:
                return _normalize_hex(sys_clr.get("lastClr") or sys_clr.get("val"))
    return None


def build_color_comparison(
    *,
    profile: dict[str, Any],
    baseline_plan_path: Path,
    personalized_plan_path: Path,
) -> dict[str, Any] | None:
    target_summary = build_color_target_summary(profile)
    target_theme_hex = target_summary.get("target_theme_hex")
    if not target_theme_hex:
        return None

    baseline_pptx = infer_themed_pptx_path(baseline_plan_path)
    personalized_pptx = infer_themed_pptx_path(personalized_plan_path)
    baseline_theme_hex = extract_theme_color_from_pptx(baseline_pptx) if baseline_pptx else None
    personalized_theme_hex = extract_theme_color_from_pptx(personalized_pptx) if personalized_pptx else None

    baseline_distance = _color_distance(baseline_theme_hex, target_theme_hex)
    personalized_distance = _color_distance(personalized_theme_hex, target_theme_hex)

    if baseline_distance is None and personalized_distance is None:
        winner = "tie"
    elif baseline_distance is None:
        winner = "personalized"
    elif personalized_distance is None:
        winner = "baseline"
    elif abs(baseline_distance - personalized_distance) <= 1e-9:
        winner = "tie"
    elif personalized_distance < baseline_distance:
        winner = "personalized"
    else:
        winner = "baseline"

    return {
        "target_theme_hex": target_theme_hex,
        "target_base_hex": target_summary.get("target_base_hex"),
        "baseline_theme_hex": baseline_theme_hex,
        "personalized_theme_hex": personalized_theme_hex,
        "baseline_distance": baseline_distance,
        "personalized_distance": personalized_distance,
        "closer_to_palette_preferences": winner,
        "baseline_themed_pptx": str(baseline_pptx) if baseline_pptx else None,
        "personalized_themed_pptx": str(personalized_pptx) if personalized_pptx else None,
    }


def build_font_comparison(
    *,
    profile: dict[str, Any],
    baseline_plan_path: Path,
    personalized_plan_path: Path,
) -> dict[str, Any] | None:
    target_summary = build_font_target_summary(profile)
    if not target_summary:
        return None

    baseline_pptx = infer_themed_pptx_path(baseline_plan_path)
    personalized_pptx = infer_themed_pptx_path(personalized_plan_path)
    baseline_fonts = extract_font_preferences_from_pptx(baseline_pptx) if baseline_pptx else {}
    personalized_fonts = extract_font_preferences_from_pptx(personalized_pptx) if personalized_pptx else {}

    baseline_title_distance = _font_match_distance(baseline_fonts.get("title_font_name"), target_summary.get("title_font_name"))
    personalized_title_distance = _font_match_distance(personalized_fonts.get("title_font_name"), target_summary.get("title_font_name"))
    baseline_body_distance = _font_match_distance(baseline_fonts.get("body_font_name"), target_summary.get("body_font_name"))
    personalized_body_distance = _font_match_distance(personalized_fonts.get("body_font_name"), target_summary.get("body_font_name"))

    baseline_total = sum(value for value in (baseline_title_distance, baseline_body_distance) if value is not None)
    personalized_total = sum(value for value in (personalized_title_distance, personalized_body_distance) if value is not None)
    baseline_available = any(value is not None for value in (baseline_title_distance, baseline_body_distance))
    personalized_available = any(value is not None for value in (personalized_title_distance, personalized_body_distance))

    if not baseline_available and not personalized_available:
        winner = "tie"
    elif not baseline_available:
        winner = "personalized"
    elif not personalized_available:
        winner = "baseline"
    elif abs(baseline_total - personalized_total) <= 1e-9:
        winner = "tie"
    elif personalized_total < baseline_total:
        winner = "personalized"
    else:
        winner = "baseline"

    return {
        "target_title_font_name": target_summary.get("title_font_name"),
        "target_body_font_name": target_summary.get("body_font_name"),
        "baseline_title_font_name": baseline_fonts.get("title_font_name"),
        "baseline_body_font_name": baseline_fonts.get("body_font_name"),
        "personalized_title_font_name": personalized_fonts.get("title_font_name"),
        "personalized_body_font_name": personalized_fonts.get("body_font_name"),
        "baseline_title_distance": baseline_title_distance,
        "personalized_title_distance": personalized_title_distance,
        "baseline_body_distance": baseline_body_distance,
        "personalized_body_distance": personalized_body_distance,
        "closer_to_font_preferences": winner,
        "baseline_themed_pptx": str(baseline_pptx) if baseline_pptx else None,
        "personalized_themed_pptx": str(personalized_pptx) if personalized_pptx else None,
    }


def build_metric_comparison(
    *,
    target_summary: dict[str, float],
    baseline_summary: dict[str, Any],
    personalized_summary: dict[str, Any],
) -> dict[str, Any]:
    metrics: dict[str, Any] = {}
    personalized_metric_wins = 0
    baseline_metric_wins = 0
    tied_metrics = 0

    for target_key, target_value in target_summary.items():
        spec = RETRIEVAL_METRIC_SPECS[target_key]
        summary_key = str(spec["summary_key"])
        baseline_value = float(baseline_summary.get(summary_key, 0.0))
        personalized_value = float(personalized_summary.get(summary_key, 0.0))

        baseline_distance = abs(baseline_value - target_value)
        personalized_distance = abs(personalized_value - target_value)
        normalizer = max(
            abs(target_value),
            baseline_distance,
            personalized_distance,
            float(spec["normalizer_floor"]),
        )
        if abs(baseline_distance - personalized_distance) <= 1e-9:
            closer = "tie"
            tied_metrics += 1
        elif personalized_distance < baseline_distance:
            closer = "personalized"
            personalized_metric_wins += 1
        else:
            closer = "baseline"
            baseline_metric_wins += 1

        raw_gain = baseline_distance - personalized_distance

        metrics[target_key] = {
            "label": spec["label"],
            "target": round(target_value, 4),
            "baseline": round(baseline_value, 4),
            "personalized": round(personalized_value, 4),
            "baseline_distance": round(baseline_distance, 4),
            "personalized_distance": round(personalized_distance, 4),
            "distance_improvement": round(raw_gain, 4),
            "closer_to_target": closer,
        }

    metric_count = len(metrics)

    if personalized_metric_wins == baseline_metric_wins:
        metric_vote_winner = "tie"
    elif personalized_metric_wins > baseline_metric_wins:
        metric_vote_winner = "personalized"
    else:
        metric_vote_winner = "baseline"

    return {
        "metrics": metrics,
        "aggregate": {
            "metric_count": metric_count,
            "personalized_metric_wins": personalized_metric_wins,
            "baseline_metric_wins": baseline_metric_wins,
            "tied_metrics": tied_metrics,
            "winner": metric_vote_winner,
        },
    }


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Evaluate baseline vs personalized plans against retrieval-profile numeric targets."
    )
    parser.add_argument("--profile", type=Path, required=True)
    parser.add_argument("--baseline-plan", type=Path, required=True)
    parser.add_argument("--personalized-plan", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=None)
    args = parser.parse_args()

    for path_arg in (args.profile, args.baseline_plan, args.personalized_plan):
        if not path_arg.exists():
            raise FileNotFoundError(f"Required input not found: {path_arg}")

    profile = load_json(args.profile)
    baseline_plan = load_json(args.baseline_plan)
    personalized_plan = load_json(args.personalized_plan)

    target_summary = build_retrieval_target_summary(profile)
    baseline_summary = summarize_slide_plan(baseline_plan)
    personalized_summary = summarize_slide_plan(personalized_plan)
    baseline_image_slide_count = sum(1 for slide in list(baseline_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("images"))
    baseline_table_slide_count = sum(1 for slide in list(baseline_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("tables"))
    baseline_formula_slide_count = sum(1 for slide in list(baseline_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("formulas"))
    personalized_image_slide_count = sum(1 for slide in list(personalized_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("images"))
    personalized_table_slide_count = sum(1 for slide in list(personalized_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("tables"))
    personalized_formula_slide_count = sum(1 for slide in list(personalized_plan.get("slides") or []) if isinstance(slide, dict) and slide.get("formulas"))
    baseline_summary["image_slide_count"] = baseline_image_slide_count
    baseline_summary["table_slide_count"] = baseline_table_slide_count
    baseline_summary["formula_slide_count"] = baseline_formula_slide_count
    personalized_summary["image_slide_count"] = personalized_image_slide_count
    personalized_summary["table_slide_count"] = personalized_table_slide_count
    personalized_summary["formula_slide_count"] = personalized_formula_slide_count

    comparison = build_metric_comparison(
        target_summary=target_summary,
        baseline_summary=baseline_summary,
        personalized_summary=personalized_summary,
    )
    color_comparison = build_color_comparison(
        profile=profile,
        baseline_plan_path=args.baseline_plan,
        personalized_plan_path=args.personalized_plan,
    )
    font_comparison = build_font_comparison(
        profile=profile,
        baseline_plan_path=args.baseline_plan,
        personalized_plan_path=args.personalized_plan,
    )

    report = {
        "inputs": {
            "profile": str(args.profile),
            "baseline_plan": str(args.baseline_plan),
            "personalized_plan": str(args.personalized_plan),
            "scoring_mode": "retrieval_numeric_target_distance",
        },
        "target_summary": target_summary,
        "observed_metrics": {
            "baseline": {
                "slide_count": int(baseline_summary.get("slide_count", 0)),
                "avg_bullets_per_slide": round(float(baseline_summary.get("avg_bullets_per_slide", 0.0)), 4),
                "avg_words_per_slide": round(float(baseline_summary.get("avg_words_per_slide", 0.0)), 4),
                "image_slide_count": baseline_image_slide_count,
                "table_slide_count": baseline_table_slide_count,
                "formula_slide_count": baseline_formula_slide_count,
            },
            "personalized": {
                "slide_count": int(personalized_summary.get("slide_count", 0)),
                "avg_bullets_per_slide": round(float(personalized_summary.get("avg_bullets_per_slide", 0.0)), 4),
                "avg_words_per_slide": round(float(personalized_summary.get("avg_words_per_slide", 0.0)), 4),
                "image_slide_count": personalized_image_slide_count,
                "table_slide_count": personalized_table_slide_count,
                "formula_slide_count": personalized_formula_slide_count,
            },
        },
        "comparison": comparison["metrics"],
        "color_palette_eval": color_comparison,
        "font_eval": font_comparison,
        "summary": comparison["aggregate"],
    }
    if color_comparison is not None:
        report["summary"]["color_palette_winner"] = color_comparison["closer_to_palette_preferences"]
    if font_comparison is not None:
        report["summary"]["font_winner"] = font_comparison["closer_to_font_preferences"]

    output_text = json.dumps(report, indent=2, ensure_ascii=False)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(output_text, encoding="utf-8")
    print(output_text)


if __name__ == "__main__":
    main()
