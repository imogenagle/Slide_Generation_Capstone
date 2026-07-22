#!/usr/bin/env python3
"""Batch test post-render user-prompt refinement.

This harness intentionally excludes LLM visual review. It tests the cheaper
user-prompt repair loop: prompt decisions -> slide-plan repair -> PPTX
regeneration -> deterministic validation.
"""

from __future__ import annotations

import argparse
import csv
import json
import os
import re
import sys
import time
from pathlib import Path
from types import SimpleNamespace
from typing import Any

try:
    from dotenv import load_dotenv
except ModuleNotFoundError:
    def load_dotenv(*_args: Any, **_kwargs: Any) -> bool:
        return False

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:
    raise SystemExit("python-pptx is required. Install project dependencies first.") from exc


REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from SlidesAgent.post_render_refinement import (
    DEFAULT_REPAIR_PROMPT,
    deterministic_render_edits,
    infer_asset_paper_name,
    infer_local_pptx,
    infer_local_slide_plan,
    infer_model_names_from_plan,
    infer_output_root_from_run_dir,
    inspect_deck,
    load_json,
    normalize_slide_ids,
    plan_index_for_slide_id,
    repair_deck,
    slide_change_summary,
    visual_counts,
)


DEFAULT_OUTPUT_ROOT = REPO_ROOT / "Capstone" / "evaluations" / "post_render_refinement"
DEFAULT_CASE_TEMPLATES = [
    {
        "case_id": "title_clarity",
        "target_selector": "second",
        "instruction": (
            "On {slide_id}, rewrite the visible slide title to be clearer and more direct. "
            "Keep the slide's core meaning and do not change other slides."
        ),
        "expected": {
            "target_slide_changed": True,
            "only_target_slide_changed": True,
            "title_changed": True,
            "pptx_target_title_visible": True,
            "no_external_asset_changes": True,
        },
    },
    {
        "case_id": "three_concise_bullets",
        "target_selector": "few_bullets",
        "instruction": (
            "On {slide_id}, keep exactly three concise bullets, preserve the main claim, "
            "and do not add external assets or change other slides."
        ),
        "expected": {
            "target_slide_changed": True,
            "only_target_slide_changed": True,
            "bullet_count": 3,
            "no_external_asset_changes": True,
        },
    },
    {
        "case_id": "general_technical_audience",
        "target_selector": "middle",
        "instruction": (
            "On {slide_id}, make this slide easier to follow for a general technical audience. "
            "Use plain wording, keep the title concise, and avoid adding assets or changing other slides."
        ),
        "expected": {
            "target_slide_changed": True,
            "only_target_slide_changed": True,
            "target_word_count_not_increased": True,
            "no_external_asset_changes": True,
        },
    },
    {
        "case_id": "scope_protection",
        "target_selector": "last",
        "instruction": (
            "On {slide_id}, tighten the title and bullets while preserving paper fidelity. "
            "This is a scope test: do not change any protected slide."
        ),
        "expected": {
            "target_slide_changed": True,
            "only_target_slide_changed": True,
            "protected_slides_unchanged": True,
            "max_bullet_count": 3,
            "no_external_asset_changes": True,
        },
    },
]


def safe_name(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(value or "")).strip("_")
    return cleaned or "item"


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def find_run_dirs(contents_dir: Path) -> list[Path]:
    if not contents_dir.exists():
        return []
    run_dirs: list[Path] = []
    for path in sorted(item for item in contents_dir.iterdir() if item.is_dir()):
        try:
            infer_local_pptx(path)
            infer_local_slide_plan(path)
        except Exception:
            continue
        run_dirs.append(path)
    return run_dirs


def load_manifest(path: Path) -> dict[str, Any]:
    manifest = load_json(path, {})
    if not isinstance(manifest, dict):
        raise ValueError(f"Manifest must be a JSON object: {path}")
    return manifest


def normalize_deck_specs(args: argparse.Namespace) -> list[dict[str, Any]]:
    manifest: dict[str, Any] = {}
    explicit_specs = bool(args.manifest or args.run_dir)
    if args.manifest:
        manifest = load_manifest(args.manifest)

    case_templates = manifest.get("case_templates") or DEFAULT_CASE_TEMPLATES
    deck_specs: list[dict[str, Any]] = []

    for raw_deck in manifest.get("decks") or []:
        if not isinstance(raw_deck, dict) or not raw_deck.get("run_dir"):
            continue
        deck_specs.append(
            {
                "run_dir": Path(raw_deck["run_dir"]),
                "pptx_path": Path(raw_deck["pptx_path"]) if raw_deck.get("pptx_path") else None,
                "slide_plan_path": Path(raw_deck["slide_plan_path"]) if raw_deck.get("slide_plan_path") else None,
                "case_templates": raw_deck.get("cases") or case_templates,
                "author_profile_path": Path(raw_deck["author_profile_path"]) if raw_deck.get("author_profile_path") else args.author_profile_path,
            }
        )

    for run_dir in args.run_dir or []:
        deck_specs.append(
            {
                "run_dir": run_dir,
                "case_templates": case_templates,
                "author_profile_path": args.author_profile_path,
            }
        )

    if not deck_specs:
        for run_dir in find_run_dirs(args.contents_dir):
            deck_specs.append(
                {
                    "run_dir": run_dir,
                    "case_templates": case_templates,
                    "author_profile_path": args.author_profile_path,
                }
            )

    if args.limit_decks > 0:
        deck_specs = deck_specs[: args.limit_decks]
    if args.require_deck_count > 0 and len(deck_specs) < args.require_deck_count:
        raise ValueError(
            f"Found {len(deck_specs)} qualifying deck(s), but --require-deck-count "
            f"requires {args.require_deck_count}. Generate more source decks or pass "
            "explicit --run-dir values."
        )
    if not explicit_specs and args.limit_decks > 0 and len(deck_specs) < args.limit_decks and not args.allow_fewer_decks:
        raise ValueError(
            f"Auto-discovery found only {len(deck_specs)} qualifying deck(s), but "
            f"--limit-decks is {args.limit_decks}. Re-run with --allow-fewer-decks "
            "for a partial batch, or generate/pass more source deck folders."
        )
    return deck_specs


def slide_ids(plan: dict[str, Any]) -> list[str]:
    normalized = normalize_slide_ids(plan)
    return [str(slide.get("slide_id")) for slide in normalized.get("slides") or [] if isinstance(slide, dict)]


def select_target_slide_id(plan: dict[str, Any], selector: str | None) -> str:
    ids = slide_ids(plan)
    if not ids:
        raise ValueError("Slide plan has no slides.")
    selector = selector or "middle"
    slides = list(normalize_slide_ids(plan).get("slides") or [])

    if selector in ids:
        return selector
    if selector == "first":
        return ids[0]
    if selector == "second":
        return ids[1] if len(ids) > 1 else ids[0]
    if selector == "middle":
        return ids[len(ids) // 2]
    if selector == "last":
        return ids[-1]
    if selector == "few_bullets":
        for slide in slides:
            bullets = [item for item in slide.get("bullets") or [] if isinstance(item, dict)]
            if len(bullets) != 3:
                return str(slide.get("slide_id"))
        return ids[len(ids) // 2]
    if selector == "visual":
        for slide in slides:
            counts = visual_counts(slide)
            if counts["image_count"] + counts["table_count"] + counts["formula_count"] > 0:
                return str(slide.get("slide_id"))
        return ids[len(ids) // 2]
    raise ValueError(f"Unknown target_selector: {selector}")


def resolve_case_slide_ids(plan: dict[str, Any], case: dict[str, Any]) -> tuple[list[str], str]:
    if case.get("target_slide_ids"):
        targets = [str(item) for item in case.get("target_slide_ids") or []]
    else:
        targets = [
            str(case.get("target_slide_id") or select_target_slide_id(plan, case.get("target_selector")))
        ]
    if not targets:
        raise ValueError("Case must resolve at least one target slide.")
    active_slide_id = str(case.get("active_slide_id") or targets[0])
    return targets, active_slide_id


def bullet_texts(slide: dict[str, Any]) -> list[str]:
    return [str(item.get("text") or "") for item in slide.get("bullets") or [] if isinstance(item, dict)]


def slide_word_count(slide: dict[str, Any]) -> int:
    parts = [str(slide.get("section") or ""), str(slide.get("subsection") or slide.get("title") or "")]
    for bullet in slide.get("bullets") or []:
        if not isinstance(bullet, dict):
            continue
        parts.append(str(bullet.get("text") or ""))
        parts.extend(str(item) for item in bullet.get("sub") or [])
    return len(re.findall(r"\b[\w'-]+\b", " ".join(parts)))


def get_slide(plan: dict[str, Any], slide_id: str) -> dict[str, Any] | None:
    index = plan_index_for_slide_id(slide_id)
    slides = list(normalize_slide_ids(plan).get("slides") or [])
    if index is None or index >= len(slides):
        return None
    slide = slides[index]
    return slide if isinstance(slide, dict) else None


def make_user_decisions(plan: dict[str, Any], case: dict[str, Any], target_slide_ids: list[str], active_slide_id: str) -> dict[str, Any]:
    ids = slide_ids(plan)
    editable = set(target_slide_ids)
    protected = list(case.get("protected_slide_ids") or [slide_id for slide_id in ids if slide_id not in editable])
    constraints = {
        "max_changed_slides": int(case.get("max_changed_slides", len(editable))),
        "allow_slide_splits": False,
        "allow_slide_merges": False,
        "preserve_section_order": True,
        "preserve_paper_fidelity": True,
        "do_not_introduce_external_assets": True,
    }
    constraints.update(case.get("constraints") or {})
    instruction = str(case.get("instruction") or "").format(
        slide_id=active_slide_id,
        active_slide_id=active_slide_id,
        target_slide_ids=", ".join(target_slide_ids),
    )
    edit_tasks = list(case.get("edit_tasks") or [])
    return {
        "user_requested_changes": [
            {
                "slide_id": active_slide_id,
                "active_slide_id": active_slide_id,
                "editable_slide_ids": target_slide_ids,
                "instruction": instruction,
                "edit_tasks": edit_tasks,
            }
        ],
        "edit_tasks": edit_tasks,
        "candidate_edit_decisions": list(case.get("candidate_edit_decisions") or []),
        "protected_slide_ids": protected,
        "protected_slide_overrides": list(case.get("protected_slide_overrides") or []),
        "constraints": constraints,
    }


def pptx_text(path: Path) -> str:
    if not path.exists():
        return ""
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def rendered_issue_count(render_dir: Path) -> tuple[int | None, int]:
    if not render_dir.exists() or not render_dir.is_dir():
        return None, 0
    images = sorted(list(render_dir.glob("*.jpg")) + list(render_dir.glob("*.png")))
    if not images:
        return None, 0
    edits, _summaries = deterministic_render_edits(images)
    return len(edits), len(images)


def add_check(checks: list[dict[str, Any]], name: str, passed: bool, details: Any = None) -> None:
    checks.append({"name": name, "passed": bool(passed), "details": details})


def validation_blocked_error(exc: Exception, repair_report: dict[str, Any] | None = None) -> bool:
    message = str(exc)
    if message.startswith("Refinement validation failed; refusing to regenerate PPTX."):
        return True
    validation = (repair_report or {}).get("validation") or {}
    return validation.get("accepted") is False


def load_repair_report(case_dir: Path, report_name: str = "post_render_refinement_report.json") -> dict[str, Any]:
    report_path = case_dir / report_name
    if not report_path.exists():
        return {}
    report = load_json(report_path, {})
    return report if isinstance(report, dict) else {}


def token_fields_from_report(repair_report: dict[str, Any]) -> dict[str, int]:
    token_usage = ((repair_report.get("cost_ledger") or {}).get("token_usage") or {})
    repair_input = int(token_usage.get("repair_input_tokens", 0) or 0)
    repair_output = int(token_usage.get("repair_output_tokens", 0) or 0)
    total_input = int(token_usage.get("total_added_input_tokens", repair_input) or 0)
    total_output = int(token_usage.get("total_added_output_tokens", repair_output) or 0)
    return {
        "repair_input_tokens": repair_input,
        "repair_output_tokens": repair_output,
        "repair_total_tokens": repair_input + repair_output,
        "total_added_input_tokens": total_input,
        "total_added_output_tokens": total_output,
        "total_added_tokens": total_input + total_output,
    }


def combine_token_fields(*token_sets: dict[str, Any]) -> dict[str, int]:
    combined = {
        "repair_input_tokens": 0,
        "repair_output_tokens": 0,
        "repair_total_tokens": 0,
        "total_added_input_tokens": 0,
        "total_added_output_tokens": 0,
        "total_added_tokens": 0,
    }
    for tokens in token_sets:
        for key in combined:
            combined[key] += int((tokens or {}).get(key, 0) or 0)
    return combined


def task_summary_from_report(repair_report: dict[str, Any]) -> dict[str, Any]:
    validation = repair_report.get("validation") or {}
    applied_unchanged = [
        item
        for item in validation.get("applied_task_unchanged_slides") or []
        if isinstance(item, dict)
    ]
    applied_unchanged_ids = [
        str(item.get("task_id"))
        for item in applied_unchanged
        if item.get("task_id")
    ]
    return {
        "required_task_count": validation.get("required_edit_task_count"),
        "edit_task_checklist_count": validation.get("edit_task_checklist_count"),
        "missing_edit_task_ids": [str(item) for item in validation.get("missing_edit_task_ids") or []],
        "skipped_edit_task_ids": [str(item) for item in validation.get("skipped_edit_task_ids") or []],
        "applied_but_unchanged_task_ids": applied_unchanged_ids,
        "applied_but_unchanged_count": len(applied_unchanged),
        "unaccounted_editable_slide_ids": [str(item) for item in validation.get("unaccounted_editable_slide_ids") or []],
        "protected_slide_violations": [str(item) for item in validation.get("protected_slide_violations") or []],
        "validation_accepted": validation.get("accepted"),
    }


def retry_decisions_payload(user_decisions: dict[str, Any], repair_report: dict[str, Any]) -> dict[str, Any]:
    retry_decisions = json.loads(json.dumps(user_decisions))
    validation = repair_report.get("validation") or {}
    retry_decisions["validation_retry_context"] = {
        "previous_repair_failed_validation": True,
        "applied_task_unchanged_slides": validation.get("applied_task_unchanged_slides") or [],
        "unaccounted_editable_slide_ids": validation.get("unaccounted_editable_slide_ids") or [],
        "protected_slide_violations": validation.get("protected_slide_violations") or [],
        "missing_edit_task_ids": validation.get("missing_edit_task_ids") or [],
        "skipped_edit_task_ids": validation.get("skipped_edit_task_ids") or [],
        "accepted": validation.get("accepted"),
        "instruction": (
            "Retry once by correcting only the validation failures. Preserve the same "
            "editable/protected slide constraints. Do not mark a task applied unless "
            "the revised slide differs from the original slide plan."
        ),
    }
    return retry_decisions


def contains_any(text: str, terms: list[str]) -> bool:
    lower = text.lower()
    return any(str(term).lower() in lower for term in terms)


def contains_all(text: str, terms: list[str]) -> bool:
    lower = text.lower()
    return all(str(term).lower() in lower for term in terms)


def evaluate_case(
    *,
    original_plan: dict[str, Any],
    revised_plan: dict[str, Any],
    repair_report: dict[str, Any],
    target_slide_id: str,
    target_slide_ids: list[str] | None = None,
    expected: dict[str, Any],
) -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    validation = repair_report.get("validation") or {}
    paths = repair_report.get("paths") or {}
    regenerated_pptx = Path(paths["regenerated_pptx"]) if paths.get("regenerated_pptx") else None
    refined_render_dir = Path(paths["refined_render_dir"]) if paths.get("refined_render_dir") else None
    change_summary = slide_change_summary(normalize_slide_ids(original_plan), normalize_slide_ids(revised_plan))
    changed_ids = set(change_summary.get("changed_slide_ids") or [])
    target_ids = list(target_slide_ids or [target_slide_id])
    target_id_set = set(target_ids)
    original_slide = get_slide(original_plan, target_slide_id) or {}
    revised_slide = get_slide(revised_plan, target_slide_id) or {}
    original_title = str(original_slide.get("subsection") or original_slide.get("title") or "")
    revised_title = str(revised_slide.get("subsection") or revised_slide.get("title") or "")
    original_visual_counts = [visual_counts(slide) for slide in normalize_slide_ids(original_plan).get("slides") or [] if isinstance(slide, dict)]
    revised_visual_counts = [visual_counts(slide) for slide in normalize_slide_ids(revised_plan).get("slides") or [] if isinstance(slide, dict)]
    deck_text = pptx_text(regenerated_pptx) if regenerated_pptx else ""
    rendered_issues, rendered_count = rendered_issue_count(refined_render_dir) if refined_render_dir else (None, 0)

    add_check(checks, "repair_validation_accepted", bool(validation.get("accepted")), validation)
    add_check(checks, "regenerated_pptx_exists", bool(regenerated_pptx and regenerated_pptx.exists()), str(regenerated_pptx) if regenerated_pptx else None)
    add_check(
        checks,
        "slide_count_preserved",
        change_summary.get("slide_count_before") == change_summary.get("slide_count_after"),
        change_summary,
    )
    add_check(checks, "no_refined_render_issues", rendered_issues in (0, None), {"rendered_issues": rendered_issues, "rendered_slide_count": rendered_count})

    if expected.get("target_slide_changed", True):
        add_check(checks, "target_slide_changed", target_slide_id in changed_ids, sorted(changed_ids))
    if expected.get("only_target_slide_changed", True):
        add_check(checks, "only_target_slide_changed", changed_ids <= target_id_set, sorted(changed_ids))
    if expected.get("all_target_slides_changed"):
        add_check(checks, "all_target_slides_changed", target_id_set <= changed_ids, {"expected": target_ids, "changed": sorted(changed_ids)})
    if expected.get("only_target_slides_changed"):
        add_check(checks, "only_target_slides_changed", changed_ids <= target_id_set, {"expected": target_ids, "changed": sorted(changed_ids)})
    if "changed_slide_ids_exact" in expected:
        exact = {str(item) for item in expected.get("changed_slide_ids_exact") or []}
        add_check(checks, "changed_slide_ids_exact", changed_ids == exact, {"expected": sorted(exact), "changed": sorted(changed_ids)})
    if expected.get("protected_slides_unchanged", True):
        protected_violations = validation.get("protected_slide_violations") or []
        add_check(checks, "protected_slides_unchanged", len(protected_violations) == 0, protected_violations)
    if expected.get("title_changed"):
        add_check(checks, "target_title_changed", revised_title != original_title, {"before": original_title, "after": revised_title})
    if expected.get("title_changed_slide_id"):
        title_slide_id = str(expected["title_changed_slide_id"])
        before_title_slide = get_slide(original_plan, title_slide_id) or {}
        after_title_slide = get_slide(revised_plan, title_slide_id) or {}
        before_title = str(before_title_slide.get("subsection") or before_title_slide.get("title") or "")
        after_title = str(after_title_slide.get("subsection") or after_title_slide.get("title") or "")
        add_check(checks, "specified_slide_title_changed", before_title != after_title, {"slide_id": title_slide_id, "before": before_title, "after": after_title})
        if expected.get("pptx_specified_title_visible") and after_title:
            add_check(checks, "pptx_contains_specified_revised_title", after_title in deck_text, {"slide_id": title_slide_id, "title": after_title})
    if expected.get("pptx_target_title_visible") and revised_title:
        add_check(checks, "pptx_contains_revised_title", revised_title in deck_text, revised_title)
    if "title_contains" in expected:
        add_check(checks, "target_title_contains_any", contains_any(revised_title, list(expected["title_contains"])), revised_title)
    if "title_contains_all" in expected:
        add_check(checks, "target_title_contains_all", contains_all(revised_title, list(expected["title_contains_all"])), revised_title)
    if "bullet_count" in expected:
        add_check(checks, "target_bullet_count", len(bullet_texts(revised_slide)) == int(expected["bullet_count"]), bullet_texts(revised_slide))
    if "bullet_count_by_slide" in expected:
        for slide_id, expected_count in dict(expected["bullet_count_by_slide"]).items():
            slide = get_slide(revised_plan, str(slide_id)) or {}
            add_check(checks, f"bullet_count_{slide_id}", len(bullet_texts(slide)) == int(expected_count), bullet_texts(slide))
    if "bullet_count_delta_by_slide" in expected:
        for slide_id, expected_delta in dict(expected["bullet_count_delta_by_slide"]).items():
            before_slide = get_slide(original_plan, str(slide_id)) or {}
            after_slide = get_slide(revised_plan, str(slide_id)) or {}
            before_count = len(bullet_texts(before_slide))
            after_count = len(bullet_texts(after_slide))
            add_check(
                checks,
                f"bullet_count_delta_{slide_id}",
                after_count - before_count == int(expected_delta),
                {"before": before_count, "after": after_count, "expected_delta": int(expected_delta)},
            )
    if "max_bullet_count" in expected:
        add_check(checks, "target_max_bullet_count", len(bullet_texts(revised_slide)) <= int(expected["max_bullet_count"]), bullet_texts(revised_slide))
    if expected.get("target_word_count_not_increased"):
        add_check(
            checks,
            "target_word_count_not_increased",
            slide_word_count(revised_slide) <= slide_word_count(original_slide),
            {"before": slide_word_count(original_slide), "after": slide_word_count(revised_slide)},
        )
    if "must_include_any" in expected:
        target_text = revised_title + "\n" + "\n".join(bullet_texts(revised_slide))
        add_check(checks, "target_contains_any_expected_term", contains_any(target_text, list(expected["must_include_any"])), target_text)
    if "must_include_all" in expected:
        target_text = revised_title + "\n" + "\n".join(bullet_texts(revised_slide))
        add_check(checks, "target_contains_all_expected_terms", contains_all(target_text, list(expected["must_include_all"])), target_text)
    if "forbidden_terms" in expected:
        target_text = revised_title + "\n" + "\n".join(bullet_texts(revised_slide))
        forbidden_found = [term for term in expected["forbidden_terms"] if str(term).lower() in target_text.lower()]
        add_check(checks, "target_forbidden_terms_absent", not forbidden_found, forbidden_found)
    if expected.get("no_external_asset_changes", True):
        add_check(checks, "visual_asset_counts_unchanged", original_visual_counts == revised_visual_counts, {"before": original_visual_counts, "after": revised_visual_counts})

    passed = sum(1 for check in checks if check["passed"])
    total = len(checks)
    return {
        "target_slide_id": target_slide_id,
        "target_slide_ids": target_ids,
        "score": round(passed / total, 4) if total else 0.0,
        "passed": passed == total,
        "passed_checks": passed,
        "total_checks": total,
        "failed_checks": [check for check in checks if not check["passed"]],
        "checks": checks,
        "changed_slide_ids": sorted(changed_ids),
        "target_before": {
            "title": original_title,
            "bullets": bullet_texts(original_slide),
            "word_count": slide_word_count(original_slide),
        },
        "target_after": {
            "title": revised_title,
            "bullets": bullet_texts(revised_slide),
            "word_count": slide_word_count(revised_slide),
        },
        "target_summaries": {
            slide_id: {
                "before": {
                    "title": str((get_slide(original_plan, slide_id) or {}).get("subsection") or (get_slide(original_plan, slide_id) or {}).get("title") or ""),
                    "bullets": bullet_texts(get_slide(original_plan, slide_id) or {}),
                    "word_count": slide_word_count(get_slide(original_plan, slide_id) or {}),
                },
                "after": {
                    "title": str((get_slide(revised_plan, slide_id) or {}).get("subsection") or (get_slide(revised_plan, slide_id) or {}).get("title") or ""),
                    "bullets": bullet_texts(get_slide(revised_plan, slide_id) or {}),
                    "word_count": slide_word_count(get_slide(revised_plan, slide_id) or {}),
                },
            }
            for slide_id in target_ids
        },
    }


def inspect_once(
    *,
    deck_dir: Path,
    pptx_path: Path,
    slide_plan_path: Path,
    output_dir: Path,
    author_profile_path: Path | None,
    args: argparse.Namespace,
) -> None:
    if args.resume and (output_dir / "post_render_deck_summary.json").exists() and (output_dir / "post_render_candidate_edits.json").exists():
        return
    inspect_deck(
        pptx_path=pptx_path,
        slide_plan_path=slide_plan_path,
        output_dir=output_dir,
        author_profile_path=author_profile_path,
        render_dpi=args.render_dpi,
        force_render=args.force_render,
        timeout=args.timeout,
    )


def run_case(
    *,
    batch_dir: Path,
    deck_spec: dict[str, Any],
    case_template: dict[str, Any],
    deck_index: int,
    case_index: int,
    args: argparse.Namespace,
) -> dict[str, Any]:
    run_dir = Path(deck_spec["run_dir"]).resolve()
    pptx_path = Path(deck_spec.get("pptx_path") or infer_local_pptx(run_dir)).resolve()
    slide_plan_path = Path(deck_spec.get("slide_plan_path") or infer_local_slide_plan(run_dir)).resolve()
    author_profile_path = deck_spec.get("author_profile_path")
    if author_profile_path:
        author_profile_path = Path(author_profile_path).resolve()

    original_plan = normalize_slide_ids(load_json(slide_plan_path, {}))
    model_name_t, model_name_v = infer_model_names_from_plan(slide_plan_path)
    target_slide_ids, active_slide_id = resolve_case_slide_ids(original_plan, case_template)
    target_slide_id = active_slide_id
    case_id = safe_name(str(case_template.get("case_id") or f"case_{case_index:02d}"))
    deck_id = safe_name(run_dir.name)
    case_dir = batch_dir / deck_id / f"{case_index:02d}_{case_id}"
    case_result_path = case_dir / "case_result.json"
    if args.resume and case_result_path.exists():
        return load_json(case_result_path, {})

    user_decisions = make_user_decisions(original_plan, case_template, target_slide_ids, active_slide_id)
    write_json(case_dir / "user_decisions.json", user_decisions)

    if args.dry_run:
        result = {
            "status": "dry_run",
            "deck": run_dir.name,
            "case_id": case_id,
            "prompt_family": case_template.get("prompt_family"),
            "target_slide_id": target_slide_id,
            "target_slide_ids": target_slide_ids,
            "instruction": user_decisions["user_requested_changes"][0]["instruction"],
            "case_dir": str(case_dir),
        }
        write_json(case_result_path, result)
        return result

    inspection_dir = batch_dir / deck_id / "_inspection"
    inspect_once(
        deck_dir=run_dir,
        pptx_path=pptx_path,
        slide_plan_path=slide_plan_path,
        output_dir=inspection_dir,
        author_profile_path=author_profile_path,
        args=args,
    )

    started = time.time()
    output_variant_suffix = f"_refined_{case_id}"
    base_repair_args = {
        "slide_plan_path": slide_plan_path,
        "candidate_edits_path": inspection_dir / "post_render_candidate_edits.json",
        "rendered_deck_summary_path": inspection_dir / "post_render_deck_summary.json",
        "author_profile_path": author_profile_path,
        "repair_mode": "llm",
        "model": args.model,
        "repair_prompt_path": args.repair_prompt_path,
        "timeout": args.timeout,
        "regenerate": True,
        "validate_render": True,
        "render_dpi": args.render_dpi,
        "paper_name": run_dir.name,
        "model_name_t": model_name_t,
        "model_name_v": model_name_v,
        "output_dir": str(args.pipeline_output_root or infer_output_root_from_run_dir(run_dir)),
        "asset_paper_name": infer_asset_paper_name(run_dir, model_name_t, model_name_v),
        "formula_mode": args.formula_mode,
        "template": args.template,
    }
    repair_args = SimpleNamespace(
        **base_repair_args,
        user_decisions_path=case_dir / "user_decisions.json",
        output_dir_path=case_dir,
        revised_plan_path=case_dir / "post_render_refined_slide_plan.json",
        output_variant_suffix=output_variant_suffix,
    )

    def blocked_result(exc: Exception, report: dict[str, Any], *, retry_attempted: bool = False, retry_report: dict[str, Any] | None = None) -> dict[str, Any]:
        tokens = token_fields_from_report(report)
        if retry_report:
            tokens = combine_token_fields(tokens, token_fields_from_report(retry_report))
        final_report = retry_report or report
        result = {
            "status": "validation_blocked" if validation_blocked_error(exc, final_report) else "error",
            "deck_index": deck_index,
            "case_index": case_index,
            "deck": run_dir.name,
            "case_id": case_id,
            "prompt_family": case_template.get("prompt_family"),
            "case_dir": str(case_dir),
            "target_slide_id": target_slide_id,
            "target_slide_ids": target_slide_ids,
            "instruction": user_decisions["user_requested_changes"][0]["instruction"],
            "runtime_seconds": round(time.time() - started, 3),
            "paths": final_report.get("paths") or {},
            "repair_validation": final_report.get("validation") or {},
            "cost_ledger": final_report.get("cost_ledger") or {},
            "tokens": tokens,
            "retry_validation_blocks": bool(args.retry_validation_blocks),
            "retry_attempted": retry_attempted,
            "retry_count": 1 if retry_attempted else 0,
            "error": str(exc),
        }
        result.update(task_summary_from_report(final_report))
        write_json(case_result_path, result)
        return result

    try:
        repair_report = repair_deck(repair_args)
    except Exception as exc:
        repair_report = load_repair_report(case_dir)
        if validation_blocked_error(exc, repair_report):
            if args.retry_validation_blocks:
                retry_decisions_path = case_dir / "user_decisions.retry_1.json"
                retry_case_dir = case_dir / "retry_1"
                write_json(retry_decisions_path, retry_decisions_payload(user_decisions, repair_report))
                retry_args = SimpleNamespace(
                    **base_repair_args,
                    user_decisions_path=retry_decisions_path,
                    output_dir_path=retry_case_dir,
                    revised_plan_path=retry_case_dir / "post_render_refined_slide_plan.retry_1.json",
                    output_variant_suffix=f"{output_variant_suffix}_retry_1",
                )
                try:
                    retry_report = repair_deck(retry_args)
                    retry_report_path = retry_case_dir / "post_render_refinement_report.json"
                    retry_copy_path = case_dir / "post_render_refinement_report.retry_1.json"
                    if retry_report_path.exists():
                        write_json(retry_copy_path, retry_report)
                    repair_report = retry_report
                    revised_plan_path = retry_case_dir / "post_render_refined_slide_plan.retry_1.json"
                    retry_tokens = token_fields_from_report(retry_report)
                    original_tokens = token_fields_from_report(load_repair_report(case_dir))
                    combined_tokens = combine_token_fields(original_tokens, retry_tokens)
                    revised_plan = normalize_slide_ids(load_json(revised_plan_path, {}))
                    metrics = evaluate_case(
                        original_plan=original_plan,
                        revised_plan=revised_plan,
                        repair_report=repair_report,
                        target_slide_id=target_slide_id,
                        target_slide_ids=target_slide_ids,
                        expected=dict(case_template.get("expected") or {}),
                    )
                    result = {
                        "status": "ok" if metrics["passed"] else "failed_checks",
                        "deck_index": deck_index,
                        "case_index": case_index,
                        "deck": run_dir.name,
                        "case_id": case_id,
                        "prompt_family": case_template.get("prompt_family"),
                        "case_dir": str(case_dir),
                        "target_slide_id": target_slide_id,
                        "target_slide_ids": target_slide_ids,
                        "instruction": user_decisions["user_requested_changes"][0]["instruction"],
                        "runtime_seconds": round(time.time() - started, 3),
                        "paths": repair_report.get("paths") or {},
                        "metrics": metrics,
                        "repair_validation": repair_report.get("validation") or {},
                        "cost_ledger": repair_report.get("cost_ledger") or {},
                        "tokens": combined_tokens,
                        "retry_validation_blocks": True,
                        "retry_attempted": True,
                        "retry_count": 1,
                    }
                    result.update(task_summary_from_report(repair_report))
                    write_json(case_result_path, result)
                    return result
                except Exception as retry_exc:
                    retry_report = load_repair_report(retry_case_dir)
                    return blocked_result(retry_exc, repair_report, retry_attempted=True, retry_report=retry_report)
            return blocked_result(exc, repair_report)
        raise
    revised_plan = normalize_slide_ids(load_json(case_dir / "post_render_refined_slide_plan.json", {}))
    metrics = evaluate_case(
        original_plan=original_plan,
        revised_plan=revised_plan,
        repair_report=repair_report,
        target_slide_id=target_slide_id,
        target_slide_ids=target_slide_ids,
        expected=dict(case_template.get("expected") or {}),
    )
    result = {
        "status": "ok" if metrics["passed"] else "failed_checks",
        "deck_index": deck_index,
        "case_index": case_index,
        "deck": run_dir.name,
        "case_id": case_id,
        "prompt_family": case_template.get("prompt_family"),
        "case_dir": str(case_dir),
        "target_slide_id": target_slide_id,
        "target_slide_ids": target_slide_ids,
        "instruction": user_decisions["user_requested_changes"][0]["instruction"],
        "runtime_seconds": round(time.time() - started, 3),
        "paths": repair_report.get("paths") or {},
        "metrics": metrics,
        "repair_validation": repair_report.get("validation") or {},
        "cost_ledger": repair_report.get("cost_ledger") or {},
        "tokens": token_fields_from_report(repair_report),
        "retry_validation_blocks": bool(args.retry_validation_blocks),
        "retry_attempted": False,
        "retry_count": 0,
    }
    result.update(task_summary_from_report(repair_report))
    write_json(case_result_path, result)
    return result


def summarize_results(batch_dir: Path, results: list[dict[str, Any]]) -> dict[str, Any]:
    completed = [item for item in results if item.get("status") in {"ok", "failed_checks"}]
    token_counted = [item for item in results if item.get("status") != "dry_run"]
    failures = [
        item
        for item in results
        if item.get("status") in {"failed_checks", "validation_blocked", "error"}
        or not item.get("metrics", {}).get("passed", item.get("status") in {"dry_run", "ok"})
    ]
    status_counts = {
        status: sum(1 for item in results if item.get("status") == status)
        for status in ["ok", "failed_checks", "validation_blocked", "error", "dry_run"]
    }
    total_input = sum(int((item.get("tokens") or {}).get("total_added_input_tokens", 0) or 0) for item in token_counted)
    total_output = sum(int((item.get("tokens") or {}).get("total_added_output_tokens", 0) or 0) for item in token_counted)
    scores = [float((item.get("metrics") or {}).get("score", 0.0) or 0.0) for item in completed]
    prompt_family_summary: dict[str, dict[str, Any]] = {}
    for item in results:
        family = str(item.get("prompt_family") or item.get("case_id") or "unknown")
        bucket = prompt_family_summary.setdefault(
            family,
            {
                "case_count": 0,
                "completed_count": 0,
                "pass_count": 0,
                "failure_count": 0,
                "validation_blocked_count": 0,
                "error_count": 0,
                "total_added_input_tokens": 0,
                "total_added_output_tokens": 0,
                "scores": [],
            },
        )
        status = item.get("status")
        metrics = item.get("metrics") or {}
        tokens = item.get("tokens") or {}
        bucket["case_count"] += 1
        if status in {"ok", "failed_checks"}:
            bucket["completed_count"] += 1
            bucket["scores"].append(float(metrics.get("score", 0.0) or 0.0))
        if metrics.get("passed"):
            bucket["pass_count"] += 1
        if status in {"failed_checks", "validation_blocked", "error"}:
            bucket["failure_count"] += 1
        if status == "validation_blocked":
            bucket["validation_blocked_count"] += 1
        if status == "error":
            bucket["error_count"] += 1
        if status != "dry_run":
            bucket["total_added_input_tokens"] += int(tokens.get("total_added_input_tokens", 0) or 0)
            bucket["total_added_output_tokens"] += int(tokens.get("total_added_output_tokens", 0) or 0)
    for bucket in prompt_family_summary.values():
        bucket_scores = list(bucket.pop("scores"))
        bucket["mean_score"] = round(sum(bucket_scores) / len(bucket_scores), 4) if bucket_scores else None
        bucket["total_added_tokens"] = bucket["total_added_input_tokens"] + bucket["total_added_output_tokens"]
    summary = {
        "batch_dir": str(batch_dir),
        "case_count": len(results),
        "completed_count": len(completed),
        "validation_blocked_count": status_counts["validation_blocked"],
        "error_count": status_counts["error"],
        "dry_run_count": status_counts["dry_run"],
        "failure_count": len(failures),
        "pass_count": sum(1 for item in completed if (item.get("metrics") or {}).get("passed")),
        "mean_score": round(sum(scores) / len(scores), 4) if scores else None,
        "status_counts": status_counts,
        "total_added_input_tokens": total_input,
        "total_added_output_tokens": total_output,
        "total_added_tokens": total_input + total_output,
        "prompt_family_summary": prompt_family_summary,
        "retry_attempt_count": sum(int(item.get("retry_count", 0) or 0) for item in results),
        "applied_but_unchanged_count": sum(int(item.get("applied_but_unchanged_count", 0) or 0) for item in results),
        "failed_cases": [
            {
                "deck": item.get("deck"),
                "case_id": item.get("case_id"),
                "prompt_family": item.get("prompt_family"),
                "status": item.get("status"),
                "failed_checks": [check.get("name") for check in (item.get("metrics") or {}).get("failed_checks", [])],
                "error": item.get("error"),
                "required_task_count": item.get("required_task_count"),
                "edit_task_checklist_count": item.get("edit_task_checklist_count"),
                "missing_edit_task_ids": item.get("missing_edit_task_ids"),
                "skipped_edit_task_ids": item.get("skipped_edit_task_ids"),
                "applied_but_unchanged_task_ids": item.get("applied_but_unchanged_task_ids"),
                "applied_but_unchanged_count": item.get("applied_but_unchanged_count"),
                "unaccounted_editable_slide_ids": item.get("unaccounted_editable_slide_ids"),
                "protected_slide_violations": item.get("protected_slide_violations"),
                "validation_accepted": item.get("validation_accepted"),
                "retry_attempted": item.get("retry_attempted"),
            }
            for item in failures
        ],
    }
    write_json(batch_dir / "summary.json", summary)
    return summary


def write_csv_summary(batch_dir: Path, results: list[dict[str, Any]]) -> None:
    rows = []
    for item in results:
        metrics = item.get("metrics") or {}
        tokens = item.get("tokens") or {}
        rows.append(
            {
                "deck": item.get("deck"),
                "case_id": item.get("case_id"),
                "prompt_family": item.get("prompt_family"),
                "target_slide_id": item.get("target_slide_id"),
                "target_slide_ids": ";".join(item.get("target_slide_ids") or ([item.get("target_slide_id")] if item.get("target_slide_id") else [])),
                "status": item.get("status"),
                "passed": metrics.get("passed"),
                "score": metrics.get("score"),
                "failed_checks": ";".join(check.get("name", "") for check in metrics.get("failed_checks", [])),
                "changed_slide_ids": ";".join(metrics.get("changed_slide_ids", [])),
                "repair_input_tokens": tokens.get("repair_input_tokens"),
                "repair_output_tokens": tokens.get("repair_output_tokens"),
                "repair_total_tokens": tokens.get("repair_total_tokens"),
                "total_added_input_tokens": tokens.get("total_added_input_tokens"),
                "total_added_output_tokens": tokens.get("total_added_output_tokens"),
                "total_added_tokens": tokens.get("total_added_tokens"),
                "required_task_count": item.get("required_task_count"),
                "edit_task_checklist_count": item.get("edit_task_checklist_count"),
                "missing_edit_task_ids": ";".join(item.get("missing_edit_task_ids") or []),
                "skipped_edit_task_ids": ";".join(item.get("skipped_edit_task_ids") or []),
                "applied_but_unchanged_task_ids": ";".join(item.get("applied_but_unchanged_task_ids") or []),
                "applied_but_unchanged_count": item.get("applied_but_unchanged_count"),
                "unaccounted_editable_slide_ids": ";".join(item.get("unaccounted_editable_slide_ids") or []),
                "protected_slide_violations": ";".join(item.get("protected_slide_violations") or []),
                "validation_accepted": item.get("validation_accepted"),
                "retry_attempted": item.get("retry_attempted"),
                "retry_count": item.get("retry_count"),
                "runtime_seconds": item.get("runtime_seconds"),
                "case_dir": item.get("case_dir"),
                "regenerated_pptx": (item.get("paths") or {}).get("regenerated_pptx"),
            }
        )
    output_path = batch_dir / "summary.csv"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()) if rows else ["deck"])
        writer.writeheader()
        writer.writerows(rows)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Batch-run post-render user-prompt refinement tests.")
    parser.add_argument("--manifest", type=Path, default=None, help="Optional JSON manifest with decks/cases.")
    parser.add_argument("--run-dir", type=Path, action="append", default=[], help="Generated contents/<paper> directory. Can be repeated.")
    parser.add_argument("--contents-dir", type=Path, default=REPO_ROOT / "contents")
    parser.add_argument("--limit-decks", type=int, default=5)
    parser.add_argument("--require-deck-count", type=int, default=0, help="Fail unless at least this many deck run directories are selected.")
    parser.add_argument("--allow-fewer-decks", action="store_true", help="Allow auto-discovery to run fewer than --limit-decks decks.")
    parser.add_argument("--case-limit", type=int, default=0)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    parser.add_argument("--batch-id", default=None)
    parser.add_argument("--author-profile-path", type=Path, default=None)
    parser.add_argument("--model", default="gpt-4o-mini")
    parser.add_argument("--repair-prompt-path", type=Path, default=DEFAULT_REPAIR_PROMPT)
    parser.add_argument("--render-dpi", type=int, default=120)
    parser.add_argument("--force-render", action="store_true")
    parser.add_argument("--timeout", type=float, default=180.0)
    parser.add_argument("--pipeline-output-root", type=Path, default=None)
    parser.add_argument("--formula-mode", type=int, default=1)
    parser.add_argument("--template", type=int, default=3)
    parser.add_argument("--resume", action="store_true", help="Skip case folders that already have case_result.json.")
    parser.add_argument("--dry-run", action="store_true", help="Write planned cases without calling the repair model.")
    parser.add_argument("--retry-validation-blocks", action="store_true", help="Retry validation-blocked repairs once with failure details in the prompt context.")
    parser.add_argument("--stop-on-error", action="store_true")
    return parser


def main() -> None:
    args = build_parser().parse_args()
    load_dotenv(REPO_ROOT / ".env")
    if not args.dry_run and not os.environ.get("OPENAI_API_KEY"):
        raise SystemExit("OPENAI_API_KEY is required for non-dry-run repair batches. Set it in the shell or repo .env, or use --dry-run.")

    batch_id = safe_name(args.batch_id or time.strftime("batch_%Y%m%d_%H%M%S"))
    batch_dir = (args.output_root / batch_id).resolve()
    batch_dir.mkdir(parents=True, exist_ok=True)

    deck_specs = normalize_deck_specs(args)
    if not deck_specs:
        raise SystemExit("No generated deck run directories found. Pass --run-dir or check --contents-dir.")

    write_json(
        batch_dir / "resolved_batch_manifest.json",
        {
            "batch_id": batch_id,
            "model": args.model,
            "llm_visual_review": False,
            "decks": [
                {
                    "run_dir": str(Path(spec["run_dir"]).resolve()),
                    "case_count": len(spec.get("case_templates") or []),
                    "author_profile_path": str(spec.get("author_profile_path")) if spec.get("author_profile_path") else None,
                }
                for spec in deck_specs
            ],
        },
    )

    results: list[dict[str, Any]] = []
    planned_cases = 0
    for deck_index, deck_spec in enumerate(deck_specs, start=1):
        case_templates = list(deck_spec.get("case_templates") or DEFAULT_CASE_TEMPLATES)
        for case_index, case_template in enumerate(case_templates, start=1):
            planned_cases += 1
            if args.case_limit > 0 and len(results) >= args.case_limit:
                break
            try:
                print(f"[batch] deck {deck_index}/{len(deck_specs)} case {case_index}/{len(case_templates)}: {Path(deck_spec['run_dir']).name} / {case_template.get('case_id')}", flush=True)
                result = run_case(
                    batch_dir=batch_dir,
                    deck_spec=deck_spec,
                    case_template=case_template,
                    deck_index=deck_index,
                    case_index=case_index,
                    args=args,
                )
            except Exception as exc:
                deck_name = Path(deck_spec["run_dir"]).name
                case_id = safe_name(str(case_template.get("case_id") or f"case_{case_index:02d}"))
                case_dir = batch_dir / safe_name(deck_name) / f"{case_index:02d}_{case_id}"
                repair_report = load_repair_report(case_dir)
                status = "validation_blocked" if validation_blocked_error(exc, repair_report) else "error"
                result = {
                    "status": status,
                    "deck_index": deck_index,
                    "case_index": case_index,
                    "deck": deck_name,
                    "case_id": case_id,
                    "prompt_family": case_template.get("prompt_family"),
                    "case_dir": str(case_dir),
                    "paths": repair_report.get("paths") or {},
                    "repair_validation": repair_report.get("validation") or {},
                    "cost_ledger": repair_report.get("cost_ledger") or {},
                    "tokens": token_fields_from_report(repair_report),
                    "error": str(exc),
                }
                result.update(task_summary_from_report(repair_report))
                if args.stop_on_error:
                    results.append(result)
                    write_json(batch_dir / "partial_results.json", results)
                    raise
            results.append(result)
        if args.case_limit > 0 and len(results) >= args.case_limit:
            break

    write_json(batch_dir / "case_results.json", results)
    write_csv_summary(batch_dir, results)
    summary = summarize_results(batch_dir, results)
    print(json.dumps(summary, indent=2), flush=True)


if __name__ == "__main__":
    main()
