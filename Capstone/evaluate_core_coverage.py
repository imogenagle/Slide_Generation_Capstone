#!/usr/bin/env python3
"""Evaluate generated decks against original slide decks using topic-set IoU."""

from __future__ import annotations

import argparse
import base64
import difflib
import json
import re
import sys
from pathlib import Path
from typing import Any

from openai import BadRequestError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from dotenv import load_dotenv
except ModuleNotFoundError:
    def load_dotenv(*_args: Any, **_kwargs: Any) -> bool:
        return False

REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from slidegen_openai_utils import build_openai_client, resolve_direct_model_name

DEFAULT_OUTPUT_DIR = REPO_ROOT / "Capstone" / "evaluations"
SLIDE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}
MAX_IMAGES_PER_REQUEST = 50
MATCH_STOPWORDS = {
    "a",
    "an",
    "and",
    "as",
    "at",
    "based",
    "by",
    "for",
    "from",
    "in",
    "into",
    "of",
    "on",
    "or",
    "over",
    "the",
    "to",
    "using",
    "via",
    "with",
}
TOKEN_SYNONYMS = {
    "activities": "activity",
    "forecast": "forecasting",
    "forecasts": "forecasting",
    "generalizes": "transfer",
    "generalize": "transfer",
    "generalized": "transfer",
    "goals": "goal",
    "metrics": "metric",
    "models": "model",
    "multi": "multi",
    "multi goal": "multigoal",
    "noisy": "noise",
    "novel": "new",
    "observations": "observation",
    "physicalfeature": "physical",
    "probabilities": "probabilistic",
    "probability": "probabilistic",
    "results": "result",
    "scenes": "scene",
    "semantics": "semantic",
    "trajectories": "trajectory",
}


def encode_image_data_uri(image_path: Path) -> str:
    mime_type = "image/jpeg"
    if image_path.suffix.lower() == ".png":
        mime_type = "image/png"
    elif image_path.suffix.lower() == ".webp":
        mime_type = "image/webp"
    data = base64.b64encode(image_path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{data}"


def numeric_slide_sort_key(path: Path) -> tuple[int, str]:
    try:
        return (int(path.stem), path.name)
    except ValueError:
        return (10**9, path.name)


def sample_slide_paths(slide_paths: list[Path], max_images: int) -> list[Path]:
    if max_images <= 0 or len(slide_paths) <= max_images:
        return slide_paths

    if max_images <= 4:
        indices = [0, len(slide_paths) // 3, (2 * len(slide_paths)) // 3, len(slide_paths) - 1]
        return [slide_paths[i] for i in sorted(set(max(0, min(i, len(slide_paths) - 1)) for i in indices))]

    selected = {0, 1, len(slide_paths) - 2, len(slide_paths) - 1}
    remaining_slots = max_images - len(selected)
    if remaining_slots > 0:
        span = len(slide_paths) - 1
        for idx in range(remaining_slots):
            position = round(((idx + 1) * span) / (remaining_slots + 1))
            selected.add(max(0, min(position, len(slide_paths) - 1)))
    return [slide_paths[i] for i in sorted(selected)[:max_images]]


def extract_text_from_shape(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text and paragraph.text.strip()
        ).strip()
        if text:
            return text
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [
                cell.text.strip()
                for cell in row.cells
                if cell.text and cell.text.strip()
            ]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            return "\n".join(rows)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        parts = [extract_text_from_shape(subshape) for subshape in shape.shapes]
        return "\n".join(part for part in parts if part).strip()
    return ""


def extract_pptx_text(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    slide_blocks = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []
        notes_text = ""
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                texts.append(text)
        if hasattr(slide, "notes_slide") and slide.notes_slide is not None:
            for shape in slide.notes_slide.shapes:
                if getattr(shape, "is_placeholder", False) and getattr(shape.placeholder_format, "idx", None) == 1:
                    continue
                note_piece = extract_text_from_shape(shape)
                if note_piece:
                    notes_text = f"Speaker notes\n{note_piece}"
        if texts:
            block_lines = [f"Slide {slide_index}"] + [f"- {line}" for line in texts]
            if notes_text:
                block_lines.append(f"- {notes_text}")
            slide_blocks.append("\n".join(block_lines))
    return "\n\n".join(slide_blocks)


def extract_json_object(raw_text: str) -> dict[str, Any]:
    text = (raw_text or "").strip()
    fenced = re.search(r"```(?:json)?\s*(\{.*\})\s*```", text, flags=re.DOTALL)
    if fenced:
        text = fenced.group(1).strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        first = text.find("{")
        last = text.rfind("}")
        if first != -1 and last != -1 and last > first:
            return json.loads(text[first:last + 1])
        raise


def extract_message_text(message: Any) -> str:
    content = getattr(message, "content", None)
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        text_parts: list[str] = []
        for part in content:
            if isinstance(part, dict):
                text_value = part.get("text")
                if isinstance(text_value, str) and text_value.strip():
                    text_parts.append(text_value)
            else:
                text_value = getattr(part, "text", None)
                if isinstance(text_value, str) and text_value.strip():
                    text_parts.append(text_value)
        return "\n".join(text_parts).strip()
    return ""


def normalize_topic_label(raw_value: Any) -> str:
    text = str(raw_value or "").strip().lower()
    text = re.sub(r"[_/\-]+", " ", text)
    text = re.sub(r"[^a-z0-9+ -]", "", text)
    text = re.sub(r"\s+", " ", text).strip(" -")
    return text


def normalize_match_token(token: str) -> str:
    token = token.strip().lower()
    if not token:
        return ""
    if token.endswith("ies") and len(token) > 4:
        token = token[:-3] + "y"
    elif token.endswith("ing") and len(token) > 5:
        token = token[:-3]
    elif token.endswith("ed") and len(token) > 4:
        token = token[:-2]
    elif token.endswith("es") and len(token) > 4:
        token = token[:-2]
    elif token.endswith("s") and len(token) > 3:
        token = token[:-1]
    token = TOKEN_SYNONYMS.get(token, token)
    return token


def topic_match_tokens(topic: str) -> set[str]:
    normalized = normalize_topic_label(topic)
    tokens: set[str] = set()
    for token in normalized.split():
        if token in MATCH_STOPWORDS:
            continue
        token = normalize_match_token(token)
        if not token or token in MATCH_STOPWORDS:
            continue
        tokens.add(token)
    return tokens


def topic_similarity(topic_a: str, topic_b: str) -> float:
    normalized_a = normalize_topic_label(topic_a)
    normalized_b = normalize_topic_label(topic_b)
    if not normalized_a or not normalized_b:
        return 0.0
    if normalized_a == normalized_b:
        return 1.0

    tokens_a = topic_match_tokens(normalized_a)
    tokens_b = topic_match_tokens(normalized_b)
    if not tokens_a or not tokens_b:
        return 0.0

    intersection = tokens_a & tokens_b
    union = tokens_a | tokens_b
    jaccard = len(intersection) / len(union) if union else 0.0
    containment = len(intersection) / min(len(tokens_a), len(tokens_b))
    char_ratio = difflib.SequenceMatcher(None, normalized_a, normalized_b).ratio()

    bonus = 0.0
    if "virat" in intersection:
        bonus += 0.08
    if {"inverse", "optimal", "control"} <= intersection:
        bonus += 0.08
    if {"knowledge", "transfer"} <= intersection:
        bonus += 0.08
    if {"destination", "forecast"} <= intersection or {"destination", "forecasting"} <= intersection:
        bonus += 0.08
    if {"maximum", "entropy"} <= intersection:
        bonus += 0.08

    return min(1.0, max(jaccard, containment * 0.92, char_ratio * 0.75) + bonus)


def topics_are_semantic_match(topic_a: str, topic_b: str) -> tuple[bool, float]:
    score = topic_similarity(topic_a, topic_b)
    tokens_a = topic_match_tokens(topic_a)
    tokens_b = topic_match_tokens(topic_b)
    overlap = len(tokens_a & tokens_b)
    strong = score >= 0.72
    moderate = score >= 0.58 and overlap >= 3
    high_overlap = overlap >= 4 and score >= 0.5
    return strong or moderate or high_overlap, round(score, 3)


def semantic_topic_matches(reference_topics: list[str], generated_topics: list[str]) -> list[dict[str, Any]]:
    candidates: list[tuple[float, str, str]] = []
    for reference_topic in reference_topics:
        for generated_topic in generated_topics:
            is_match, score = topics_are_semantic_match(reference_topic, generated_topic)
            if is_match:
                candidates.append((score, reference_topic, generated_topic))

    matches: list[dict[str, Any]] = []
    used_reference: set[str] = set()
    used_generated: set[str] = set()
    for score, reference_topic, generated_topic in sorted(candidates, key=lambda item: (-item[0], item[1], item[2])):
        if reference_topic in used_reference or generated_topic in used_generated:
            continue
        matches.append(
            {
                "reference_topic": reference_topic,
                "generated_topic": generated_topic,
                "similarity": round(score, 3),
            }
        )
        used_reference.add(reference_topic)
        used_generated.add(generated_topic)
    return matches


def sanitize_topic_list(raw_topics: Any) -> list[str]:
    clean_topics: list[str] = []
    seen: set[str] = set()
    for item in raw_topics or []:
        if isinstance(item, dict):
            candidate = item.get("canonical_topic") or item.get("topic") or item.get("label") or ""
        else:
            candidate = item
        normalized = normalize_topic_label(candidate)
        if normalized and normalized not in seen:
            clean_topics.append(normalized)
            seen.add(normalized)
    return clean_topics


def sanitize_result(result: dict[str, Any], paper_id: str, title: str) -> dict[str, Any]:
    reference_topics = sanitize_topic_list(result.get("reference_topics"))
    generated_topics = sanitize_topic_list(result.get("generated_topics"))
    semantic_matches = semantic_topic_matches(reference_topics, generated_topics)
    matched_reference_topics = {item["reference_topic"] for item in semantic_matches}
    matched_generated_topics = {item["generated_topic"] for item in semantic_matches}
    matched_topics = sorted(item["reference_topic"] for item in semantic_matches)
    reference_only_topics = sorted(topic for topic in reference_topics if topic not in matched_reference_topics)
    generated_only_topics = sorted(topic for topic in generated_topics if topic not in matched_generated_topics)
    intersection_count = len(semantic_matches)
    union_count = len(reference_topics) + len(generated_topics) - intersection_count
    diff = result.get("difference_summary", {}) or {}
    return {
        "paper_id": paper_id,
        "title": title,
        "reference_topics": reference_topics,
        "generated_topics": generated_topics,
        "matched_topics": matched_topics,
        "matched_topic_pairs": semantic_matches,
        "reference_only_topics": reference_only_topics,
        "generated_only_topics": generated_only_topics,
        "reference_topic_count": len(reference_topics),
        "generated_topic_count": len(generated_topics),
        "topic_intersection_count": intersection_count,
        "topic_union_count": union_count,
        "topic_iou": round((intersection_count / union_count), 3) if union_count > 0 else None,
        "coverage_ratio": round((intersection_count / union_count), 3) if union_count > 0 else None,
        "missing_from_generated": reference_only_topics,
        "difference_summary": {
            "missing_content": [str(item) for item in diff.get("missing_content", [])],
            "compressed_content": [str(item) for item in diff.get("compressed_content", [])],
            "overemphasized_or_extra_content": [
                str(item) for item in diff.get("overemphasized_or_extra_content", [])
            ],
            "structure_changes": [str(item) for item in diff.get("structure_changes", [])],
            "overall_summary": str(diff.get("overall_summary", "")).strip(),
        },
        "notes": str(result.get("notes", "")).strip(),
    }


def create_chat_completion(client: Any, *, resolved_model_name: str, messages: list[dict[str, Any]]) -> dict[str, Any]:
    request_kwargs: dict[str, Any] = {
        "model": resolved_model_name,
        "messages": messages,
        "temperature": 0.1,
    }
    if "gpt-5" in resolved_model_name.lower():
        request_kwargs["max_completion_tokens"] = 1800
    else:
        request_kwargs["max_tokens"] = 1800
    response = client.chat.completions.create(**request_kwargs)
    raw_text = extract_message_text(response.choices[0].message)
    return extract_json_object(raw_text)


def is_content_policy_violation(exc: Exception) -> bool:
    if not isinstance(exc, BadRequestError):
        return False
    body = getattr(exc, "body", None)
    if isinstance(body, dict):
        error = body.get("error") or {}
        code = str(error.get("code") or "").strip().lower()
        message = str(error.get("message") or "").strip().lower()
        if code == "content_policy_violation":
            return True
        if "content safety" in message:
            return True
    return "content_policy_violation" in str(exc).lower()


def normalize_topics_with_llm(
    *,
    client: Any,
    resolved_model_name: str,
    paper_id: str,
    title: str,
    reference_topics: list[str],
    generated_topics: list[str],
) -> dict[str, Any]:
    system_prompt = (
        "You are an evaluation agent for scientific slide decks.\n"
        "Your job is to normalize two topic lists into a shared canonical vocabulary.\n"
        "Map semantically equivalent topics across the two lists to the same short canonical label.\n"
        "Preserve real distinctions when topics are materially different.\n"
        "Use concise labels of roughly 2 to 8 words.\n"
        "Return only valid JSON."
    )
    user_prompt = f"""
Paper ID: {paper_id}
Paper Title: {title}

Reference topics:
{json.dumps(reference_topics, ensure_ascii=False, indent=2)}

Generated topics:
{json.dumps(generated_topics, ensure_ascii=False, indent=2)}

Task:
1. For every reference topic, provide a concise canonical topic label.
2. For every generated topic, provide a concise canonical topic label.
3. Use exactly the same canonical label whenever the two topics express the same main scientific idea.
4. Keep labels broad enough to merge paraphrases, but specific enough to avoid collapsing distinct contributions.

Return exactly this JSON schema:
{{
  "reference_topics": [
    {{
      "raw_topic": "original reference topic",
      "canonical_topic": "shared canonical label"
    }}
  ],
  "generated_topics": [
    {{
      "raw_topic": "original generated topic",
      "canonical_topic": "shared canonical label"
    }}
  ],
  "notes": ""
}}
""".strip()
    return create_chat_completion(
        client,
        resolved_model_name=resolved_model_name,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    )


def evaluate_core_coverage(
    *,
    paper_id: str,
    title: str,
    original_slide_dir: Path,
    generated_pptx: Path,
    model: str,
    max_original_slides: int,
) -> dict[str, Any]:
    client = build_openai_client()
    resolved_model_name = resolve_direct_model_name(model)

    original_slide_paths = sorted(
        [path for path in original_slide_dir.iterdir() if path.suffix.lower() in SLIDE_EXTENSIONS],
        key=numeric_slide_sort_key,
    )
    effective_max_original_slides = max_original_slides
    if effective_max_original_slides <= 0 or effective_max_original_slides > MAX_IMAGES_PER_REQUEST:
        effective_max_original_slides = MAX_IMAGES_PER_REQUEST
    sampled_slide_paths = sample_slide_paths(original_slide_paths, effective_max_original_slides)
    generated_text = extract_pptx_text(generated_pptx)

    system_prompt = (
        "You are an evaluation agent for scientific slide decks.\n"
        "Compare a generated presentation against the original presentation and extract semantic topics from each.\n"
        "Do not rely on exact slide titles or exact section headers.\n"
        "Treat topics as the main scientific ideas actually covered in the deck, not generic slide functions like agenda or thank you.\n"
        "Canonicalize semantically equivalent topics to the same short label even when phrased differently.\n"
        "Keep topics concise, specific, and deduplicated.\n"
        "Only include topics that are materially present in the deck.\n"
        "Return only valid JSON."
    )
    user_prompt = f"""
Paper ID: {paper_id}
Paper Title: {title}

Original presentation:
- You are given sampled images from the original slide deck.
- Treat them as representative evidence of the original deck's content and structure.

Generated presentation:
- You are given extracted text from the generated PPTX below.

Generated deck text:
{generated_text}

Task:
1. Extract a set of canonical reference topics from the original deck.
2. Extract a set of canonical generated topics from the generated deck.
3. Use the same wording whenever two topics are semantically the same across decks.
4. Summarize how the decks differ in what they emphasize, omit, compress, or restructure.

Return exactly this JSON schema:
{{
  "paper_id": "{paper_id}",
  "title": "{title}",
  "reference_topics": [
    "topic one"
  ],
  "generated_topics": [
    "topic one"
  ],
  "difference_summary": {{
    "missing_content": [],
    "compressed_content": [],
    "overemphasized_or_extra_content": [],
    "structure_changes": [],
    "overall_summary": ""
  }},
  "notes": ""
}}
""".strip()

    content: list[dict[str, Any]] = [{"type": "text", "text": user_prompt}]
    for slide_path in sampled_slide_paths:
        content.append(
            {
                "type": "text",
                "text": f"Original slide image: {slide_path.name}",
            }
        )
        content.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": encode_image_data_uri(slide_path),
                    "detail": "low",
                },
            }
        )

    try:
        result = create_chat_completion(
            client,
            resolved_model_name=resolved_model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": content},
            ],
        )
    except Exception as exc:
        if is_content_policy_violation(exc):
            return {
                "paper_id": paper_id,
                "title": title,
                "reference_topics": [],
                "generated_topics": [],
                "matched_topics": [],
                "matched_topic_pairs": [],
                "reference_only_topics": [],
                "generated_only_topics": [],
                "reference_topic_count": 0,
                "generated_topic_count": 0,
                "topic_intersection_count": 0,
                "topic_union_count": 0,
                "topic_iou": None,
                "coverage_ratio": None,
                "missing_from_generated": [],
                "difference_summary": {
                    "missing_content": [],
                    "compressed_content": [],
                    "overemphasized_or_extra_content": [],
                    "structure_changes": [],
                    "overall_summary": "",
                },
                "notes": "Skipped core coverage because one or more sampled original slide images triggered the model content-safety filter.",
                "sampled_original_slides": [path.name for path in sampled_slide_paths],
                "generated_pptx": str(generated_pptx),
                "skipped": True,
                "skip_reason": "content_policy_violation",
            }
        raise

    raw_reference_topics = sanitize_topic_list(result.get("reference_topics"))
    raw_generated_topics = sanitize_topic_list(result.get("generated_topics"))
    normalized_result = normalize_topics_with_llm(
        client=client,
        resolved_model_name=resolved_model_name,
        paper_id=paper_id,
        title=title,
        reference_topics=raw_reference_topics,
        generated_topics=raw_generated_topics,
    )
    normalized_result["difference_summary"] = result.get("difference_summary", {}) or {}
    normalized_result["notes"] = str(normalized_result.get("notes") or result.get("notes") or "").strip()

    sanitized = sanitize_result(normalized_result, paper_id=paper_id, title=title)
    sanitized["raw_reference_topics"] = raw_reference_topics
    sanitized["raw_generated_topics"] = raw_generated_topics
    sanitized["sampled_original_slides"] = [path.name for path in sampled_slide_paths]
    sanitized["generated_pptx"] = str(generated_pptx)
    return sanitized


def main() -> None:
    parser = argparse.ArgumentParser(description="Evaluate generated decks with semantic topic-set IoU.")
    parser.add_argument("--paper-id", required=True, help="Paper id, e.g. eccv20:589")
    parser.add_argument("--title", required=True, help="Paper title")
    parser.add_argument("--original-slide-dir", type=Path, required=True, help="Directory of original slide images")
    parser.add_argument("--generated-pptx", type=Path, required=True, help="Generated PPTX path")
    parser.add_argument("--model", default="gpt-5.4-nano", help="Model or Azure deployment alias")
    parser.add_argument(
        "--max-original-slides",
        type=int,
        default=0,
        help="Max original slide images to include; use 0 to include all slides (default: 0)",
    )
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    args = parser.parse_args()

    load_dotenv(REPO_ROOT / ".env")
    args.output_dir.mkdir(parents=True, exist_ok=True)

    result = evaluate_core_coverage(
        paper_id=args.paper_id,
        title=args.title,
        original_slide_dir=args.original_slide_dir,
        generated_pptx=args.generated_pptx,
        model=args.model,
        max_original_slides=args.max_original_slides,
    )

    output_path = args.output_dir / f"{args.paper_id.replace(':', '_')}.core_coverage.json"
    output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Saved evaluation to {output_path}")


if __name__ == "__main__":
    main()
