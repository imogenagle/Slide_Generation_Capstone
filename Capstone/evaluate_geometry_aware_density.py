#!/usr/bin/env python3
"""Evaluate Geometry-Aware Density (GAD) for generated slide decks."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

try:
    from dotenv import load_dotenv
except ModuleNotFoundError:
    def load_dotenv(*_args: Any, **_kwargs: Any) -> bool:
        return False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ModuleNotFoundError as exc:  # pragma: no cover - import guard for environments without deps
    raise SystemExit(
        "python-pptx is required for evaluate_geometry_aware_density.py. "
        "Install project dependencies first."
    ) from exc


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_DIR = REPO_ROOT / "Capstone" / "evaluations"


Rect = tuple[float, float, float, float]


def clamp(value: float, low: float = 0.0, high: float = 1.0) -> float:
    return max(low, min(high, value))


def flatten_shapes(shapes: Iterable[Any]) -> list[Any]:
    flat: list[Any] = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            flat.extend(flatten_shapes(shape.shapes))
        else:
            flat.append(shape)
    return flat


def shape_rect(shape: Any) -> Rect | None:
    width = float(getattr(shape, "width", 0.0) or 0.0)
    height = float(getattr(shape, "height", 0.0) or 0.0)
    if width <= 0 or height <= 0:
        return None
    left = float(getattr(shape, "left", 0.0) or 0.0)
    top = float(getattr(shape, "top", 0.0) or 0.0)
    return (left, top, left + width, top + height)


def rect_area(rect: Rect) -> float:
    left, top, right, bottom = rect
    return max(0.0, right - left) * max(0.0, bottom - top)


def has_meaningful_content(shape: Any) -> bool:
    if getattr(shape, "has_table", False):
        return True
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return True
    if getattr(shape, "has_text_frame", False):
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if getattr(paragraph, "text", "").strip()
        ).strip()
        if text:
            return True
    return False


def is_background_like(shape: Any, rect: Rect, slide_area: float, background_area_ratio: float) -> bool:
    if has_meaningful_content(shape):
        return False
    area_ratio = rect_area(rect) / slide_area if slide_area > 0 else 0.0
    if area_ratio < background_area_ratio:
        return False
    return True


def extract_content_rects(
    slide: Any,
    *,
    slide_width: float,
    slide_height: float,
    area_min_ratio: float,
    background_area_ratio: float,
) -> tuple[list[Rect], list[dict[str, Any]]]:
    slide_area = slide_width * slide_height
    min_area = slide_area * area_min_ratio
    rects: list[Rect] = []
    shape_debug: list[dict[str, Any]] = []

    for shape in flatten_shapes(slide.shapes):
        rect = shape_rect(shape)
        if rect is None:
            continue
        area = rect_area(rect)
        if area < min_area:
            continue
        if is_background_like(shape, rect, slide_area, background_area_ratio):
            continue

        shape_type = getattr(shape, "shape_type", None)
        shape_name = getattr(shape, "name", None) or "shape"
        contentful = has_meaningful_content(shape)
        if not contentful and shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        rects.append(rect)
        shape_debug.append(
            {
                "name": str(shape_name),
                "shape_type": int(shape_type) if shape_type is not None else None,
                "left": round(rect[0], 2),
                "top": round(rect[1], 2),
                "right": round(rect[2], 2),
                "bottom": round(rect[3], 2),
                "width": round(rect[2] - rect[0], 2),
                "height": round(rect[3] - rect[1], 2),
                "area_ratio": round(area / slide_area, 5) if slide_area > 0 else 0.0,
            }
        )

    return rects, shape_debug


def union_area(rects: list[Rect]) -> float:
    if not rects:
        return 0.0

    xs = sorted({rect[0] for rect in rects} | {rect[2] for rect in rects})
    ys = sorted({rect[1] for rect in rects} | {rect[3] for rect in rects})
    if len(xs) < 2 or len(ys) < 2:
        return 0.0

    area = 0.0
    for x1, x2 in zip(xs, xs[1:]):
        if x2 <= x1:
            continue
        overlapping = [rect for rect in rects if rect[0] < x2 and rect[2] > x1]
        if not overlapping:
            continue
        for y1, y2 in zip(ys, ys[1:]):
            if y2 <= y1:
                continue
            if any(rect[1] < y2 and rect[3] > y1 for rect in overlapping):
                area += (x2 - x1) * (y2 - y1)
    return area


def effective_region_count(rects: list[Rect]) -> int:
    # The paper definition provided here treats each qualifying content box as one
    # effective region once it clears the a_min threshold.
    return len(rects)


def score_slide(
    rects: list[Rect],
    *,
    slide_width: float,
    slide_height: float,
    tau: float,
    m_star: float,
    kappa: float,
    lambda_occupancy: float,
    lambda_fragmentation: float,
) -> dict[str, Any]:
    slide_area = slide_width * slide_height
    occupancy = union_area(rects) / slide_area if slide_area > 0 else 0.0
    effective_regions = effective_region_count(rects)
    occupancy_matching = clamp(1.0 - abs(occupancy - tau))
    fragmentation_reward = clamp(1.0 - ((effective_regions - m_star) ** 2) / kappa) if kappa > 0 else 0.0
    geometry_score = (lambda_occupancy * occupancy_matching) + (lambda_fragmentation * fragmentation_reward)
    return {
        "area_occupancy": round(occupancy, 5),
        "occupancy_matching": round(occupancy_matching, 5),
        "effective_region_count": effective_regions,
        "fragmentation_reward": round(fragmentation_reward, 5),
        "geometry_score": round(geometry_score, 5),
    }


def evaluate_geometry_aware_density(
    *,
    pptx_path: Path,
    tau: float,
    m_star: float,
    kappa: float,
    lambda_occupancy: float,
    lambda_fragmentation: float,
    area_min_ratio: float,
    background_area_ratio: float,
) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_width = float(prs.slide_width)
    slide_height = float(prs.slide_height)

    per_slide: list[dict[str, Any]] = []
    geometry_scores: list[float] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        rects, shape_debug = extract_content_rects(
            slide,
            slide_width=slide_width,
            slide_height=slide_height,
            area_min_ratio=area_min_ratio,
            background_area_ratio=background_area_ratio,
        )
        slide_result = score_slide(
            rects,
            slide_width=slide_width,
            slide_height=slide_height,
            tau=tau,
            m_star=m_star,
            kappa=kappa,
            lambda_occupancy=lambda_occupancy,
            lambda_fragmentation=lambda_fragmentation,
        )
        slide_result["slide_index"] = slide_index
        slide_result["content_shape_count"] = len(rects)
        slide_result["content_shapes"] = shape_debug
        per_slide.append(slide_result)
        geometry_scores.append(slide_result["geometry_score"])

    deck_score = sum(geometry_scores) / len(geometry_scores) if geometry_scores else 0.0
    return {
        "pptx_path": str(pptx_path),
        "metric": "geometry_aware_density",
        "parameters": {
            "tau": tau,
            "m_star": m_star,
            "kappa": kappa,
            "lambda_occupancy": lambda_occupancy,
            "lambda_fragmentation": lambda_fragmentation,
            "area_min_ratio": area_min_ratio,
            "background_area_ratio": background_area_ratio,
        },
        "slide_size_emu": {
            "width": round(slide_width, 2),
            "height": round(slide_height, 2),
        },
        "slide_count": len(per_slide),
        "gad_geom": round(deck_score, 5),
        "per_slide": per_slide,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Evaluate Geometry-Aware Density (GAD) for a PPTX deck.")
    parser.add_argument("--pptx-path", type=Path, required=True, help="Path to the generated PPTX.")
    parser.add_argument("--tau", type=float, default=0.5, help="Target area occupancy ratio.")
    parser.add_argument("--m-star", type=float, default=4.0, help="Target effective region count.")
    parser.add_argument("--kappa", type=float, default=6.3, help="Quadratic fragmentation width.")
    parser.add_argument(
        "--lambda-occupancy",
        type=float,
        default=0.5,
        help="Weight for occupancy matching reward.",
    )
    parser.add_argument(
        "--lambda-fragmentation",
        type=float,
        default=0.5,
        help="Weight for fragmentation reward.",
    )
    parser.add_argument(
        "--area-min-ratio",
        type=float,
        default=0.005,
        help="Minimum shape area ratio to count toward content occupancy.",
    )
    parser.add_argument(
        "--background-area-ratio",
        type=float,
        default=0.9,
        help="Ignore non-content shapes at or above this slide-area ratio as background-like.",
    )
    parser.add_argument("--output", type=Path, default=None, help="Optional output JSON path.")
    args = parser.parse_args()

    load_dotenv(REPO_ROOT / ".env")

    lambda_sum = args.lambda_occupancy + args.lambda_fragmentation
    if lambda_sum <= 0:
        raise SystemExit("lambda weights must sum to a positive value.")
    lambda_occupancy = args.lambda_occupancy / lambda_sum
    lambda_fragmentation = args.lambda_fragmentation / lambda_sum

    result = evaluate_geometry_aware_density(
        pptx_path=args.pptx_path,
        tau=args.tau,
        m_star=args.m_star,
        kappa=args.kappa,
        lambda_occupancy=lambda_occupancy,
        lambda_fragmentation=lambda_fragmentation,
        area_min_ratio=args.area_min_ratio,
        background_area_ratio=args.background_area_ratio,
    )

    output_text = json.dumps(result, indent=2, ensure_ascii=False)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(output_text, encoding="utf-8")
    print(output_text)


if __name__ == "__main__":
    main()
