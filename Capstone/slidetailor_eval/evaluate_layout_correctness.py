#!/usr/bin/env python3
"""Deterministic layout-correctness evaluation for SlideGen decks."""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if __package__ in {None, ""}:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[3]))
    from SlideGen.Capstone.slidetailor_eval.common import REPO_ROOT, resolve_output_path
else:
    from .common import REPO_ROOT, resolve_output_path


EMU_PER_PT = 12700.0
CONTENT_TYPES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,
    MSO_SHAPE_TYPE.CALLOUT,
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.GROUP,
    MSO_SHAPE_TYPE.PICTURE,
    MSO_SHAPE_TYPE.PLACEHOLDER,
    MSO_SHAPE_TYPE.TABLE,
    MSO_SHAPE_TYPE.TEXT_BOX,
}


def _paragraph_text(paragraph) -> str:
    if paragraph is None:
        return ""
    text = paragraph.text or ""
    if text:
        return text
    return "".join(run.text for run in paragraph.runs if getattr(run, "text", ""))


def _paragraph_font_size_pt(paragraph, fallback: float) -> float:
    if paragraph is None:
        return fallback
    if paragraph.font is not None and paragraph.font.size is not None:
        return max(float(paragraph.font.size.pt), 1.0)
    for run in paragraph.runs:
        if run.font is not None and run.font.size is not None:
            return max(float(run.font.size.pt), 1.0)
    return fallback


def _estimate_wrapped_line_count(text: str, usable_width_pt: float, font_size_pt: float) -> int:
    text = (text or "").strip()
    if not text:
        return 1
    chars_per_line = max(4, int(usable_width_pt / max(font_size_pt * 0.5, 1.0)))
    line_count = 0
    for logical_line in text.splitlines() or [""]:
        line_count += max(1, math.ceil(max(len(logical_line), 1) / chars_per_line))
    return line_count


def _text_frame_fits_estimate(shape, text_frame, min_size_pt: float) -> float:
    width_pt = max(shape.width / EMU_PER_PT - 2 * 3.0, 24.0)
    height_pt = max(shape.height / EMU_PER_PT - 2 * 3.0, 24.0)
    total_height_pt = 0.0
    paragraphs = list(text_frame.paragraphs or [])

    for idx, paragraph in enumerate(paragraphs):
        para_text = _paragraph_text(paragraph)
        font_size_pt = max(_paragraph_font_size_pt(paragraph, min_size_pt), min_size_pt)
        level = max(int(getattr(paragraph, "level", 0) or 0), 0)
        indent_penalty_pt = min(level * 18.0, width_pt * 0.45)
        usable_width_pt = max(width_pt - indent_penalty_pt, 24.0)
        line_count = _estimate_wrapped_line_count(para_text, usable_width_pt, font_size_pt)
        line_height_pt = font_size_pt * (1.18 if level == 0 else 1.12)
        total_height_pt += line_count * line_height_pt
        if idx < len(paragraphs) - 1:
            total_height_pt += max(font_size_pt * 0.18, 1.5)

    return total_height_pt - height_pt


def _shape_rect(shape) -> tuple[int, int, int, int]:
    return (int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height))


def _rects_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int], padding: int = 0) -> bool:
    return not (
        a[2] <= b[0] + padding
        or a[0] >= b[2] - padding
        or a[3] <= b[1] + padding
        or a[1] >= b[3] - padding
    )


def _shape_has_content(shape) -> bool:
    if getattr(shape, "has_text_frame", False):
        tf = shape.text_frame
        if tf and any(_paragraph_text(p).strip() for p in tf.paragraphs):
            return True
    if getattr(shape, "has_table", False):
        return True
    return getattr(shape, "shape_type", None) in CONTENT_TYPES


def _is_text_shape(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    return bool(tf and any(_paragraph_text(p).strip() for p in tf.paragraphs))


def _is_title_shape(shape, slide_height: int) -> bool:
    name = (getattr(shape, "name", "") or "").lower()
    if "title" in name:
        return True
    top = int(getattr(shape, "top", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return top < int(slide_height * 0.2) and height < int(slide_height * 0.2)


def _shape_kind(shape) -> str:
    if _is_text_shape(shape):
        return "text"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    return "other"


def _collect_overflow_issues(shape, slide_height: int) -> list[dict[str, Any]]:
    if not _is_text_shape(shape):
        return []
    min_size_pt = 18.0 if _is_title_shape(shape, slide_height) else 15.0
    overflow_pt = _text_frame_fits_estimate(shape, shape.text_frame, min_size_pt)
    if overflow_pt <= 0:
        return []
    preview = " ".join(
        _paragraph_text(paragraph).strip()
        for paragraph in shape.text_frame.paragraphs
        if _paragraph_text(paragraph).strip()
    )[:160]
    return [
        {
            "type": "title_text_overflow" if _is_title_shape(shape, slide_height) else "body_text_overflow",
            "shape_name": getattr(shape, "name", ""),
            "overflow_pt": round(float(overflow_pt), 2),
            "text_preview": preview,
        }
    ]


def _collect_offslide_issues(shape, slide_width: int, slide_height: int) -> list[dict[str, Any]]:
    left, top, right, bottom = _shape_rect(shape)
    tolerance = int(2 * EMU_PER_PT)
    issues: list[dict[str, Any]] = []
    if left < -tolerance or top < -tolerance or right > slide_width + tolerance or bottom > slide_height + tolerance:
        issues.append(
            {
                "type": "offslide_text" if _is_text_shape(shape) else "offslide_visual",
                "shape_name": getattr(shape, "name", ""),
                "bounds": {"left": left, "top": top, "right": right, "bottom": bottom},
            }
        )
    return issues


def _collect_overlap_issues(content_shapes: list[Any]) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    padding = int(1.5 * EMU_PER_PT)
    for idx, shape_a in enumerate(content_shapes):
        rect_a = _shape_rect(shape_a)
        for shape_b in content_shapes[idx + 1 :]:
            rect_b = _shape_rect(shape_b)
            if not _rects_overlap(rect_a, rect_b, padding=padding):
                continue
            kind_a = _shape_kind(shape_a)
            kind_b = _shape_kind(shape_b)
            issues.append(
                {
                    "type": "overlap",
                    "shape_a": getattr(shape_a, "name", ""),
                    "shape_b": getattr(shape_b, "name", ""),
                    "kind_a": kind_a,
                    "kind_b": kind_b,
                }
            )
    return issues


def _count_slide_defects(issues: list[dict[str, Any]]) -> dict[str, int]:
    counts = {
        "body_text_overflow": 0,
        "title_text_overflow": 0,
        "offslide_text": 0,
        "offslide_visual": 0,
        "overlap_text_related": 0,
        "overlap_visual_only": 0,
    }

    for issue in issues:
        issue_type = issue["type"]
        if issue_type == "body_text_overflow":
            counts["body_text_overflow"] += 1
        elif issue_type == "title_text_overflow":
            counts["title_text_overflow"] += 1
        elif issue_type == "offslide_text":
            counts["offslide_text"] += 1
        elif issue_type == "offslide_visual":
            counts["offslide_visual"] += 1
        elif issue_type == "overlap":
            if "text" in (issue.get("kind_a"), issue.get("kind_b")):
                counts["overlap_text_related"] += 1
            else:
                counts["overlap_visual_only"] += 1
    return counts


def _sum_defects(issue_counts: dict[str, int]) -> int:
    return sum(int(value) for value in issue_counts.values())


def _major_defect_count(issue_counts: dict[str, int]) -> int:
    return (
        int(issue_counts.get("body_text_overflow", 0))
        + int(issue_counts.get("title_text_overflow", 0))
        + int(issue_counts.get("offslide_text", 0))
        + int(issue_counts.get("overlap_text_related", 0))
    )


def _aggregate_counts(per_slide: list[dict[str, Any]]) -> dict[str, int]:
    totals = {
        "body_text_overflow": 0,
        "title_text_overflow": 0,
        "offslide_text": 0,
        "offslide_visual": 0,
        "overlap_text_related": 0,
        "overlap_visual_only": 0,
    }
    for item in per_slide:
        counts = item["issue_counts"]
        for key in totals:
            totals[key] += int(counts.get(key, 0))
    return totals


def evaluate_layout_correctness(*, pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    per_slide: list[dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        content_shapes = [shape for shape in slide.shapes if _shape_has_content(shape)]
        issues: list[dict[str, Any]] = []

        for shape in content_shapes:
            issues.extend(_collect_offslide_issues(shape, slide_width, slide_height))
            issues.extend(_collect_overflow_issues(shape, slide_height))

        issues.extend(_collect_overlap_issues(content_shapes))
        issue_counts = _count_slide_defects(issues)
        defect_count = _sum_defects(issue_counts)
        major_defect_count = _major_defect_count(issue_counts)

        per_slide.append(
            {
                "slide_index": slide_idx,
                "issue_counts": issue_counts,
                "defect_count": defect_count,
                "major_defect_count": major_defect_count,
                "issues": issues,
            }
        )

    slide_count = len(per_slide)
    total_issue_counts = _aggregate_counts(per_slide)
    total_defect_count = sum(total_issue_counts.values())
    total_major_defect_count = sum(item["major_defect_count"] for item in per_slide)
    slides_with_any_defect = sum(1 for item in per_slide if item["defect_count"] > 0)
    slides_with_major_defect = sum(1 for item in per_slide if item["major_defect_count"] > 0)
    defects_per_slide = round(total_defect_count / slide_count, 4) if slide_count else 0.0
    major_defects_per_slide = round(total_major_defect_count / slide_count, 4) if slide_count else 0.0

    return {
        "source": "SlideGen custom",
        "metric": "layout_correctness",
        "pptx_path": str(pptx_path),
        "slide_count": slide_count,
        "total_issue_counts": total_issue_counts,
        "total_defect_count": total_defect_count,
        "total_major_defect_count": total_major_defect_count,
        "slides_with_any_defect": slides_with_any_defect,
        "slides_with_major_defect": slides_with_major_defect,
        "defects_per_slide": defects_per_slide,
        "major_defects_per_slide": major_defects_per_slide,
        "deck_score": defects_per_slide,
        "per_slide": per_slide,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Evaluate deterministic layout correctness.")
    parser.add_argument("--pptx-path", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=None)
    args = parser.parse_args()

    result = evaluate_layout_correctness(pptx_path=args.pptx_path)
    output_path = resolve_output_path("layout_correctness", args.output, args.pptx_path.stem)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
